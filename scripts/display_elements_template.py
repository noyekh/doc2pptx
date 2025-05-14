from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def dump_layouts(prs):
    print("=== LAYOUTS ===")
    for li, layout in enumerate(prs.slide_layouts):
        print(f"\nLayout {li}: '{layout.name}'")
        for shape in layout.shapes:
            parts = []
            if shape.is_placeholder:
                idx   = shape.placeholder_format.idx
                ptype = str(shape.placeholder_format.type).split(" ")[0]
                parts.append(f"Placeholder idx={idx} type={ptype}")
            else:
                # shape_type is just an int; you could map it back to names via MSO_SHAPE_TYPE if desired
                parts.append(f"ShapeType={shape.shape_type}")
            parts.append(f"name='{shape.name}'")
            parts.append(f"pos=({shape.left.pt:.1f},{shape.top.pt:.1f})")
            parts.append(f"size=({shape.width.pt:.1f}×{shape.height.pt:.1f})")
            text = shape.text.strip() if hasattr(shape, "text") else ""
            if text:
                parts.append(f"text=\"{text}\"")
            print("  • " + "; ".join(parts))

def dump_slides(prs):
    print("\n=== SLIDES ===")
    for si, slide in enumerate(prs.slides, start=1):
        print(f"\nSlide {si}: uses layout '{slide.slide_layout.name}'")
        for shape in slide.shapes:
            parts = []
            if shape.is_placeholder:
                idx   = shape.placeholder_format.idx
                ptype = str(shape.placeholder_format.type).split(" ")[0]
                parts.append(f"Placeholder idx={idx} type={ptype}")
            else:
                parts.append(f"ShapeType={shape.shape_type}")
            parts.append(f"name='{shape.name}'")
            text = shape.text.strip() if hasattr(shape, "text") else ""
            if text:
                parts.append(f"text=\"{text}\"")
            print("  • " + "; ".join(parts))

if __name__ == "__main__":
    prs = Presentation("base_template.pptx")
    dump_layouts(prs)
    dump_slides(prs)