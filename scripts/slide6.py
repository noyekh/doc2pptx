#!/usr/bin/env python3

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL

def create_editable_template(pptx_path, output_path=None):
    """
    Analyze a PPTX file and provide information on how to edit all text boxes
    
    Args:
        pptx_path: Path to the input PPTX file
        output_path: Optional path to save a copy of the PPTX
    
    Returns:
        Dictionary mapping shape indices to their content and properties
    """
    prs = Presentation(pptx_path)
    
    # Dictionary to store text box information
    text_boxes = {}
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_text_boxes = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                box_info = {
                    "slide_idx": slide_idx + 1,
                    "shape_idx": shape_idx + 1,
                    "text": shape.text,
                    "name": shape.name,
                    "fill_type": str(shape.fill.type) if hasattr(shape, "fill") else "None"
                }
                
                # Get text color if available
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if hasattr(run.font.color, "rgb"):
                                box_info["text_color"] = (
                                    run.font.color.rgb.r,
                                    run.font.color.rgb.g,
                                    run.font.color.rgb.b
                                )
                                break
                        if "text_color" in box_info:
                            break
                
                slide_text_boxes.append(box_info)
        
        text_boxes[f"slide_{slide_idx+1}"] = slide_text_boxes
    
    # Print information about all text boxes
    print("\nText Boxes in the Presentation:")
    for slide_key, boxes in text_boxes.items():
        print(f"\n{slide_key}:")
        for i, box in enumerate(boxes):
            print(f"  Box {i+1}:")
            print(f"    Shape Index: {box['shape_idx']}")
            print(f"    Name: {box['name']}")
            print(f"    Fill Type: {box['fill_type']}")
            if "text_color" in box:
                print(f"    Text Color: RGB{box['text_color']}")
            print(f"    Text: {box['text'][:50]}..." if len(box['text']) > 50 else f"    Text: {box['text']}")
    
    # Save a copy if output_path is provided
    if output_path:
        prs.save(output_path)
        print(f"\nSaved a copy to {output_path}")
    
    print("\nTo update text and color, use:")
    print("update_text_box(pptx_path, slide_index, shape_index, new_text, new_color)")
    
    return text_boxes

def update_text_box(pptx_path, output_path, slide_index, shape_index, new_text=None, new_color=None):
    """
    Update the text and/or color of a specific text box
    
    Args:
        pptx_path: Path to the input PPTX file
        output_path: Path where the modified PPTX will be saved
        slide_index: 1-based index of the slide
        shape_index: 1-based index of the shape
        new_text: New text content (None to keep existing)
        new_color: RGBColor for the text (None to keep existing)
    """
    prs = Presentation(pptx_path)
    
    # Convert to 0-based indices
    slide_idx = slide_index - 1
    shape_idx = shape_index - 1
    
    if slide_idx < 0 or slide_idx >= len(prs.slides):
        print(f"Error: Slide index {slide_index} is out of range")
        return False
    
    slide = prs.slides[slide_idx]
    
    if shape_idx < 0 or shape_idx >= len(slide.shapes):
        print(f"Error: Shape index {shape_index} on slide {slide_index} is out of range")
        return False
    
    shape = slide.shapes[shape_idx]
    
    if not hasattr(shape, "text"):
        print(f"Error: Shape {shape_index} on slide {slide_index} does not have text")
        return False
    
    # Update text if specified
    if new_text is not None:
        # Clear existing text
        text_frame = shape.text_frame
        for i in range(len(text_frame.paragraphs)-1, -1, -1):
            p = text_frame.paragraphs[i]
            if i == 0:
                p.text = ""
                run = p.add_run()
                run.text = new_text
                
                # Set color if specified
                if new_color:
                    run.font.color.rgb = new_color
            else:
                # Remove extra paragraphs
                tr = text_frame._txBody.remove(text_frame._txBody.p_lst[i])
        
        print(f"Updated text of shape {shape_index} on slide {slide_index}")
    
    # Update color only if not updating text but color is specified
    elif new_color:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = new_color
        
        print(f"Updated color of shape {shape_index} on slide {slide_index}")
    
    prs.save(output_path)
    return True

def update_all_solid_box_text_colors(pptx_path, output_path, new_color):
    """
    Update the text color in all boxes with a SOLID fill
    
    Args:
        pptx_path: Path to the input PPTX file
        output_path: Path where the modified PPTX will be saved
        new_color: RGBColor for the text
    """
    prs = Presentation(pptx_path)
    
    boxes_modified = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            # Check if this is a text box with SOLID fill
            if (hasattr(shape, "text") and shape.text.strip() and
                hasattr(shape, "fill") and shape.fill.type == MSO_FILL.SOLID):
                
                # Update the text color
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = new_color
                        boxes_modified += 1
                
                print(f"Updated color in shape {shape_idx+1} on slide {slide_idx+1}")
    
    print(f"Updated text color in {boxes_modified} runs across all SOLID fill boxes")
    prs.save(output_path)
    return boxes_modified

# Path to your PPTX file
pptx_file = "slide6.pptx"
template_file = "slide6_template.pptx"
modified_file = "slide6_modified.pptx"

# Create an editable template with all text boxes identified
text_boxes = create_editable_template(pptx_file, template_file)

# Example: Update all solid fill boxes to have blue text
update_all_solid_box_text_colors(template_file, modified_file, RGBColor(0, 0, 255))

# Example: Update a specific text box
# To update the text in the first box of slide 1 (assuming it's shape index 2 based on your analysis)
# update_text_box(template_file, "slide6_updated.pptx", 1, 2, "Nouveau texte pour cette boîte", RGBColor(255, 0, 0))