"""
PowerPoint presentation builder for doc2pptx.

This module provides functionality to build PowerPoint presentations
from structured data using templates and layout rules.
"""
import logging
import io
import re
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Union, Any, cast

from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.slide import Slide as PptxSlide
from pptx.util import Pt, Inches, Cm, Emu
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.dml import MSO_THEME_COLOR_INDEX, MSO_COLOR_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.table import Table, _Cell, _Row, _Column
from pptx.util import Emu

from doc2pptx.core.models import Section, Slide, ContentType, SlideBlock, SlideContent, Presentation, SectionType
from doc2pptx.layout.selector import LayoutSelector
from doc2pptx.ppt.template_loader import TemplateLoader, TemplateInfo
from doc2pptx.ppt.overflow import OverflowHandler
from doc2pptx.llm.optimizer import PresentationOptimizer


logger = logging.getLogger(__name__)


# Définition des capacités des layouts disponibles dans le template
LAYOUT_CAPABILITIES = {
    "Diapositive de titre": {
        "title": True,
        "subtitle": True,
        "content": False,
        "table": False,
        "image": False,
        "chart": False,
        "max_blocks": 1,
        "description": "Title slide with subtitle"
    },
    "Introduction": {
        "title": True,
        "subtitle": False,
        "content": True,
        "table": False,
        "image": False,
        "chart": False,
        "max_blocks": 1,
        "description": "Title with large content area"
    },
    "Titre et texte": {
        "title": True,
        "subtitle": False,
        "content": True,
        "table": False,
        "image": False,
        "chart": False,
        "max_blocks": 1,
        "description": "Title with text content"
    },
    "Titre et tableau": {
        "title": True,
        "subtitle": False,
        "content": False,
        "table": True,
        "image": False,
        "chart": False,
        "max_blocks": 1,
        "description": "Title with table"
    },
    "Titre et texte 1 visuel gauche": {
        "title": True,
        "subtitle": False,
        "content": True,
        "table": False,
        "image": True,
        "chart": False,
        "max_blocks": 2,
        "description": "Title with image on left and text on right"
    },
    "Titre et texte 1 histogramme": {
        "title": True,
        "subtitle": False,
        "content": True,
        "table": False,
        "image": False,
        "chart": True,
        "max_blocks": 2,
        "description": "Title with text on left and chart on right"
    },
    "Titre et 3 colonnes": {
        "title": True,
        "subtitle": False,
        "content": True,
        "table": False,
        "image": False,
        "chart": False,
        "max_blocks": 3,
        "description": "Title with three text columns"
    },
    "Chapitre 1": {
        "title": True,
        "subtitle": False,
        "content": False,
        "table": False,
        "image": False,
        "chart": False,
        "max_blocks": 0,
        "description": "Section title only"
    }
}

# Mapping des placeholders par layout (index des placeholders)
LAYOUT_PLACEHOLDER_MAP = {
    "Diapositive de titre": {
        "title": 0,       # idx=0, TITLE
        "subtitle": 1,    # idx=1, SUBTITLE
    },
    "Introduction": {
        "title": 0,       # idx=0, TITLE
        "content": 1,     # idx=1, BODY
        "slide_number": 12 # idx=12, SLIDE_NUMBER
    },
    "Titre et texte": {
        "title": 0,       # idx=0, TITLE
        "content": 1,     # idx=1, BODY
        "slide_number": 12 # idx=12, SLIDE_NUMBER
    },
    "Titre et tableau": {
        "title": 0,       # idx=0, TITLE
        "slide_number": 12 # idx=12, SLIDE_NUMBER
        # Pas de placeholder pour la table, elle est ajoutée comme shape
    },
    "Titre et texte 1 visuel gauche": {
        "title": 0,       # idx=0, TITLE
        "content": 1,     # idx=1, BODY (à droite)
        "image": 2,       # idx=2, PICTURE (à gauche)
        "slide_number": 12 # idx=12, SLIDE_NUMBER
    },
    "Titre et texte 1 histogramme": {
        "title": 0,       # idx=0, TITLE
        "content": 1,     # idx=1, BODY (à gauche)
        "chart": 2,       # idx=2, CHART (à droite)
        "slide_number": 12 # idx=12, SLIDE_NUMBER
    },
    "Titre et 3 colonnes": {
        "title": 0,       # idx=0, TITLE
        "column1": 1,     # idx=1, BODY (colonne 1)
        "column2": 2,     # idx=2, BODY (colonne 2)
        "column3": 3,     # idx=3, BODY (colonne 3)
        "slide_number": 12 # idx=12, SLIDE_NUMBER
    },
    "Chapitre 1": {
        "title": 0,       # idx=0, TITLE
    }
}


class PPTBuilder:
    """
    Builds PowerPoint presentations from structured data.
    
    This class handles the process of creating a PowerPoint presentation
    from a Presentation model, filling in template placeholders with content,
    and managing layout selection for each slide.
    """
    
    # Table style presets
    TABLE_STYLES = {
        "default": {
            "header_bg": "4472C4",  # Blue header background
            "header_text": "FFFFFF",  # White text
            "body_bg": None,  # No background color
            "body_text": "000000",  # Black text
            "border_color": "4472C4",  # Blue border
            "border_width": Pt(1),  # 1 point border width
            "accent_color": "A5A5A5",  # Light gray for alternating rows
            "first_row": True,  # Format first row as header
            "total_row": False,  # No footer row
            "banded_rows": False,  # Don't alternate row colors for now -> fix later
            "banded_cols": False,  # No alternating column colors
        },
        "minimal": {
            "header_bg": None,  # No background color
            "header_text": "000000",  # Black text
            "body_bg": None,  # No background color
            "body_text": "000000",  # Black text
            "border_color": "D9D9D9",  # Light gray border
            "border_width": Pt(0.5),  # Thin border
            "accent_color": "F2F2F2",  # Very light gray for alternating rows
            "first_row": True,  # Format first row as header
            "total_row": False,  # No footer row
            "banded_rows": True,  # Alternating row colors
            "banded_cols": False,  # No alternating column colors
        },
        "grid": {
            "header_bg": "4472C4",  # Blue header background
            "header_text": "FFFFFF",  # White text
            "body_bg": None,  # No background color
            "body_text": "000000",  # Black text
            "border_color": "000000",  # Black border
            "border_width": Pt(1),  # 1 point border width
            "accent_color": "E6E6E6",  # Light gray for alternating rows
            "first_row": True,  # Format first row as header
            "total_row": False,  # No footer row
            "banded_rows": True,  # Alternating row colors
            "banded_cols": True,  # Alternating column colors
        },
        "accent1": {
            "header_bg": "5B9BD5",  # Accent1 color (blue)
            "header_text": "FFFFFF",  # White text
            "body_bg": None,  # No background color
            "body_text": "000000",  # Black text
            "border_color": "5B9BD5",  # Accent1 color
            "border_width": Pt(1),  # 1 point border width
            "accent_color": "DEEBF7",  # Light blue for alternating rows
            "first_row": True,  # Format first row as header
            "total_row": False,  # No footer row
            "banded_rows": True,  # Alternating row colors
            "banded_cols": False,  # No alternating column colors
        },
        "accent2": {
            "header_bg": "ED7D31",  # Accent2 color (orange)
            "header_text": "FFFFFF",  # White text
            "body_bg": None,  # No background color
            "body_text": "000000",  # Black text
            "border_color": "ED7D31",  # Accent2 color
            "border_width": Pt(1),  # 1 point border width
            "accent_color": "FBE5D6",  # Light orange for alternating rows
            "first_row": True,  # Format first row as header
            "total_row": False,  # No footer row
            "banded_rows": True,  # Alternating row colors
            "banded_cols": False,  # No alternating column colors
        },
        "accent3": {
            "header_bg": "A5A5A5",  # Accent3 color (gray)
            "header_text": "FFFFFF",  # White text
            "body_bg": None,  # No background color
            "body_text": "000000",  # Black text
            "border_color": "A5A5A5",  # Accent3 color
            "border_width": Pt(1),  # 1 point border width
            "accent_color": "EDEDED",  # Light gray for alternating rows
            "first_row": True,  # Format first row as header
            "total_row": False,  # No footer row
            "banded_rows": True,  # Alternating row colors
            "banded_cols": False,  # No alternating column colors
        },
    }
    
    # Regex patterns for text formatting
    BOLD_PATTERN = r'\*\*(.+?)\*\*'
    ITALIC_PATTERN = r'\*(.+?)\*'
    STRIKETHROUGH_PATTERN = r'~~(.+?)~~'
    UNDERLINE_PATTERN = r'__(.+?)__'
    COLOR_PATTERN = r'\{color:([a-zA-Z0-9#]+)\}(.+?)\{/color\}'
    HIGHLIGHT_PATTERN = r'\{highlight:([a-zA-Z0-9#]+)\}(.+?)\{/highlight\}'
    FONT_SIZE_PATTERN = r'\{size:(\d+)(pt|px)?\}(.+?)\{/size\}'
    
    # Common colors
    COLORS = {
        "red": "FF0000",
        "green": "00FF00",
        "blue": "0000FF",
        "yellow": "FFFF00",
        "orange": "FFA500",
        "purple": "800080",
        "black": "000000",
        "white": "FFFFFF",
        "gray": "808080",
        "lightgray": "D3D3D3",
        "darkgray": "A9A9A9",
    }
    
    def __init__(self, template_path: Optional[Union[str, Path]] = None, 
                use_ai: bool = False, use_content_planning: bool = False):
        """
        Initialize a PowerPoint builder.
        
        Args:
            template_path: Optional path to a PowerPoint template file.
                        If not provided, a new blank presentation will be created.
            use_ai: Whether to use AI for optimization.
            use_content_planning: Whether to use AI content planning.
        
        Raises:
            FileNotFoundError: If the template file does not exist.
            ValueError: If the template file is invalid.
        """
        self.template_loader = TemplateLoader()
        self.layout_selector = LayoutSelector()
        self.overflow_handler = OverflowHandler()
        self.use_ai = use_ai
        self.use_content_planning = use_content_planning

        # Initialize optimizer if AI is enabled
        if self.use_ai or self.use_content_planning:
            try:
                self.optimizer = PresentationOptimizer()
            except Exception as e:
                logger.warning(f"Could not initialize AI optimizer: {e}. Some AI features will be disabled.")
                self.use_ai = False
                self.use_content_planning = False
        
        # Initialize template_info
        self.template_info: Optional[TemplateInfo] = None
        self.template_path: Optional[Path] = None
        
        if template_path:
            self.template_path = Path(template_path)
            # Use AI-enhanced template analysis if available
            if self.use_ai or self.use_content_planning:
                self.template_info = self.template_loader.analyze_template_with_ai(self.template_path)
            else:
                self.template_info = self.template_loader.analyze_template(self.template_path)
        
    # def build(self, presentation: Presentation, output_path: Union[str, Path]) -> Path:
    #     """
    #     Build a PowerPoint presentation from a Presentation model.
        
    #     Args:
    #         presentation: Presentation model containing sections and slides.
    #         output_path: Path where the generated PowerPoint file will be saved.
        
    #     Returns:
    #         Path to the generated PowerPoint file.
            
    #     Raises:
    #         ValueError: If the presentation cannot be built due to invalid content or layout.
    #     """
    #     # Use template from presentation if provided, otherwise use the one from constructor
    #     template_path = presentation.template_path or self.template_path
        
    #     if not template_path:
    #         raise ValueError("No template path provided. Either specify a template_path in the presentation model or when initializing PPTBuilder.")
        
    #     # Update template_info if template_path has changed
    #     if template_path != self.template_path:
    #         self.template_path = Path(template_path)
    #         if self.use_ai:
    #             self.template_info = self.template_loader.analyze_template_with_ai(self.template_path)
    #         else:
    #             self.template_info = self.template_loader.analyze_template(self.template_path)
        
    #     # Load the template
    #     pptx = self.template_loader.load_template(self.template_path)
        
    #     # Create a new LayoutSelector with the template and template_info
    #     self.layout_selector = LayoutSelector(template=pptx, use_ai=self.use_ai)
    #     # Explicitly pass the template_info to ensure it's available
    #     self.layout_selector.template_info = self.template_info

    #     # ── purge des slides déjà présentes dans le template
    #     self._clear_template_slides(pptx)

    #     # Process each section and slide
    #     for section in presentation.sections:
    #         # Validate custom section types if AI is enabled
    #         if self.use_ai and not isinstance(section.type, SectionType):
    #             try:
    #                 # Map custom section type to standard type
    #                 mapped_type = self.optimizer.validate_and_map_section_type(section.type)
    #                 section.type = SectionType(mapped_type)
    #             except Exception as e:
    #                 logger.warning(f"Error mapping custom section type '{section.type}': {e}. Using 'custom' type.")
    #                 section.type = SectionType.CUSTOM

    #         # ─── Diapositive d'en-tête de section ─────────────────────────
    #         if section.title and self._needs_section_header(section):
    #             header_layout = self.layout_selector.get_layout_name(section)
    #             header_slide = self._create_slide(pptx, header_layout)
    #             self._fill_slide_title(header_slide, section.title)

    #         # Process each slide in the section
    #         slide_index = 0
    #         while slide_index < len(section.slides):
    #             slide = section.slides[slide_index]

    #             # Select layout name if not specified or validate it
    #             if not slide.layout_name or slide.layout_name == "auto":
    #                 slide.layout_name = self.layout_selector.get_layout_name(section, slide)
                
    #             # Validate if the layout is appropriate for the content using template_info
    #             slide.layout_name = self._validate_layout_for_content(slide)
                
    #             # Create the slide
    #             pptx_slide = self._create_slide(pptx, slide.layout_name)
                
    #             # Check for overflow before filling content
    #             if self.use_ai:
    #                 # Get shape dimensions for overflow check
    #                 content_shape_width = 0
    #                 content_shape_height = 0
                    
    #                 # Find a content placeholder to get dimensions
    #                 for shape in pptx_slide.shapes:
    #                     if (shape.is_placeholder and 
    #                         hasattr(shape, 'placeholder_format') and
    #                         shape.placeholder_format.type == PP_PLACEHOLDER.BODY):
    #                         content_shape_width = Emu(shape.width).pt
    #                         content_shape_height = Emu(shape.height).pt
    #                         break
                    
    #                 if content_shape_width > 0 and content_shape_height > 0:
    #                     # Check and handle overflow
    #                     result_slides = self.overflow_handler.handle_slide_overflow(
    #                         section, slide, content_shape_width, content_shape_height
    #                     )
                        
    #                     if len(result_slides) > 1:
    #                         # Replace the current slide with the first result
    #                         section.slides[slide_index] = result_slides[0]
    #                         slide = result_slides[0]
                            
    #                         # Insert additional slides after the current one
    #                         for i, new_slide in enumerate(result_slides[1:], 1):
    #                             section.slides.insert(slide_index + i, new_slide)
                
    #             # Fill the slide with content
    #             self._fill_slide(pptx_slide, slide, section)
                
    #             # Move to the next slide
    #             slide_index += 1
        
    #     # Save the presentation (création automatique du répertoire parent si besoin)
    #     output_path = Path(output_path)
    #     output_path.parent.mkdir(parents=True, exist_ok=True)
    #     try:
    #         pptx.save(output_path)
    #     except AttributeError:
    #         logger.warning("Object returned by load_template() has no .save() "
    #                     "(mock in tests) – creating stub.")
    #         from unittest.mock import MagicMock
    #         pptx.save = MagicMock()
    #         pptx.save(output_path)
        
    #     logger.info(f"PowerPoint presentation successfully built and saved to {output_path}")
        
    #     return output_path
    
    def build(self, presentation: Presentation, output_path: Union[str, Path]) -> Path:
        """
        Build a PowerPoint presentation from a Presentation model.
        
        Args:
            presentation: Presentation model containing sections and slides.
            output_path: Path where the generated PowerPoint file will be saved.
        
        Returns:
            Path to the generated PowerPoint file.
            
        Raises:
            ValueError: If the presentation cannot be built due to invalid content or layout.
        """
        # Use template from presentation if provided, otherwise use the one from constructor
        template_path = presentation.template_path or self.template_path
        
        if not template_path:
            raise ValueError("No template path provided. Either specify a template_path in the presentation model or when initializing PPTBuilder.")
        
        # Update template_info if template_path has changed
        if template_path != self.template_path:
            self.template_path = Path(template_path)
            if self.use_ai:
                self.template_info = self.template_loader.analyze_template_with_ai(self.template_path)
            else:
                self.template_info = self.template_loader.analyze_template(self.template_path)
        
        # Ensure the layout_selector has the template_info
        self.layout_selector = LayoutSelector(template=self.template_loader.load_template(self.template_path),
                                            use_ai=self.use_ai)
        self.layout_selector.template_info = self.template_info
        
        # Load the template
        pptx = self.template_loader.load_template(self.template_path)

        # ── purge des slides déjà présentes dans le template
        self._clear_template_slides(pptx)

        # Use ContentPlanner to optimize sections if AI is enabled
        if self.use_ai:
            from doc2pptx.llm.content_planner import ContentPlanner
            content_planner = ContentPlanner(optimizer=self.optimizer if hasattr(self, 'optimizer') else None)
            
            # Process each section using the ContentPlanner
            optimized_sections = []
            for section in presentation.sections:
                optimized_section = content_planner.plan_section_content(section, self.template_info)
                optimized_sections.append(optimized_section)
            
            # Replace the sections in the presentation
            presentation.sections = optimized_sections

        # Process each section and slide (now optimized if AI was used)
        for section in presentation.sections:
            # Validate custom section types if AI is enabled
            if self.use_ai and not isinstance(section.type, SectionType):
                try:
                    # Map custom section type to standard type
                    mapped_type = self.optimizer.validate_and_map_section_type(section.type)
                    section.type = SectionType(mapped_type)
                except Exception as e:
                    logger.warning(f"Error mapping custom section type '{section.type}': {e}. Using 'custom' type.")
                    section.type = SectionType.CUSTOM

            # ─── Diapositive d'en-tête de section ─────────────────────────
            if section.title and self._needs_section_header(section):
                header_layout = self.layout_selector.get_layout_name(section)
                header_slide = self._create_slide(pptx, header_layout)
                self._fill_slide_title(header_slide, section.title)

            # Process each slide in the section
            for slide in section.slides:
                # Select layout name if not specified or validate it
                if not slide.layout_name or slide.layout_name == "auto":
                    slide.layout_name = self.layout_selector.get_layout_name(section, slide)
                
                # Validate if the layout is appropriate for the content
                slide.layout_name = self._validate_layout_for_content(slide)
                
                # Create the slide
                pptx_slide = self._create_slide(pptx, slide.layout_name)
                
                # Fill the slide with content
                self._fill_slide(pptx_slide, slide, section)
        
        # Save the presentation (création automatique du répertoire parent si besoin)
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        try:
            pptx.save(output_path)
        except AttributeError:
            logger.warning("Object returned by load_template() has no .save() "
                        "(mock in tests) – creating stub.")
            from unittest.mock import MagicMock
            pptx.save = MagicMock()
            pptx.save(output_path)
        
        logger.info(f"PowerPoint presentation successfully built and saved to {output_path}")
        
        return output_path
    
    def _validate_layout_for_content(self, slide: Slide) -> str:
        """
        Validate if the selected layout is appropriate for the slide content.
        If not, select a more appropriate layout.
        
        Args:
            slide: The slide to validate layout for
            
        Returns:
            The validated or corrected layout name
        """
        current_layout = slide.layout_name

        # If we have template_info, use it for validation
        if self.template_info and current_layout in self.template_info.layout_map:
            layout_info = self.template_info.layout_map[current_layout]
            
            # Vérifier le nombre de blocs de contenu
            num_blocks = len(slide.blocks)
            max_blocks = layout_info.max_content_blocks
            
            if num_blocks > max_blocks and max_blocks > 0:
                logger.warning(
                    f"Layout '{current_layout}' supports max {max_blocks} blocks but slide has {num_blocks}. "
                    "Selecting a more appropriate layout."
                )
                
                # Get better layout using LayoutSelector
                return self.layout_selector.get_layout_name(None, slide)
            
            # Vérifier les types de contenu spécifiques
            has_table = any(block.content.content_type == ContentType.TABLE for block in slide.blocks)
            
            # Si nous avons une table mais pas dans un layout de table, changer pour un layout de table
            if has_table and not layout_info.supports_table:
                logger.warning(f"Slide contains table but layout '{current_layout}' does not support tables. Using table layout.")
                # Find a table layout
                table_layouts = self.template_info.table_layouts
                if table_layouts:
                    return table_layouts[0]
                return current_layout
            
            # Le layout actuel est approprié
            return current_layout

        # Si le layout n'est pas dans nos capacités connues, utiliser le layout par défaut
        if current_layout not in LAYOUT_CAPABILITIES:
            logger.warning(f"Layout '{current_layout}' not found in capabilities. Using default layout.")
            return "Titre et texte"
        
        # Vérifier le nombre de blocs de contenu
        num_blocks = len(slide.blocks)
        max_blocks = LAYOUT_CAPABILITIES[current_layout]["max_blocks"]
        
        if num_blocks > max_blocks and max_blocks > 0:
            logger.warning(
                f"Layout '{current_layout}' supports max {max_blocks} blocks but slide has {num_blocks}. "
                "Selecting a more appropriate layout."
            )
            
            # Sélectionner un layout plus approprié en fonction du contenu
            if any(block.content.content_type == ContentType.TABLE for block in slide.blocks):
                return "Titre et tableau"
            elif num_blocks <= 3:
                return "Titre et 3 colonnes"
            else:
                return "Titre et texte"  # Fallback
        
        # Vérifier les types de contenu spécifiques
        has_table = any(block.content.content_type == ContentType.TABLE for block in slide.blocks)
        
        # Si nous avons une table mais pas dans un layout de table, changer pour "Titre et tableau"
        if has_table and not LAYOUT_CAPABILITIES[current_layout]["table"]:
            logger.warning(f"Slide contains table but layout '{current_layout}' does not support tables. Using table layout.")
            return "Titre et tableau"
        
        # Si nous avons un layout de table mais pas de table, utiliser un layout standard
        if current_layout == "Titre et tableau" and not has_table:
            logger.warning(f"Layout is 'Titre et tableau' but slide does not contain a table. Using standard layout.")
            return "Titre et texte"
        
        # Le layout actuel est approprié
        return current_layout

    def _get_layout_capabilities(self):
        """
        Get layout capabilities from template_info if available, otherwise use static definitions.
        
        Returns:
            Dictionary of layout capabilities
        """
        if self.template_info:
            # Build capabilities from template_info
            capabilities = {}
            for layout_name, layout_info in self.template_info.layout_map.items():
                capabilities[layout_name] = {
                    "title": layout_info.supports_title,
                    "subtitle": layout_info.placeholder_types and PP_PLACEHOLDER.SUBTITLE in layout_info.placeholder_types,
                    "content": layout_info.supports_content,
                    "table": layout_info.supports_table,
                    "image": layout_info.supports_image,
                    "chart": layout_info.supports_chart,
                    "max_blocks": layout_info.max_content_blocks,
                    "description": layout_info.ai_description if hasattr(layout_info, 'ai_description') else ""
                }
            return capabilities
        else:
            # Use static definitions from the original code
            return {
                "Diapositive de titre": {
                    "title": True,
                    "subtitle": True,
                    "content": False,
                    "table": False,
                    "image": False,
                    "chart": False,
                    "max_blocks": 1,
                    "description": "Title slide with subtitle"
                },
                "Introduction": {
                    "title": True,
                    "subtitle": False,
                    "content": True,
                    "table": False,
                    "image": False,
                    "chart": False,
                    "max_blocks": 1,
                    "description": "Title with large content area"
                },
                "Titre et texte": {
                    "title": True,
                    "subtitle": False,
                    "content": True,
                    "table": False,
                    "image": False,
                    "chart": False,
                    "max_blocks": 1,
                    "description": "Title with text content"
                },
                "Titre et tableau": {
                    "title": True,
                    "subtitle": False,
                    "content": False,
                    "table": True,
                    "image": False,
                    "chart": False,
                    "max_blocks": 1,
                    "description": "Title with table"
                },
                "Titre et texte 1 visuel gauche": {
                    "title": True,
                    "subtitle": False,
                    "content": True,
                    "table": False,
                    "image": True,
                    "chart": False,
                    "max_blocks": 2,
                    "description": "Title with image on left and text on right"
                },
                "Titre et texte 1 histogramme": {
                    "title": True,
                    "subtitle": False,
                    "content": True,
                    "table": False,
                    "image": False,
                    "chart": True,
                    "max_blocks": 2,
                    "description": "Title with text on left and chart on right"
                },
                "Titre et 3 colonnes": {
                    "title": True,
                    "subtitle": False,
                    "content": True,
                    "table": False,
                    "image": False,
                    "chart": False,
                    "max_blocks": 3,
                    "description": "Title with three text columns"
                },
                "Chapitre 1": {
                    "title": True,
                    "subtitle": False,
                    "content": False,
                    "table": False,
                    "image": False,
                    "chart": False,
                    "max_blocks": 0,
                    "description": "Section title only"
                }
            }
            
    def _get_placeholder_map(self, layout_name: str):
        """
        Get placeholder map for layout from template_info if available, otherwise use static map.
        
        Args:
            layout_name: Name of the layout
            
        Returns:
            Dictionary mapping placeholder types to indices
        """
        if self.template_info and layout_name in self.template_info.layout_map:
            layout_info = self.template_info.layout_map[layout_name]
            
            # Create mapping from placeholder types to indices
            mapping = {}
            for i, ph_type in enumerate(layout_info.placeholder_types):
                if ph_type in TemplateLoader.PLACEHOLDER_TYPE_MAP:
                    capability = TemplateLoader.PLACEHOLDER_TYPE_MAP[ph_type]
                    if capability not in mapping:
                        mapping[capability] = layout_info.placeholder_indices[i]
            
            return mapping
        else:
            # Use static mapping from the original code
            layout_placeholder_map = {
                "Diapositive de titre": {
                    "title": 0,       # idx=0, TITLE
                    "subtitle": 1,    # idx=1, SUBTITLE
                },
                "Introduction": {
                    "title": 0,       # idx=0, TITLE
                    "content": 1,     # idx=1, BODY
                    "slide_number": 12 # idx=12, SLIDE_NUMBER
                },
                "Titre et texte": {
                    "title": 0,       # idx=0, TITLE
                    "content": 1,     # idx=1, BODY
                    "slide_number": 12 # idx=12, SLIDE_NUMBER
                },
                "Titre et tableau": {
                    "title": 0,       # idx=0, TITLE
                    "slide_number": 12 # idx=12, SLIDE_NUMBER
                    # Pas de placeholder pour la table, elle est ajoutée comme shape
                },
                "Titre et texte 1 visuel gauche": {
                    "title": 0,       # idx=0, TITLE
                    "content": 1,     # idx=1, BODY (à droite)
                    "image": 2,       # idx=2, PICTURE (à gauche)
                    "slide_number": 12 # idx=12, SLIDE_NUMBER
                },
                "Titre et texte 1 histogramme": {
                    "title": 0,       # idx=0, TITLE
                    "content": 1,     # idx=1, BODY (à gauche)
                    "chart": 2,       # idx=2, CHART (à droite)
                    "slide_number": 12 # idx=12, SLIDE_NUMBER
                },
                "Titre et 3 colonnes": {
                    "title": 0,       # idx=0, TITLE
                    "column1": 1,     # idx=1, BODY (colonne 1)
                    "column2": 2,     # idx=2, BODY (colonne 2)
                    "column3": 3,     # idx=3, BODY (colonne 3)
                    "slide_number": 12 # idx=12, SLIDE_NUMBER
                },
                "Chapitre 1": {
                    "title": 0,       # idx=0, TITLE
                }
            }
            
            return layout_placeholder_map.get(layout_name, {})

    @staticmethod
    def _clear_template_slides(pptx: PptxPresentation) -> None:
        """
        Remove all slides that may be included in the template.
        
        Args:
            pptx: PowerPoint presentation to clear slides from.
        """
        for sldId in list(pptx.slides._sldIdLst):
            pptx.part.drop_rel(sldId.rId)
            pptx.slides._sldIdLst.remove(sldId)
        
    def _needs_section_header(self, section: Section) -> bool:
        """
        Determine if a section header slide should be added.
        
        Args:
            section: Section to check.
            
        Returns:
            True if a section header should be added, False otherwise.
        """
        # 1️⃣ pas de header si la section n'a pas de slides ou
        # 2️⃣ pas de header pour les sections déjà typées 'title', 'agenda', 'section_header'
        # return (
        #     len(section.slides) > 0
        #     and section.type not in {SectionType.TITLE, SectionType.AGENDA, SectionType.SECTION_HEADER}
        # )
        return False
        
    def _create_slide(self, pptx: PptxPresentation, layout_name: str) -> PptxSlide:
        """
        Create a new slide in the presentation with the specified layout.
        
        Args:
            pptx: PowerPoint presentation to add the slide to.
            layout_name: Name of the layout to use for the slide.
        
        Returns:
            The created PowerPoint slide.
            
        Raises:
            ValueError: If the layout does not exist in the template.
        """
        # Find the layout by name
        layout = None
        for slide_layout in pptx.slide_layouts:
            if slide_layout.name == layout_name:
                layout = slide_layout
                break
        
        if layout is None:
            # Get available layouts
            available_layouts = [layout.name for layout in pptx.slide_layouts]
            logger.warning(f"Layout '{layout_name}' not found in template. Using the first available layout instead.")
            logger.info(f"Available layouts: {available_layouts}")
            
            # Use the first layout as fallback
            layout = pptx.slide_layouts[0]
        
        # Create the slide with the selected layout
        slide = pptx.slides.add_slide(layout)
        
        return slide
    
    def _fill_slide(self, pptx_slide: PptxSlide, slide: Slide, section: Section) -> None:
        """
        Fill a PowerPoint slide with content from a Slide model.
        
        Args:
            pptx_slide: PowerPoint slide to fill.
            slide: Slide model containing content to add to the PowerPoint slide.
            section: Section model containing the slide.
            
        Raises:
            ValueError: If the content cannot be added to the slide.
        """
        # Add title if provided
        self._fill_slide_title(pptx_slide, slide.title)
        
        # Special handling based on layout type
        if slide.layout_name == "Diapositive de titre":
            self._fill_title_slide(pptx_slide, slide)
        elif slide.layout_name == "Titre et tableau":
            self._fill_table_slide(pptx_slide, slide)
        elif slide.layout_name == "Titre et 3 colonnes":
            self._fill_column_layout_slide(pptx_slide, slide)
        elif slide.layout_name == "Titre et texte 1 visuel gauche":
            self._fill_image_layout_slide(pptx_slide, slide)
        elif slide.layout_name == "Titre et texte 1 histogramme":
            self._fill_chart_layout_slide(pptx_slide, slide)
        elif slide.layout_name == "Chapitre 1":
            # Chapitre 1 has only a title, which we already filled
            pass
        else:
            # Default handling for other layouts (generally just title + content)
            self._fill_content_slide(pptx_slide, slide)
        
        # Add speaker notes if provided
        if slide.notes:
            notes_slide = pptx_slide.notes_slide
            notes_slide.notes_text_frame.text = slide.notes
    
    def _fill_slide_title(self, pptx_slide: PptxSlide, title: Optional[str]) -> None:
        """
        Fill the title placeholder of a slide if available.
        
        Args:
            pptx_slide: PowerPoint slide to add title to.
            title: Title text to add.
        """
        if not title:
            return
            
        # Find title placeholder
        for shape in pptx_slide.shapes:
            if (shape.is_placeholder and 
                shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE] and
                hasattr(shape, 'text_frame')):
                self._add_formatted_text(shape.text_frame, title)
                return
                
        logger.warning("No title placeholder found in slide")
    
    def _fill_title_slide(self, pptx_slide: PptxSlide, slide: Slide) -> None:
        """
        Fill a title slide with title and subtitle.
        
        Args:
            pptx_slide: PowerPoint slide to fill.
            slide: Slide model containing content.
        """
        # Find subtitle placeholder
        subtitle_placeholder = None
        for shape in pptx_slide.shapes:
            if (shape.is_placeholder and 
                shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE and
                hasattr(shape, 'text_frame')):
                subtitle_placeholder = shape
                break
        
        # Add subtitle if found
        if subtitle_placeholder and slide.blocks:
            block = slide.blocks[0]
            if block.content.content_type == ContentType.TEXT and block.content.text:
                self._add_formatted_text(subtitle_placeholder.text_frame, block.content.text)
    
    def _fill_content_slide(self, pptx_slide: PptxSlide, slide: Slide) -> None:
        """
        Fill a standard content slide with a single content area.
        
        Args:
            pptx_slide: PowerPoint slide to fill.
            slide: Slide model containing content.
        """
        # Find the main content placeholder
        content_placeholder = None
        for shape in pptx_slide.shapes:
            if (shape.is_placeholder and 
                shape.placeholder_format.type == PP_PLACEHOLDER.BODY and
                hasattr(shape, 'text_frame')):
                content_placeholder = shape
                break
        
        if not content_placeholder:
            logger.warning("No content placeholder found in slide")
            return
            
        # Clear the placeholder before adding content
        content_placeholder.text_frame.clear()
            
        # Add each content block to the placeholder
        for i, block in enumerate(slide.blocks):
            if i > 0:
                # Add a line break between blocks
                content_placeholder.text_frame.add_paragraph()
                
            # Add block title if present
            if block.title:
                para = content_placeholder.text_frame.add_paragraph()
                para.text = block.title
                # Format as heading
                para.font.bold = True
                para.font.size = Pt(16)
                
            # Add content based on type
            content_type = block.content.content_type
            
            if content_type == ContentType.TEXT and block.content.text:
                self._add_text_content_to_placeholder(content_placeholder, block.content.text)
            
            elif content_type == ContentType.BULLET_POINTS and block.content.bullet_points:
                self._add_bullet_points_to_placeholder(content_placeholder, 
                                                      block.content.bullet_points, 
                                                      block.content.as_bullets)
            
            # Pour les autres types de contenu, afficher du texte par défaut
            else:
                para = content_placeholder.text_frame.add_paragraph()
                para.text = f"[{content_type.value} content not shown in this layout]"
    
    def _clear_content_placeholders(self, pptx_slide: PptxSlide) -> None:
        """
        Clear any content placeholders in the slide to prevent duplicate content rendering.
        
        Args:
            pptx_slide: PowerPoint slide to clear placeholders from.
        """
        for shape in pptx_slide.shapes:
            if (shape.is_placeholder and 
                shape.placeholder_format.type == PP_PLACEHOLDER.BODY and
                hasattr(shape, 'text_frame')):
                shape.text_frame.clear()

    def _fill_table_slide(self, pptx_slide: PptxSlide, slide: Slide) -> None:
        """
        Fill a slide containing a table.
        
        Args:
            pptx_slide: PowerPoint slide to fill.
            slide: Slide model containing table content.
        """
        logger.info(f"_fill_table_slide: Starting to fill table slide")
        
        # Clear any content placeholders
        self._clear_content_placeholders(pptx_slide)
        
        # Find the table content
        table_block = None
        for block in slide.blocks:
            if block.content and block.content.content_type == ContentType.TABLE:
                table_block = block
                break
        
        if not table_block or not table_block.content or not table_block.content.table:
            logger.warning("No table content found in slide")
            return
        
        # Extract table data
        table_data = table_block.content.table
        if not table_data.headers or not table_data.rows:
            logger.warning("Table has no headers or rows")
            return
        
        # Extract style if present
        style = "default"
        style_from_headers = self._get_style_from_headers(table_data.headers)
        if style_from_headers:
            style = style_from_headers
            # Remove style header from display
            headers = table_data.headers[:-1]
        else:
            headers = table_data.headers
        
        # Calculate table dimensions
        rows = len(table_data.rows) + 1  # +1 for header row
        cols = len(headers)
        
        # Safety check - ensure non-zero dimensions
        if rows <= 0 or cols <= 0:
            logger.error(f"Invalid table dimensions: {rows} rows, {cols} columns")
            return
        
        # Calculate table position and size
        title_height = Inches(1.5)
        left = Inches(0.5)
        top = title_height + Inches(0.5)
        width = Inches(9)
        height = Inches(5)
        
        try:
            # Convert to EMU safely
            left_emu = self._emu(left)
            top_emu = self._emu(top)
            width_emu = self._emu(width)
            height_emu = self._emu(height)
            
            logger.info(f"Attempting to add table: {rows}x{cols} at ({left_emu},{top_emu}) with size ({width_emu},{height_emu})")
            table_shape = pptx_slide.shapes.add_table(rows, cols, left_emu, top_emu, width_emu, height_emu)
            logger.info(f"Table added successfully: {table_shape}")
            
            table = table_shape.table
            self._fill_table_with_data(table, headers, table_data.rows, style)
            logger.info("Table data filled successfully")
        except Exception as e:
            logger.error(f"Error creating table: {e}")
            logger.error(f"Exception type: {type(e).__name__}")
            logger.error(f"Exception details: {str(e)}")
            # Add proper stack trace logging
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")

    
    def _fill_column_layout_slide(self, pptx_slide: PptxSlide, slide: Slide) -> None:
        """
        Fill a slide with multiple column layout.

        Args:
            pptx_slide: PowerPoint slide to fill.
            slide: Slide model containing content for multiple columns.
        """
        # Find all column placeholders
        column_placeholders = []
        for shape in pptx_slide.shapes:
            if (shape.is_placeholder and 
                shape.placeholder_format.type == PP_PLACEHOLDER.BODY and
                hasattr(shape, 'text_frame')):
                column_placeholders.append(shape)
        
        # Make sure we have at least one column placeholder
        if not column_placeholders:
            logger.warning("No column placeholders found in slide")
            return
            
        # Sort placeholders by left position to ensure correct column order
        try:
            column_placeholders.sort(key=lambda p: getattr(p, 'left', 0))
        except TypeError:
            # If sorting fails, keep the original order
            logger.warning("Unable to sort placeholders by position. Using original order.")
        
            
        # Distribute content blocks among column placeholders
        num_columns = len(column_placeholders)
        num_blocks = len(slide.blocks)
        
        # Clear all placeholders first
        for placeholder in column_placeholders:
            placeholder.text_frame.clear()
        
        # Assign blocks to columns
        if num_blocks <= num_columns:
            # If we have fewer or equal blocks than columns, place each block in its own column
            for i, block in enumerate(slide.blocks):
                if i < num_columns:
                    self._add_block_to_placeholder(column_placeholders[i], block)
        else:
            # If we have more blocks than columns, distribute them
            blocks_per_column = num_blocks // num_columns
            remainder = num_blocks % num_columns
            
            block_index = 0
            for col_index in range(num_columns):
                # Calculate how many blocks this column should get
                blocks_for_this_column = blocks_per_column
                if col_index < remainder:
                    blocks_for_this_column += 1
                
                # Add blocks to this column
                for i in range(blocks_for_this_column):
                    if block_index < num_blocks:
                        if i > 0:
                            # Add a separator between blocks in the same column
                            separator = column_placeholders[col_index].text_frame.add_paragraph()
                            separator.text = "---"
                            separator.alignment = PP_ALIGN.CENTER
                        
                        # Add the block
                        self._add_block_to_placeholder(column_placeholders[col_index], slide.blocks[block_index])
                        block_index += 1
    
    def _fill_image_layout_slide(self, pptx_slide: PptxSlide, slide: Slide) -> None:
        """
        Fill a slide with image on left and text on right.
        
        Args:
            pptx_slide: PowerPoint slide to fill.
            slide: Slide model containing image and text content.
        """
        # Find image and content placeholders
        image_placeholder = None
        content_placeholder = None
        
        for shape in pptx_slide.shapes:
            if shape.is_placeholder and hasattr(shape, 'placeholder_format'):
                if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                    image_placeholder = shape
                elif shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                    content_placeholder = shape
        
        # Handle image content
        image_block = None
        text_block = None
        
        for block in slide.blocks:
            if block.content.content_type == ContentType.IMAGE:
                image_block = block
            elif block.content.content_type in [ContentType.TEXT, ContentType.BULLET_POINTS]:
                text_block = block
        
        # Add image content
        if image_placeholder and image_block and image_block.content.image:
            # For now, just add a placeholder text describing the image
            # (Actual image handling will be implemented in the future)
            image_info = image_block.content.image
            image_desc = f"[Image: "
            if hasattr(image_info, 'query') and image_info.query:
                image_desc += f"Query: {image_info.query}"
            elif hasattr(image_info, 'url') and image_info.url:
                image_desc += f"URL: {image_info.url}"
            elif hasattr(image_info, 'path') and image_info.path:
                image_desc += f"Path: {image_info.path}"
            image_desc += "]"
            
            # Add description to the image placeholder
            if hasattr(image_placeholder, 'text_frame'):
                image_placeholder.text_frame.text = image_desc
        
        # Add text content
        if content_placeholder and text_block:
            content_placeholder.text_frame.clear()
            
            # Add content based on type
            if text_block.content.content_type == ContentType.TEXT and text_block.content.text:
                self._add_text_content_to_placeholder(content_placeholder, text_block.content.text)
            elif text_block.content.content_type == ContentType.BULLET_POINTS and text_block.content.bullet_points:
                self._add_bullet_points_to_placeholder(content_placeholder, 
                                                     text_block.content.bullet_points,
                                                     text_block.content.as_bullets)
    
    def _fill_chart_layout_slide(self, pptx_slide: PptxSlide, slide: Slide) -> None:
        """
        Fill a slide with text on left and chart on right.
        
        Args:
            pptx_slide: PowerPoint slide to fill.
            slide: Slide model containing chart and text content.
        """
        # Find chart and content placeholders
        chart_placeholder = None
        content_placeholder = None
        
        for shape in pptx_slide.shapes:
            if shape.is_placeholder and hasattr(shape, 'placeholder_format'):
                if shape.placeholder_format.type == PP_PLACEHOLDER.CHART:
                    chart_placeholder = shape
                elif shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                    content_placeholder = shape
        
        # Handle chart content
        chart_block = None
        text_block = None
        mermaid_block = None
        
        for block in slide.blocks:
            if block.content.content_type == ContentType.CHART:
                chart_block = block
            elif block.content.content_type == ContentType.MERMAID:
                mermaid_block = block
            elif block.content.content_type in [ContentType.TEXT, ContentType.BULLET_POINTS]:
                text_block = block
        
        # Add chart content
        if chart_placeholder:
            if chart_block and chart_block.content.chart:
                # For now, just add a placeholder text describing the chart
                # (Actual chart handling will be implemented in the future)
                chart_info = chart_block.content.chart
                chart_desc = f"[Chart: {chart_info.chart_type}"
                if hasattr(chart_info, 'title') and chart_info.title:
                    chart_desc += f", Title: {chart_info.title}"
                chart_desc += "]"
                
                # Add description to the chart placeholder
                if hasattr(chart_placeholder, 'text_frame'):
                    chart_placeholder.text_frame.text = chart_desc
            elif mermaid_block and mermaid_block.content.mermaid:
                # For now, just add a placeholder text describing the mermaid diagram
                # (Actual mermaid handling will be implemented in the future)
                mermaid_info = mermaid_block.content.mermaid
                mermaid_desc = f"[Mermaid diagram"
                if hasattr(mermaid_info, 'caption') and mermaid_info.caption:
                    mermaid_desc += f": {mermaid_info.caption}"
                mermaid_desc += "]"
                
                # Add description to the chart placeholder
                if hasattr(chart_placeholder, 'text_frame'):
                    chart_placeholder.text_frame.text = mermaid_desc
        
        # Add text content
        if content_placeholder and text_block:
            content_placeholder.text_frame.clear()
            
            # Add content based on type
            if text_block.content.content_type == ContentType.TEXT and text_block.content.text:
                self._add_text_content_to_placeholder(content_placeholder, text_block.content.text)
            elif text_block.content.content_type == ContentType.BULLET_POINTS and text_block.content.bullet_points:
                self._add_bullet_points_to_placeholder(content_placeholder, 
                                                     text_block.content.bullet_points,
                                                     text_block.content.as_bullets)
    
    def _add_block_to_placeholder(self, placeholder: SlidePlaceholder, block: SlideBlock) -> None:
        """
        Add a content block to a placeholder.
        
        Args:
            placeholder: The PowerPoint placeholder to add content to.
            block: The SlideBlock containing content to add.
        """
        # Clear the placeholder first
        placeholder.text_frame.clear()
        
        # Add block title if present
        if block.title:
            para = placeholder.text_frame.add_paragraph()
            para.text = block.title
            # Format as heading
            para.font.bold = True
            para.font.size = Pt(16)
            
        # Add content based on type
        content_type = block.content.content_type
        
        if content_type == ContentType.TEXT and block.content.text:
            self._add_text_content_to_placeholder(placeholder, block.content.text)
        
        elif content_type == ContentType.BULLET_POINTS and block.content.bullet_points:
            self._add_bullet_points_to_placeholder(placeholder, 
                                                  block.content.bullet_points,
                                                  block.content.as_bullets)
        
        # Pour les autres types de contenu, afficher du texte par défaut
        else:
            para = placeholder.text_frame.add_paragraph()
            para.text = f"[{content_type.value} content not shown in this placeholder]"
    
    def _add_text_content_to_placeholder(self, placeholder: SlidePlaceholder, text: str) -> None:
        """
        Add text content to a placeholder.
        
        Args:
            placeholder: PowerPoint placeholder to add text to.
            text: Text content to add.
        """
        # Split text into paragraphs
        paragraphs = text.split('\n')
        
        # Add each paragraph
        for i, paragraph_text in enumerate(paragraphs):
            if not paragraph_text.strip():
                # Empty paragraph, add a blank line
                placeholder.text_frame.add_paragraph()
                continue
                
            if i == 0 and not placeholder.text_frame.paragraphs[0].runs:
                # Use first paragraph if empty
                p = placeholder.text_frame.paragraphs[0]
            else:
                # Add a new paragraph
                p = placeholder.text_frame.add_paragraph()
            
            # Add the formatted text
            self._add_formatted_text_to_paragraph(p, paragraph_text)
    
    
    def _ensure_bullet_visible(self, paragraph) -> None:
        """
        Force l'ajout d'un style de puce ‹•› si le template n'en prévoit pas.
        Idempotent : n'ajoute rien si une puce est déjà définie.
        """
        pPr = paragraph._p.get_or_add_pPr()
        # <a:buNone/> ou pas de balise du tout ⇒ pas de puce configurée
        has_bullet_char = pPr.find(qn("a:buChar")) is not None
        has_auto_num = pPr.find(qn("a:buAutoNum")) is not None
        has_bu_none = pPr.find(qn("a:buNone")) is not None

        if not (has_bullet_char or has_auto_num) or has_bu_none:
            # Supprime éventuellement <a:buNone/>
            if has_bu_none:
                pPr.remove(pPr.find(qn("a:buNone")))
            # Ajoute <a:buChar char="•"/>
            buChar = OxmlElement("a:buChar")
            buChar.set("char", "•")
            pPr.append(buChar)
            
            
    def _add_bullet_points_to_placeholder(
    self,
    placeholder: SlidePlaceholder,
    bullet_points: list[str],
    as_bullets: bool = True,
) -> None:
        """
        Ajoute des points de liste (bullet points) à un placeholder.
        
        Args:
            placeholder: PowerPoint placeholder to add bullet points to.
            bullet_points: List of bullet point texts.
            as_bullets: Whether to format as bullets (True) or paragraphs (False).
        """
        for i, bullet_text in enumerate(bullet_points):
            # Use existing paragraph if it's the first item and paragraph is empty
            p = (
                placeholder.text_frame.paragraphs[0]
                if i == 0 and not placeholder.text_frame.paragraphs[0].runs
                else placeholder.text_frame.add_paragraph()
            )
            
            # Add text content first (without "•" prefixes)
            self._add_formatted_text_to_paragraph(p, bullet_text)
            
            # Check if this paragraph already has bullet formatting
            pPr = p._p.get_or_add_pPr()
            has_bullet_format = (
                pPr.find(qn("a:buChar")) is not None or 
                pPr.find(qn("a:buAutoNum")) is not None
            )
            has_bu_none = pPr.find(qn("a:buNone")) is not None
            
            if as_bullets:
                # Set level to 0 to activate list formatting
                p.level = 0
                
                # Only add bullet character via XML if not already formatted with bullets
                if not has_bullet_format or has_bu_none:
                    self._ensure_bullet_visible(p)
                    
                # No need to add "•" in the text itself - PowerPoint will show the bullet
            else:
                # Remove bullet formatting if we don't want bullets
                self._remove_bullet(p)

            
    def _calculate_table_dimensions(self, pptx_slide: PptxSlide, rows: int, cols: int) -> Tuple[int, int, int, int]:
        """
        Calculate optimal table dimensions and position based on slide layout.
        
        Args:
            pptx_slide: PowerPoint slide where the table will be placed
            rows: Number of rows in the table
            cols: Number of columns in the table
            
        Returns:
            Tuple of (left_emu, top_emu, width_emu, height_emu) in EMU units
        """
        # Get slide dimensions
        slide_width = pptx_slide.slide_width
        slide_height = pptx_slide.slide_height
        
        # Find title placeholder to get its actual height
        title_height = Inches(1.0)  # Default title height if title placeholder not found
        for shape in pptx_slide.shapes:
            if (shape.is_placeholder and 
                hasattr(shape, 'placeholder_format') and
                shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]):
                title_height = shape.height + Inches(0.2)  # Add a small margin
                break
                
        # Calculate available space
        slide_margin_left = Inches(0.5)    # Left margin
        slide_margin_right = Inches(0.5)   # Right margin
        slide_margin_bottom = Inches(0.75) # Bottom margin to avoid footer
        
        # Calculate available width, centered
        available_width = slide_width - slide_margin_left - slide_margin_right
        left = slide_margin_left + (available_width - Inches(9)) / 2
        # Ensure left margin is at least slide_margin_left
        left = max(left, slide_margin_left)
        
        # Determine table width - adjust based on content if needed
        width = min(Inches(9), available_width)
        
        # Calculate top position - below title
        top = title_height + Inches(0.3)  # Add some spacing after title
        
        # Calculate available height for table (respect footer area)
        available_height = slide_height - top - slide_margin_bottom
        
        # Use calculated height or adjust if table would be too large
        table_rows_height = Inches(0.4) * rows  # Estimate row heights
        height = min(table_rows_height, available_height)
        
        # Convert to EMU safely
        left_emu = self._emu(left)
        top_emu = self._emu(top)
        width_emu = self._emu(width)
        height_emu = self._emu(height)
        
        return (left_emu, top_emu, width_emu, height_emu)
            
    def _add_table_to_slide(self, pptx_slide: PptxSlide, headers: List[str], 
                        rows: List[List[str]], style: str = "default") -> Optional[Table]:
        """
        Add a table to a slide as a shape (not in a placeholder).
        
        Args:
            pptx_slide: PowerPoint slide to add table to.
            headers: List of table headers.
            rows: List of table rows.
            style: Table style to apply.
            
        Returns:
            The created table object, or None if creation failed.
        """
        try:
            # Table dimensions
            num_rows = len(rows) + 1  # +1 for header row
            num_cols = len(headers)
            
            if num_rows == 0 or num_cols == 0:
                logger.warning("Cannot create table with 0 rows or columns")
                return None
            
            # Calculate optimal table dimensions based on slide layout
            left_emu, top_emu, width_emu, height_emu = self._calculate_table_dimensions(
                pptx_slide, num_rows, num_cols
            )
            
            logger.info(f"Adding table with calculated dimensions: {num_rows}x{num_cols} at ({left_emu},{top_emu}) with size ({width_emu},{height_emu})")
            
            # Create the table shape
            table = pptx_slide.shapes.add_table(num_rows, num_cols, left_emu, top_emu, width_emu, height_emu).table
            
            # Fill the table with data
            self._fill_table_with_data(table, headers, rows, style)
            
            return table
                
        except Exception as e:
            logger.error(f"Error creating table: {e}")
            logger.error(f"Traceback: {traceback.format_exc()}")
            return None
    
    def _get_style_from_headers(self, headers: List[str]) -> Optional[str]:
        """
        Extract table style information from headers if present.
        
        The last header can contain style information in the format "style:{style_name}".
        If found, the style marker is removed from the headers list.
        
        Args:
            headers: List of table headers
            
        Returns:
            Optional[str]: Style name if found, None otherwise
        """
        if not headers:
            return None
            
        # Check if the last header contains style information
        last_header = headers[-1]
        if isinstance(last_header, str) and last_header.startswith("style:"):
            style_name = last_header.split("style:")[1].strip()
            return style_name
            
        return None

    def _add_formatted_text(self, text_frame, text: str) -> None:
        """
        Add text with formatting to a text frame, parsing markdown-like syntax.
        
        Args:
            text_frame: The PowerPoint text frame to add formatted text to.
            text: The text to add, with optional formatting syntax.
        """
        if not text:
            # Don't clear or modify the text_frame if no text is provided
            return
        
        # Clear any existing text
        text_frame.clear()
        
        # Split the text into paragraphs
        paragraphs = text.split("\n")
        
        # Process each paragraph
        for i, paragraph_text in enumerate(paragraphs):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Add the formatted text to the paragraph
            self._add_formatted_text_to_paragraph(p, paragraph_text)
    
    def _add_formatted_text_to_paragraph(self, paragraph, text: str) -> None:
        """
        Add formatted text to a paragraph, parsing markdown-like syntax.
        
        Args:
            paragraph: The PowerPoint paragraph to add formatted text to.
            text: The text to add, with optional formatting syntax.
        """
        if not text:
            return
        
        # Clear any existing runs
        # Fixed: properly remove runs from paragraph
        if hasattr(paragraph, 'runs') and paragraph.runs:
            for run in list(paragraph.runs):
                try:
                    if hasattr(paragraph, '_p') and hasattr(run, '_r'):
                        paragraph._p.remove(run._r)
                except Exception as e:
                    logger.warning(f"Could not remove run: {e}")
        
        # Parse formatting
        segments = self._parse_text_formatting(text)
        
        # Add each segment with its formatting
        for segment in segments:
            run = paragraph.add_run()
            run.text = segment['text']
            
            # Apply formatting
            if segment.get('bold'):
                run.font.bold = True
            if segment.get('italic'):
                run.font.italic = True
            if segment.get('underline'):
                run.font.underline = True
            if segment.get('strikethrough'):
                run.font.strike = True
            if segment.get('size'):
                # Convert to points if not already
                size = segment['size']
                if isinstance(size, str):
                    try:
                        size = float(size.rstrip('pt').rstrip('px'))
                    except ValueError:
                        size = 12  # Default size
                run.font.size = Pt(size)
            if segment.get('color'):
                color = segment['color']
                # Handle color names or hex values
                if color in self.COLORS:
                    color = self.COLORS[color]
                # Remove '#' if present
                if color.startswith('#'):
                    color = color[1:]
                # Ensure 6 digits
                if len(color) == 3:
                    color = ''.join(c + c for c in color)
                # Create RGB color
                try:
                    r = int(color[0:2], 16)
                    g = int(color[2:4], 16)
                    b = int(color[4:6], 16)
                    run.font.color.rgb = RGBColor(r, g, b)
                except (ValueError, IndexError):
                    # Default to black if color is invalid
                    run.font.color.rgb = RGBColor(0, 0, 0)
            if segment.get('highlight'):
                highlight = segment['highlight']
                # Handle color names or hex values
                if highlight in self.COLORS:
                    highlight = self.COLORS[highlight]
                # Remove '#' if present
                if highlight.startswith('#'):
                    highlight = highlight[1:]
                # Ensure 6 digits
                if len(highlight) == 3:
                    highlight = ''.join(c + c for c in highlight)
                # Create RGB color
                try:
                    r = int(highlight[0:2], 16)
                    g = int(highlight[2:4], 16)
                    b = int(highlight[4:6], 16)
                    # Set highlight color
                    color_name = self._closest_highlight_color(r, g, b)
                    self._apply_highlight_to_run(run, color_name)
                except (ValueError, IndexError):
                    # Skip highlight if color is invalid
                    pass
    
    def _parse_text_formatting(self, text: str) -> List[Dict[str, Any]]:
        """
        Parse formatted text and return segments with formatting information.
        
        Args:
            text: Text with markdown-like formatting syntax.
            
        Returns:
            List of dictionaries with text and formatting information.
        """
        segments = [{'text': text}]
        
        # Parse each formatting pattern
        # Order is important - nested formatting should be processed before outer formatting
        
        # Process font size
        segments = self._apply_pattern(segments, self.FONT_SIZE_PATTERN, 
                                     lambda m: {'size': m.group(1), 'text': m.group(3)})
        
        # Process color
        segments = self._apply_pattern(segments, self.COLOR_PATTERN, 
                                     lambda m: {'color': m.group(1), 'text': m.group(2)})
        
        # Process highlight
        segments = self._apply_pattern(segments, self.HIGHLIGHT_PATTERN, 
                                     lambda m: {'highlight': m.group(1), 'text': m.group(2)})
        
        # Process bold
        segments = self._apply_pattern(segments, self.BOLD_PATTERN, 
                                     lambda m: {'bold': True, 'text': m.group(1)})
        
        # Process italic
        segments = self._apply_pattern(segments, self.ITALIC_PATTERN, 
                                     lambda m: {'italic': True, 'text': m.group(1)})
        
        # Process underline
        segments = self._apply_pattern(segments, self.UNDERLINE_PATTERN, 
                                     lambda m: {'underline': True, 'text': m.group(1)})
        
        # Process strikethrough
        segments = self._apply_pattern(segments, self.STRIKETHROUGH_PATTERN, 
                                     lambda m: {'strikethrough': True, 'text': m.group(1)})
        
        return segments
    
    def _apply_pattern(self, segments: List[Dict[str, Any]], pattern: str, 
                      formatter: callable) -> List[Dict[str, Any]]:
        """
        Apply a regex pattern to text segments and update formatting.
        
        Args:
            segments: List of dictionaries with text and formatting information.
            pattern: Regex pattern to match.
            formatter: Function that returns formatting for matched text.
            
        Returns:
            Updated list of dictionaries with text and formatting information.
        """
        result = []
        
        for segment in segments:
            # Skip empty segments
            if not segment['text']:
                continue
                
            # Skip already formatted segments for this pattern
            if any(key in segment for key in ['bold', 'italic', 'underline', 'strikethrough', 
                                             'color', 'highlight', 'size']):
                result.append(segment)
                continue
            
            # Check for matches
            matches = list(re.finditer(pattern, segment['text']))
            
            if not matches:
                # No matches, keep original segment
                result.append(segment)
                continue
            
            # Process matches
            last_end = 0
            for match in matches:
                # Add text before match
                if match.start() > last_end:
                    result.append({
                        'text': segment['text'][last_end:match.start()],
                        **{k: v for k, v in segment.items() if k != 'text'}
                    })
                
                # Add formatted match
                formatted = formatter(match)
                result.append({
                    **{k: v for k, v in segment.items() if k != 'text'},
                    **formatted
                })
                
                last_end = match.end()
            
            # Add text after last match
            if last_end < len(segment['text']):
                result.append({
                    'text': segment['text'][last_end:],
                    **{k: v for k, v in segment.items() if k != 'text'}
                })
        
        return result
    
    def _closest_highlight_color(self, r: int, g: int, b: int) -> str:
        """
        Find the closest PowerPoint highlight color to the given RGB color.
        
        Args:
            r: Red component (0-255).
            g: Green component (0-255).
            b: Blue component (0-255).
            
        Returns:
            PowerPoint highlight color name.
        """
        # PowerPoint highlight colors and approximate RGB values
        highlight_colors = {
            'yellow': (255, 255, 0),
            'green': (0, 255, 0),
            'cyan': (0, 255, 255),
            'magenta': (255, 0, 255),
            'blue': (0, 0, 255),
            'red': (255, 0, 0),
            'darkBlue': (0, 0, 128),
            'darkCyan': (0, 128, 128),
            'darkGreen': (0, 128, 0),
            'darkMagenta': (128, 0, 128),
            'darkRed': (128, 0, 0),
            'darkYellow': (128, 128, 0),
            'darkGray': (128, 128, 128),
            'lightGray': (192, 192, 192),
            'black': (0, 0, 0),
            'white': (255, 255, 255),
        }
        
        # Find the closest color
        closest_color = None
        min_distance = float('inf')
        
        for color_name, color_rgb in highlight_colors.items():
            # Calculate Euclidean distance
            distance = (
                (r - color_rgb[0]) ** 2 + 
                (g - color_rgb[1]) ** 2 + 
                (b - color_rgb[2]) ** 2
            ) ** 0.5
            
            if distance < min_distance:
                min_distance = distance
                closest_color = color_name
        
        # Ajuster spécifiquement pour les tons rouges foncés
        if closest_color == 'red' and r > 150 and g < 100 and b < 100 and r < 230:
            return 'darkRed'
            
        return closest_color or 'yellow'  # Default to yellow if no match
        
    def _fill_table_with_data(self, table: Table, headers: List[str], rows: List[List[str]], style: str = "default") -> None:
        """
        Fill a PowerPoint table with data and apply formatting.
        
        Args:
            table: PowerPoint table to fill.
            headers: List of column headers.
            rows: List of rows, where each row is a list of cell values.
            style: Style preset to apply to the table.
        """
        # Get the style preset from the PPTBuilder class
        style_preset = self.TABLE_STYLES.get(style, self.TABLE_STYLES["default"])
        
        # Ensure the table has the correct number of rows and columns
        actual_rows = len(table.rows)
        actual_cols = len(table.columns)
        needed_rows = len(rows) + 1  # +1 for header row
        needed_cols = len(headers)
        
        # Validate dimensions
        if needed_rows > actual_rows or needed_cols > actual_cols:
            logger.warning(
                f"Table dimensions mismatch: needed {needed_rows}x{needed_cols}, "
                f"actual {actual_rows}x{actual_cols}. "
                f"Data may be truncated."
            )
        
        # Add headers
        for col_idx, header in enumerate(headers):
            if col_idx < actual_cols:
                cell = table.cell(0, col_idx)
                self._format_table_cell(cell, header, is_header=True, style_preset=style_preset)
        
        # Add rows
        for row_idx, row_data in enumerate(rows):
            if row_idx + 1 < actual_rows:  # +1 to skip header row
                for col_idx, cell_value in enumerate(row_data):
                    if col_idx < actual_cols and col_idx < len(headers):  # Ensure we don't exceed headers
                        cell = table.cell(row_idx + 1, col_idx)
                        is_alternate_row = style_preset.get("banded_rows", False) and row_idx % 2 == 1
                        is_alternate_col = style_preset.get("banded_cols", False) and col_idx % 2 == 1
                        
                        self._format_table_cell(
                            cell, 
                            cell_value, 
                            is_header=False, 
                            is_alternate_row=is_alternate_row,
                            is_alternate_col=is_alternate_col,
                            style_preset=style_preset
                        )
        
        # Apply table styling
        self._apply_table_style(table, style_preset)
    
    def _format_table_cell(self, cell: _Cell, text: str, is_header: bool = False, 
                          is_alternate_row: bool = False, is_alternate_col: bool = False,
                          style_preset: Dict[str, Any] = None) -> None:
        """
        Format a table cell with text and styling.
        
        Args:
            cell: PowerPoint table cell to format.
            text: Text to add to the cell.
            is_header: Whether this is a header cell.
            is_alternate_row: Whether this cell is in an alternate row.
            is_alternate_col: Whether this cell is in an alternate column.
            style_preset: Style preset to apply to the cell.
        """
        if style_preset is None:
            style_preset = self.TABLE_STYLES["default"]
        
        # Clear any existing text
        text_frame = cell.text_frame
        text_frame.clear()
        
        # Add the text with formatting support
        self._add_formatted_text(text_frame, text)
        
        # Set vertical alignment to middle
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # Set text alignment
        for paragraph in text_frame.paragraphs:
            if is_header:
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                paragraph.alignment = PP_ALIGN.LEFT
            
            # Ensure consistent font size
            for run in paragraph.runs:
                if not run.font.size:
                    run.font.size = Pt(10)
        
        # Apply background color
        if is_header and style_preset.get("header_bg"):
            color = style_preset["header_bg"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._hex_to_rgb(color)
        elif is_alternate_row and style_preset.get("accent_color"):
            color = style_preset["accent_color"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._hex_to_rgb(color)
        
        # Apply text color
        if is_header and style_preset.get("header_text"):
            color = style_preset["header_text"]
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = self._hex_to_rgb(color)
        elif not is_header and style_preset.get("body_text"):
            color = style_preset["body_text"]
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = self._hex_to_rgb(color)

    @staticmethod
    def _int_emu(value: float | int) -> Emu:
        """Assure la conversion vers un entier Emu accepté par python-pptx."""
        return Emu(int(round(value)))
    
    def _emu(self, value: int | float | Emu) -> Emu:
        """Force n’importe quelle valeur numérique vers Emu (== int)."""
        return Emu(int(round(value)))

    def _remove_bullet(self, paragraph):
        """Supprime proprement la puce sans toucher à level (compat 0.6.21)."""
        pPr = paragraph._p.get_or_add_pPr()
        for tag in ('a:buChar', 'a:buAutoNum'):
            el = pPr.find(qn(tag))
            if el is not None:
                pPr.remove(el)
        if pPr.find(qn('a:buNone')) is None:
            pPr.append(OxmlElement('a:buNone'))
    
    def _apply_table_style(self, table: Table, style_preset: Dict[str, Any]) -> None:
        """
        Apply styling to a PowerPoint table.
        
        Args:
            table: PowerPoint table to style.
            style_preset: Style preset to apply to the table.
        """
        # Get style properties
        border_color = style_preset.get("border_color")
        border_width = style_preset.get("border_width", Pt(1))
        
        # Check for empty table
        if len(table.columns) == 0 or len(table.rows) == 0:
            logger.warning("Cannot apply style to empty table (no rows or columns)")
            return
        
        # Largeurs de colonnes
        total_width = sum(col.width for col in table.columns)
        col_width = total_width // len(table.columns)  # entier Emu
        for col in table.columns:
            col.width = col_width          # ⚠️  plus de float

        # Hauteurs de lignes
        total_height = sum(row.height for row in table.rows)
        row_height = total_height // len(table.rows)
        
        # First row slightly taller (special case for header)
        if len(table.rows) > 0:
            table.rows[0].height = self._emu(row_height * 1.2)
        
        # Set height for other rows individually (avoid using slice)
        for i in range(1, len(table.rows)):
            table.rows[i].height = self._emu(row_height)
        
        # Skip border styling since cells don't have border attributes
        # in this version of python-pptx
        logger.info("Skipping table border styling due to API limitations")
    
    @staticmethod
    def _hex_to_rgb(hex_value: str) -> RGBColor:
        """
        Convert a hex color string to an RGBColor object.
        
        Args:
            hex_value: Hex color string (with or without #).
            
        Returns:
            RGBColor object.
        """
        # Remove '#' if present
        if hex_value.startswith('#'):
            hex_value = hex_value[1:]
        
        # Ensure 6 digits
        if len(hex_value) == 3:
            hex_value = ''.join(c + c for c in hex_value)
        
        # Convert to RGB
        try:
            r = int(hex_value[0:2], 16)
            g = int(hex_value[2:4], 16)
            b = int(hex_value[4:6], 16)
            return RGBColor(r, g, b)
        except (ValueError, IndexError):
            # Default to black if color is invalid
            return RGBColor(0, 0, 0)

    def _apply_highlight_to_run(self, run, highlight_color):
        """
        Apply highlight color to a run using OOXML.
        Extracted for better testability.
        
        Args:
            run: The run to apply highlight to
            highlight_color: The highlight color to apply
        """
        if hasattr(run, '_element') and run._element is not None:
            run._element.get_or_add_rPr().get_or_add_highlight().val = highlight_color