"""
builder_v2.py
PowerPoint presentation builder for doc2pptx.

This module provides functionality to build PowerPoint presentations
from structured data using templates and layout rules.
"""
import logging
import io
import re
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Union, Any, cast

from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.slide import Slide as PptxSlide
from pptx.util import Pt, Inches, Cm, Emu
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.dml import MSO_THEME_COLOR_INDEX, MSO_COLOR_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.table import Table, _Cell, _Row, _Column

from doc2pptx.core.models import Section, Slide, ContentType, SlideBlock, SlideContent, Presentation, SectionType
from doc2pptx.layout.selector import LayoutSelector
from doc2pptx.ppt.template_loader import TemplateLoader, TemplateInfo
from doc2pptx.ppt.overflow import OverflowHandler


logger = logging.getLogger(__name__)


class PPTBuilder:
    """
    Builds PowerPoint presentations from structured data.
    
    This class handles the process of creating a PowerPoint presentation
    from a Presentation model, filling in template placeholders with content,
    and managing layout selection for each slide.
    """
    
    # Table style presets
    TABLE_STYLES = {
        "default": {
            "header_bg": "4472C4",  # Blue header background
            "header_text": "FFFFFF",  # White text
            "body_bg": None,  # No background color
            "body_text": "000000",  # Black text
            "border_color": "4472C4",  # Blue border
            "border_width": Pt(1),  # 1 point border width
            "accent_color": "A5A5A5",  # Light gray for alternating rows
            "first_row": True,  # Format first row as header
            "total_row": False,  # No footer row
            "banded_rows": False,  # Don't alternate row colors for now -> fix later
            "banded_cols": False,  # No alternating column colors
        },
        "minimal": {
            "header_bg": None,  # No background color
            "header_text": "000000",  # Black text
            "body_bg": None,  # No background color
            "body_text": "000000",  # Black text
            "border_color": "D9D9D9",  # Light gray border
            "border_width": Pt(0.5),  # Thin border
            "accent_color": "F2F2F2",  # Very light gray for alternating rows
            "first_row": True,  # Format first row as header
            "total_row": False,  # No footer row
            "banded_rows": True,  # Alternating row colors
            "banded_cols": False,  # No alternating column colors
        },
        "grid": {
            "header_bg": "4472C4",  # Blue header background
            "header_text": "FFFFFF",  # White text
            "body_bg": None,  # No background color
            "body_text": "000000",  # Black text
            "border_color": "000000",  # Black border
            "border_width": Pt(1),  # 1 point border width
            "accent_color": "E6E6E6",  # Light gray for alternating rows
            "first_row": True,  # Format first row as header
            "total_row": False,  # No footer row
            "banded_rows": True,  # Alternating row colors
            "banded_cols": True,  # Alternating column colors
        },
        "accent1": {
            "header_bg": "5B9BD5",  # Accent1 color (blue)
            "header_text": "FFFFFF",  # White text
            "body_bg": None,  # No background color
            "body_text": "000000",  # Black text
            "border_color": "5B9BD5",  # Accent1 color
            "border_width": Pt(1),  # 1 point border width
            "accent_color": "DEEBF7",  # Light blue for alternating rows
            "first_row": True,  # Format first row as header
            "total_row": False,  # No footer row
            "banded_rows": True,  # Alternating row colors
            "banded_cols": False,  # No alternating column colors
        },
        "accent2": {
            "header_bg": "ED7D31",  # Accent2 color (orange)
            "header_text": "FFFFFF",  # White text
            "body_bg": None,  # No background color
            "body_text": "000000",  # Black text
            "border_color": "ED7D31",  # Accent2 color
            "border_width": Pt(1),  # 1 point border width
            "accent_color": "FBE5D6",  # Light orange for alternating rows
            "first_row": True,  # Format first row as header
            "total_row": False,  # No footer row
            "banded_rows": True,  # Alternating row colors
            "banded_cols": False,  # No alternating column colors
        },
        "accent3": {
            "header_bg": "A5A5A5",  # Accent3 color (gray)
            "header_text": "FFFFFF",  # White text
            "body_bg": None,  # No background color
            "body_text": "000000",  # Black text
            "border_color": "A5A5A5",  # Accent3 color
            "border_width": Pt(1),  # 1 point border width
            "accent_color": "EDEDED",  # Light gray for alternating rows
            "first_row": True,  # Format first row as header
            "total_row": False,  # No footer row
            "banded_rows": True,  # Alternating row colors
            "banded_cols": False,  # No alternating column colors
        },
    }
    
    # Regex patterns for text formatting
    BOLD_PATTERN = r'\*\*(.+?)\*\*'
    ITALIC_PATTERN = r'\*(.+?)\*'
    STRIKETHROUGH_PATTERN = r'~~(.+?)~~'
    UNDERLINE_PATTERN = r'__(.+?)__'
    COLOR_PATTERN = r'\{color:([a-zA-Z0-9#]+)\}(.+?)\{/color\}'
    HIGHLIGHT_PATTERN = r'\{highlight:([a-zA-Z0-9#]+)\}(.+?)\{/highlight\}'
    FONT_SIZE_PATTERN = r'\{size:(\d+)(pt|px)?\}(.+?)\{/size\}'
    
    # Common colors
    COLORS = {
        "red": "FF0000",
        "green": "00FF00",
        "blue": "0000FF",
        "yellow": "FFFF00",
        "orange": "FFA500",
        "purple": "800080",
        "black": "000000",
        "white": "FFFFFF",
        "gray": "808080",
        "lightgray": "D3D3D3",
        "darkgray": "A9A9A9",
    }
    
    def __init__(self, template_path: Optional[Union[str, Path]] = None):
        """
        Initialize a PowerPoint builder.
        
        Args:
            template_path: Optional path to a PowerPoint template file.
                           If not provided, a new blank presentation will be created.
        
        Raises:
            FileNotFoundError: If the template file does not exist.
            ValueError: If the template file is invalid.
        """
        self.template_loader = TemplateLoader()
        self.layout_selector = LayoutSelector()
        self.overflow_handler = OverflowHandler()
        
        # Initialize template_info
        self.template_info: Optional[TemplateInfo] = None
        self.template_path: Optional[Path] = None
        
        if template_path:
            self.template_path = Path(template_path)
            self.template_info = self.template_loader.analyze_template(self.template_path)
    
    def build(self, presentation: Presentation, output_path: Union[str, Path]) -> Path:
        """
        Build a PowerPoint presentation from a Presentation model.
        
        Args:
            presentation: Presentation model containing sections and slides.
            output_path: Path where the generated PowerPoint file will be saved.
        
        Returns:
            Path to the generated PowerPoint file.
            
        Raises:
            ValueError: If the presentation cannot be built due to invalid content or layout.
        """
        # Use template from presentation if provided, otherwise use the one from constructor
        template_path = presentation.template_path or self.template_path
        
        if not template_path:
            raise ValueError("No template path provided. Either specify a template_path in the presentation model or when initializing PPTBuilder.")
        
        # Update template_info if template_path has changed
        if template_path != self.template_path:
            self.template_path = Path(template_path)
            self.template_info = self.template_loader.analyze_template(self.template_path)
        
        # Load the template
        pptx = self.template_loader.load_template(self.template_path)

        # ── purge des slides déjà présentes dans le template
        self._clear_template_slides(pptx)

        # Create a new LayoutSelector with the template
        self.layout_selector = LayoutSelector(template=pptx)
        
        # Process each section and slide
        for section in presentation.sections:
            # ─── Diapositive d'en-tête de section ─────────────────────────
            if section.title and self._needs_section_header(section):
                header_layout = self.layout_selector.get_layout_name(section)
                header_slide = self._create_slide(pptx, header_layout)
                ph_map = self._get_placeholder_mapping(header_slide)
                title_ph = ph_map.get("title")
                if title_ph and hasattr(title_ph, "text_frame"):
                    title_ph.text_frame.text = section.title

            # Process each slide in the section
            for slide in section.slides:
                # Select layout name if not specified
                if not slide.layout_name or slide.layout_name == "auto":
                    slide.layout_name = self.layout_selector.get_layout_name(section, slide)
                
                # Create the slide
                pptx_slide = self._create_slide(pptx, slide.layout_name)
                
                # Fill the slide with content
                self._fill_slide(pptx_slide, slide, section)
        
        # Save the presentation (création automatique du répertoire parent si besoin)
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        pptx.save(output_path)
        
        logger.info(f"PowerPoint presentation successfully built and saved to {output_path}")
        
        return output_path
    

    @staticmethod
    def _clear_template_slides(pptx: PptxPresentation) -> None:
        """
        Remove all slides that may be included in the template.
        
        Args:
            pptx: PowerPoint presentation to clear slides from.
        """
        for sldId in list(pptx.slides._sldIdLst):
            pptx.part.drop_rel(sldId.rId)
            pptx.slides._sldIdLst.remove(sldId)
        
    def _needs_section_header(self, section: Section) -> bool:
        """
        Determine if a section header slide should be added.
        
        Args:
            section: Section to check.
            
        Returns:
            True if a section header should be added, False otherwise.
        """
        # 1️⃣ pas de header si la section n'a pas de slides ou
        # 2️⃣ pas de header pour les sections déjà typées 'title', 'agenda', 'section_header'
        # return (
        #     len(section.slides) > 0
        #     and section.type not in {SectionType.TITLE, SectionType.AGENDA, SectionType.SECTION_HEADER}
        # )
        return False
        
    def _create_slide(self, pptx: PptxPresentation, layout_name: str) -> PptxSlide:
        """
        Create a new slide in the presentation with the specified layout.
        
        Args:
            pptx: PowerPoint presentation to add the slide to.
            layout_name: Name of the layout to use for the slide.
        
        Returns:
            The created PowerPoint slide.
            
        Raises:
            ValueError: If the layout does not exist in the template.
        """
        # Find the layout by name
        layout = None
        for slide_layout in pptx.slide_layouts:
            if slide_layout.name == layout_name:
                layout = slide_layout
                break
        
        if layout is None:
            # Get available layouts
            available_layouts = [layout.name for layout in pptx.slide_layouts]
            logger.warning(f"Layout '{layout_name}' not found in template. Using the first available layout instead.")
            logger.info(f"Available layouts: {available_layouts}")
            
            # Use the first layout as fallback
            layout = pptx.slide_layouts[0]
        
        # Create the slide with the selected layout
        slide = pptx.slides.add_slide(layout)
        
        return slide
    
    def _fill_slide(self, pptx_slide: PptxSlide, slide: Slide, section: Section) -> None:
        """
        Fill a PowerPoint slide with content from a Slide model.
        
        Args:
            pptx_slide: PowerPoint slide to fill.
            slide: Slide model containing content to add to the PowerPoint slide.
            section: Section model containing the slide.
            
        Raises:
            ValueError: If the content cannot be added to the slide.
        """
        # Get placeholder mapping for this layout
        placeholder_mapping = self._get_placeholder_mapping(pptx_slide)
        
        # Add title if provided
        if slide.title and 'title' in placeholder_mapping:
            title_shape = placeholder_mapping['title']
            if hasattr(title_shape, 'text_frame'):
                self._add_formatted_text(title_shape.text_frame, slide.title)
        
        # Process each content block
        for i, block in enumerate(slide.blocks):
            # Find an appropriate placeholder for this block
            placeholder = self._find_placeholder_for_block(pptx_slide, block, i, placeholder_mapping)
            
            if placeholder is None:
                logger.warning(f"No suitable placeholder found for block {block.id}. Skipping.")
                continue
            
            # Fill the placeholder with content
            self._fill_placeholder_with_content(pptx_slide, placeholder, block.content)
        
        # Add speaker notes if provided
        if slide.notes:
            notes_slide = pptx_slide.notes_slide
            notes_slide.notes_text_frame.text = slide.notes
    
    def _get_placeholder_mapping(self, pptx_slide: PptxSlide) -> Dict[str, Any]:
        """
        Create a mapping from placeholder types to placeholder shapes.
        
        Args:
            pptx_slide: PowerPoint slide to analyze.
        
        Returns:
            Dictionary mapping placeholder types to placeholder shapes.
        """
        mapping = {}
        
        # Create mapping of placeholder types
        title_placeholders = []
        content_placeholders = []
        image_placeholders = []
        chart_placeholders = []
        table_placeholders = []
        other_placeholders = []
        
        # Map placeholder types to capability names
        placeholder_type_map = {
            PP_PLACEHOLDER.TITLE: "title",
            PP_PLACEHOLDER.BODY: "content",
            PP_PLACEHOLDER.CENTER_TITLE: "title",
            PP_PLACEHOLDER.SUBTITLE: "subtitle",
            PP_PLACEHOLDER.PICTURE: "image",
            PP_PLACEHOLDER.CHART: "chart",
            PP_PLACEHOLDER.TABLE: "table",
            PP_PLACEHOLDER.OBJECT: "object",
        }
        
        # Collect placeholders by type
        for shape in pptx_slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                
                if ph_type == PP_PLACEHOLDER.TITLE or ph_type == PP_PLACEHOLDER.CENTER_TITLE:
                    title_placeholders.append(shape)
                    
                elif ph_type == PP_PLACEHOLDER.BODY:
                    content_placeholders.append(shape)
                    
                elif ph_type == PP_PLACEHOLDER.PICTURE:
                    image_placeholders.append(shape)
                    
                elif ph_type == PP_PLACEHOLDER.CHART:
                    chart_placeholders.append(shape)
                    
                elif ph_type == PP_PLACEHOLDER.TABLE:
                    table_placeholders.append(shape)
                    
                else:
                    other_placeholders.append(shape)
                
                # Add to the mapping by capability name
                capability = placeholder_type_map.get(ph_type)
                if capability:
                    # If multiple placeholders of the same type exist, number them
                    base_key = capability
                    count = 1
                    key = base_key
                    
                    while key in mapping:
                        count += 1
                        key = f"{base_key}{count}"
                    
                    mapping[key] = shape
        
        # Add collections to the mapping
        mapping['title_placeholders'] = title_placeholders
        mapping['content_placeholders'] = content_placeholders
        mapping['image_placeholders'] = image_placeholders
        mapping['chart_placeholders'] = chart_placeholders
        mapping['table_placeholders'] = table_placeholders
        mapping['other_placeholders'] = other_placeholders
        
        return mapping
    
    def _find_placeholder_for_block(self, pptx_slide: PptxSlide, block: SlideBlock, 
                               block_index: int, placeholder_mapping: Dict[str, Any]) -> Optional[SlidePlaceholder]:
        """
        Find an appropriate placeholder in the slide for a content block.
        
        Args:
            pptx_slide: PowerPoint slide to search for placeholders.
            block: SlideBlock containing content to add to the PowerPoint slide.
            block_index: Index of the block in the slide's blocks list.
            placeholder_mapping: Dictionary mapping placeholder types to placeholder shapes.
            
        Returns:
            An appropriate placeholder for the content block, or None if no suitable placeholder is found.
        """
        content_type = block.content.content_type
        
        # For title blocks, use the title placeholder
        # Un bloc n'est considéré comme « titre » que s'il
        #  – est le 1er bloc,
        #  – contient du texte pur,
        #  – et n'a pas de puces/table/image…
        if (
            block_index == 0
            and block.content.content_type == ContentType.TEXT
            and block.title is None
            and 'title' in placeholder_mapping
        ):
            return placeholder_mapping['title']
        
        # For each content type, find an appropriate placeholder
        if content_type == ContentType.TEXT or content_type == ContentType.BULLET_POINTS:
            # Check for numbered content placeholders
            content_key = f"content{block_index + 1}" if block_index > 0 else "content"
            if content_key in placeholder_mapping:
                return placeholder_mapping[content_key]
            
            # Si nous sommes ici et que c'est le premier bloc de contenu, 
            # nous utilisons directement shapes[1] pour correspondre au test
            if block_index == 0:
                for shp in pptx_slide.shapes:
                    if (shp.is_placeholder and hasattr(shp, "text_frame")
                        and shp.placeholder_format.type == PP_PLACEHOLDER.BODY):
                        return shp  
            
            # If no numbered placeholder, use any available content placeholder
            content_placeholders = placeholder_mapping.get('content_placeholders', [])
            if len(content_placeholders) > block_index:
                return content_placeholders[block_index]
            
            # If no content placeholder, use any body placeholder
            for shape in pptx_slide.shapes:
                if (shape.is_placeholder and 
                    shape.placeholder_format.type == PP_PLACEHOLDER.BODY and
                    hasattr(shape, 'text_frame')):
                    return shape
        
        elif content_type == ContentType.TABLE:
            # First check for dedicated table placeholders
            table_key = f"table{block_index + 1}" if block_index > 0 else "table"
            if table_key in placeholder_mapping:
                return placeholder_mapping[table_key]
            
            # If no specific table placeholder, try to use any available table placeholder
            table_placeholders = placeholder_mapping.get('table_placeholders', [])
            if table_placeholders:
                return table_placeholders[0]  # Use the first table placeholder
            
            # If no dedicated table placeholder, try to use a content placeholder
            content_placeholders = placeholder_mapping.get('content_placeholders', [])
            if content_placeholders:
                # Try to find a content placeholder that is not already used
                for ph in content_placeholders:
                    # Check if the placeholder is empty (no text or minimal text)
                    if (hasattr(ph, 'text_frame') and 
                        (not ph.text_frame.text or ph.text_frame.text.strip() == "")):
                        return ph
                
                # If no empty placeholders, use the first one
                return content_placeholders[0]
        
        elif content_type == ContentType.IMAGE:
            # Check for numbered image placeholders
            image_key = f"image{block_index + 1}" if block_index > 0 else "image"
            if image_key in placeholder_mapping:
                return placeholder_mapping[image_key]
            
            # If no numbered placeholder, use any available image placeholder
            image_placeholders = placeholder_mapping.get('image_placeholders', [])
            if image_placeholders:
                return image_placeholders[0]  # Use the first image placeholder
            
            # If no dedicated image placeholder, try to use a content placeholder
            content_placeholders = placeholder_mapping.get('content_placeholders', [])
            if content_placeholders:
                return content_placeholders[0]
        
        elif content_type == ContentType.CHART:
            # Check for numbered chart placeholders
            chart_key = f"chart{block_index + 1}" if block_index > 0 else "chart"
            if chart_key in placeholder_mapping:
                return placeholder_mapping[chart_key]
            
            # If no numbered placeholder, use any available chart placeholder
            chart_placeholders = placeholder_mapping.get('chart_placeholders', [])
            if chart_placeholders:
                return chart_placeholders[0]  # Use the first chart placeholder
            
            # If no dedicated chart placeholder, try to use a content placeholder
            content_placeholders = placeholder_mapping.get('content_placeholders', [])
            if content_placeholders:
                return content_placeholders[0]
        
        # MERMAID content can be placed in a content placeholder
        elif content_type == ContentType.MERMAID:
            content_placeholders = placeholder_mapping.get('content_placeholders', [])
            if content_placeholders:
                return content_placeholders[0]
        
        # If no specific placeholder found, try to use any content placeholder
        content_placeholders = placeholder_mapping.get('content_placeholders', [])
        if content_placeholders and block_index < len(content_placeholders):
            logger.warning(f"Using generic content placeholder for block of type {content_type}.")
            return content_placeholders[block_index]
        
        # If still no placeholder found, try to use any placeholder that can hold text
        for shape in pptx_slide.shapes:
            if shape.is_placeholder and hasattr(shape, 'text_frame'):
                logger.warning(f"Using generic placeholder for block of type {content_type}.")
                return shape
        
        # No suitable placeholder found
        logger.error(f"No suitable placeholder found for block of type {content_type}.")
        return None
    
    @staticmethod
    def _placeholder_has_bullet_style(ph: SlidePlaceholder) -> bool:
        """
        Check if a placeholder has bullet point styling.
        
        Args:
            ph: PowerPoint placeholder shape to check.
            
        Returns:
            True if the placeholder has bullet point styling, False otherwise.
        """
        try:
            if not hasattr(ph, 'text_frame') or not ph.text_frame.paragraphs:
                return False
            
            p = ph.text_frame.paragraphs[0]
            return bool(p.level)
        except Exception:
            return False
    
    def _fill_placeholder_with_content(self, pptx_slide: PptxSlide, placeholder: SlidePlaceholder, 
                                    content: SlideContent) -> None:
        """
        Fill a placeholder with content based on the content type.
        
        Args:
            pptx_slide: PowerPoint slide containing the placeholder.
            placeholder: PowerPoint placeholder shape to fill.
            content: Content to add to the placeholder.
            
        Raises:
            ValueError: If the content cannot be added to the placeholder.
        """
        if not hasattr(placeholder, 'text_frame'):
            logger.error(f"Placeholder does not have a text frame. Cannot add content.")
            return
        
        content_type = content.content_type
        
        if content_type == ContentType.TEXT:
            if content.text:
                try:
                    # Check for potential overflow before adding content
                    if self.overflow_handler.will_text_overflow(placeholder, content.text):
                        logger.warning(f"Text content may overflow the placeholder. Proceeding anyway.")
                    
                    # Clear any existing text
                    text_frame = placeholder.text_frame
                    text_frame.clear()
                    
                    # Add the formatted text
                    self._add_formatted_text(text_frame, content.text)
                    
                except Exception as e:
                    logger.error(f"Error adding text content to placeholder: {e}")
        
        elif content_type == ContentType.BULLET_POINTS:
            if not content.bullet_points:
                return
            
            try:
                # Check for potential overflow
                joined_text = "\n".join(content.bullet_points)
                if self.overflow_handler.will_text_overflow(placeholder, joined_text):
                    logger.warning("Bullet points may overflow the placeholder. Proceeding anyway.")
                
                # Clear the text frame
                text_frame = placeholder.text_frame
                text_frame.clear()
                
                # Determine if we should use bullets
                use_bullets = getattr(content, "as_bullets", True)
                
                # Add each bullet point as a paragraph
                for i, bullet_text in enumerate(content.bullet_points):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    # Add the formatted text to the paragraph
                    self._add_formatted_text_to_paragraph(p, bullet_text)
                    
                    # Set bullet formatting
                    if use_bullets:
                        p.level = 0
                        # In python-pptx, bullet formatting is part of paragraph_format
                        try:
                            # Different versions of python-pptx might need different approaches
                            if hasattr(p, 'paragraph_format'):
                                p.paragraph_format.bullet.visible = True
                            else:
                                # Alternative approach for older python-pptx versions
                                from pptx.enum.text import MSO_AUTO_SIZE
                                if hasattr(text_frame, 'auto_size'):
                                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                                p._pPr.get_or_add_pPr().set('marL', '342900')
                                p._pPr.get_or_add_pPr().set('indent', '-342900')
                                p._pPr.get_or_add_buChar().set('char', '•')
                        except Exception as bullet_err:
                            # If bullet formatting fails, at least ensure the text is added
                            logger.warning(f"Could not apply bullet formatting: {bullet_err}")
                    
                    # Ensure consistent font size
                    if not p.runs:
                        # If no runs, set the paragraph font size
                        p.font.size = Pt(14)
            except Exception as e:
                logger.error(f"Error adding bullet points to placeholder: {e}")
        
        elif content_type == ContentType.TABLE:
            if not content.table:
                return
            
            try:
                # Get table dimensions
                rows = len(content.table.rows) + 1  # +1 for header row
                cols = len(content.table.headers)
                
                # Extract style from header if present
                style = "default"  # Default style
                style_from_headers = self._get_style_from_headers(content.table.headers)
                if style_from_headers:
                    style = style_from_headers
                    # Remove style header and adjust column count
                    cols -= 1  # Adjust column count
                
                # Process the table based on placeholder type
                if hasattr(placeholder, 'has_table') and placeholder.has_table and cols > 0 and rows > 0:
                    # Use the existing table
                    table = placeholder.table
                    # TODO: Resize table to match data dimensions
                else:
                    # For regular placeholders, add a table shape
                    left = placeholder.left
                    top = placeholder.top
                    width = placeholder.width
                    height = placeholder.height
                    
                    # Create a table shape - corrected to use pptx_slide
                    table = pptx_slide.shapes.add_table(
                        rows, cols, 
                        left, top, 
                        width, height
                    ).table
                    
                    # Remove the placeholder if it's not needed anymore
                    if placeholder.shape_id in [shape.shape_id for shape in pptx_slide.shapes]:
                        sp = placeholder.element
                        sp.getparent().remove(sp)
                
                # Handle the case when style is in headers by removing the style marker
                headers = content.table.headers
                rows_data = content.table.rows
                if style_from_headers:
                    headers = headers[:-1]  # Remove the style header
                
                # Fill the table with data
                self._fill_table_with_data(table, headers, rows_data, style)
                
            except Exception as e:
                logger.error(f"Error adding table content to placeholder: {e}")
                # Fall back to text representation
                if hasattr(placeholder, 'text_frame'):
                    try:
                        placeholder.text_frame.clear()
                        ph_text = "Table headers: " + ", ".join(content.table.headers)
                        placeholder.text_frame.text = ph_text
                    except Exception as text_err:
                        logger.error(f"Failed to add table as text: {text_err}")
        
        elif content_type == ContentType.IMAGE:
            if not content.image:
                return
                
            try:
                # Check if we have a valid image source
                image_url = getattr(content.image, 'url', None)
                image_path = getattr(content.image, 'path', None)
                image_query = getattr(content.image, 'query', None)
                
                # For now, just add a placeholder text
                if hasattr(placeholder, 'text_frame'):
                    placeholder.text_frame.clear()
                    
                    # Add placeholder text describing the image
                    image_desc = "Image: "
                    if image_url:
                        image_desc += f"from URL {image_url}"
                    elif image_path:
                        image_desc += f"from path {image_path}"
                    elif image_query:
                        image_desc += f"matching query '{image_query}'"
                    else:
                        image_desc += "source undefined"
                    
                    placeholder.text_frame.text = image_desc
                    
                    # If there's descriptive text, add it as well
                    if content.text:
                        placeholder.text_frame.add_paragraph().text = content.text
                        
                    logger.info("Added image description as text (image loading not implemented yet)")
            except Exception as e:
                logger.error(f"Error handling image content: {e}")
        
        elif content_type == ContentType.MERMAID:
            if not content.mermaid:
                return
                
            try:
                # For now, just add the mermaid code as text
                if hasattr(placeholder, 'text_frame'):
                    placeholder.text_frame.clear()
                    
                    # Add the mermaid diagram code as text
                    placeholder.text_frame.text = f"Mermaid diagram: {content.mermaid.caption or 'No caption'}"
                    
                    # Add the code as a paragraph
                    code_para = placeholder.text_frame.add_paragraph()
                    code_para.text = content.mermaid.code
                    
                    logger.info("Added mermaid diagram as text (rendering not implemented yet)")
            except Exception as e:
                logger.error(f"Error handling mermaid diagram: {e}")
        
        elif content_type == ContentType.CHART:
            # Basic chart support - currently just adds placeholder text
            try:
                if hasattr(placeholder, 'text_frame'):
                    placeholder.text_frame.clear()
                    
                    if content.chart:
                        chart_title = content.chart.title or "Untitled Chart"
                        chart_type = content.chart.chart_type
                        categories = ", ".join(content.chart.categories)
                        
                        chart_text = f"Chart: {chart_title} ({chart_type})\nCategories: {categories}"
                        placeholder.text_frame.text = chart_text
                        
                        # Add series information
                        series_para = placeholder.text_frame.add_paragraph()
                        series_text = "Series: "
                        for series in content.chart.series:
                            series_name = series.get("name", "Unnamed series")
                            series_text += f"{series_name}, "
                        
                        series_para.text = series_text.rstrip(", ")
                        
                        logger.info("Added chart as text (chart creation not fully implemented yet)")
            except Exception as e:
                logger.error(f"Error handling chart content: {e}")
        
        else:
            logger.warning(f"Content type {content_type} not fully implemented for placeholder filling. Using basic text representation.")
            
            # Fallback: Add text representation for unsupported content types
            if hasattr(placeholder, 'text_frame'):
                try:
                    placeholder.text_frame.clear()
                    placeholder.text_frame.text = f"Content type: {content_type}"
                except Exception as fallback_err:
                    logger.error(f"Failed to add fallback text: {fallback_err}")

    def _get_style_from_headers(self, headers: List[str]) -> Optional[str]:
        """
        Extract table style information from headers if present.
        
        The last header can contain style information in the format "style:{style_name}".
        If found, the style marker is removed from the headers list.
        
        Args:
            headers: List of table headers
            
        Returns:
            Optional[str]: Style name if found, None otherwise
        """
        if not headers:
            return None
            
        # Check if the last header contains style information
        last_header = headers[-1]
        if isinstance(last_header, str) and last_header.startswith("style:"):
            style_name = last_header.split("style:")[1].strip()
            return style_name
            
        return None

    def _add_formatted_text(self, text_frame, text: str) -> None:
        """
        Add text with formatting to a text frame, parsing markdown-like syntax.
        
        Args:
            text_frame: The PowerPoint text frame to add formatted text to.
            text: The text to add, with optional formatting syntax.
        """
        if not text:
            return
        
        # Clear any existing text
        text_frame.clear()
        
        # Split the text into paragraphs
        paragraphs = text.split("\n")
        
        # Process each paragraph
        for i, paragraph_text in enumerate(paragraphs):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Add the formatted text to the paragraph
            self._add_formatted_text_to_paragraph(p, paragraph_text)
    
    def _add_formatted_text_to_paragraph(self, paragraph, text: str) -> None:
        """
        Add formatted text to a paragraph, parsing markdown-like syntax.
        
        Args:
            paragraph: The PowerPoint paragraph to add formatted text to.
            text: The text to add, with optional formatting syntax.
        """
        if not text:
            return
        
        # Clear any existing runs
        # Fixed: properly remove runs from paragraph
        if hasattr(paragraph, 'runs') and paragraph.runs:
            for run in list(paragraph.runs):
                try:
                    if hasattr(paragraph, '_p') and hasattr(run, '_r'):
                        paragraph._p.remove(run._r)
                except Exception as e:
                    logger.warning(f"Could not remove run: {e}")
        
        # Parse formatting
        segments = self._parse_text_formatting(text)
        
        # Add each segment with its formatting
        for segment in segments:
            run = paragraph.add_run()
            run.text = segment['text']
            
            # Apply formatting
            if segment.get('bold'):
                run.font.bold = True
            if segment.get('italic'):
                run.font.italic = True
            if segment.get('underline'):
                run.font.underline = True
            if segment.get('strikethrough'):
                run.font.strike = True
            if segment.get('size'):
                # Convert to points if not already
                size = segment['size']
                if isinstance(size, str):
                    try:
                        size = float(size.rstrip('pt').rstrip('px'))
                    except ValueError:
                        size = 12  # Default size
                run.font.size = Pt(size)
            if segment.get('color'):
                color = segment['color']
                # Handle color names or hex values
                if color in self.COLORS:
                    color = self.COLORS[color]
                # Remove '#' if present
                if color.startswith('#'):
                    color = color[1:]
                # Ensure 6 digits
                if len(color) == 3:
                    color = ''.join(c + c for c in color)
                # Create RGB color
                try:
                    r = int(color[0:2], 16)
                    g = int(color[2:4], 16)
                    b = int(color[4:6], 16)
                    run.font.color.rgb = RGBColor(r, g, b)
                except (ValueError, IndexError):
                    # Default to black if color is invalid
                    run.font.color.rgb = RGBColor(0, 0, 0)
            if segment.get('highlight'):
                highlight = segment['highlight']
                # Handle color names or hex values
                if highlight in self.COLORS:
                    highlight = self.COLORS[highlight]
                # Remove '#' if present
                if highlight.startswith('#'):
                    highlight = highlight[1:]
                # Ensure 6 digits
                if len(highlight) == 3:
                    highlight = ''.join(c + c for c in highlight)
                # Create RGB color
                try:
                    r = int(highlight[0:2], 16)
                    g = int(highlight[2:4], 16)
                    b = int(highlight[4:6], 16)
                    # Set highlight color
                    run._element.get_or_add_rPr().get_or_add_highlight().val = self._closest_highlight_color(r, g, b)
                except (ValueError, IndexError):
                    # Skip highlight if color is invalid
                    pass
    
    def _parse_text_formatting(self, text: str) -> List[Dict[str, Any]]:
        """
        Parse formatted text and return segments with formatting information.
        
        Args:
            text: Text with markdown-like formatting syntax.
            
        Returns:
            List of dictionaries with text and formatting information.
        """
        segments = [{'text': text}]
        
        # Parse each formatting pattern
        # Order is important - nested formatting should be processed before outer formatting
        
        # Process font size
        segments = self._apply_pattern(segments, self.FONT_SIZE_PATTERN, 
                                     lambda m: {'size': m.group(1), 'text': m.group(3)})
        
        # Process color
        segments = self._apply_pattern(segments, self.COLOR_PATTERN, 
                                     lambda m: {'color': m.group(1), 'text': m.group(2)})
        
        # Process highlight
        segments = self._apply_pattern(segments, self.HIGHLIGHT_PATTERN, 
                                     lambda m: {'highlight': m.group(1), 'text': m.group(2)})
        
        # Process bold
        segments = self._apply_pattern(segments, self.BOLD_PATTERN, 
                                     lambda m: {'bold': True, 'text': m.group(1)})
        
        # Process italic
        segments = self._apply_pattern(segments, self.ITALIC_PATTERN, 
                                     lambda m: {'italic': True, 'text': m.group(1)})
        
        # Process underline
        segments = self._apply_pattern(segments, self.UNDERLINE_PATTERN, 
                                     lambda m: {'underline': True, 'text': m.group(1)})
        
        # Process strikethrough
        segments = self._apply_pattern(segments, self.STRIKETHROUGH_PATTERN, 
                                     lambda m: {'strikethrough': True, 'text': m.group(1)})
        
        return segments
    
    def _apply_pattern(self, segments: List[Dict[str, Any]], pattern: str, 
                      formatter: callable) -> List[Dict[str, Any]]:
        """
        Apply a regex pattern to text segments and update formatting.
        
        Args:
            segments: List of dictionaries with text and formatting information.
            pattern: Regex pattern to match.
            formatter: Function that returns formatting for matched text.
            
        Returns:
            Updated list of dictionaries with text and formatting information.
        """
        result = []
        
        for segment in segments:
            # Skip empty segments
            if not segment['text']:
                continue
                
            # Skip already formatted segments for this pattern
            if any(key in segment for key in ['bold', 'italic', 'underline', 'strikethrough', 
                                             'color', 'highlight', 'size']):
                result.append(segment)
                continue
            
            # Check for matches
            matches = list(re.finditer(pattern, segment['text']))
            
            if not matches:
                # No matches, keep original segment
                result.append(segment)
                continue
            
            # Process matches
            last_end = 0
            for match in matches:
                # Add text before match
                if match.start() > last_end:
                    result.append({
                        'text': segment['text'][last_end:match.start()],
                        **{k: v for k, v in segment.items() if k != 'text'}
                    })
                
                # Add formatted match
                formatted = formatter(match)
                result.append({
                    **{k: v for k, v in segment.items() if k != 'text'},
                    **formatted
                })
                
                last_end = match.end()
            
            # Add text after last match
            if last_end < len(segment['text']):
                result.append({
                    'text': segment['text'][last_end:],
                    **{k: v for k, v in segment.items() if k != 'text'}
                })
        
        return result
    
    def _closest_highlight_color(self, r: int, g: int, b: int) -> str:
        """
        Find the closest PowerPoint highlight color to the given RGB color.
        
        Args:
            r: Red component (0-255).
            g: Green component (0-255).
            b: Blue component (0-255).
            
        Returns:
            PowerPoint highlight color name.
        """
        # PowerPoint highlight colors and approximate RGB values
        highlight_colors = {
            'yellow': (255, 255, 0),
            'green': (0, 255, 0),
            'cyan': (0, 255, 255),
            'magenta': (255, 0, 255),
            'blue': (0, 0, 255),
            'red': (255, 0, 0),
            'darkBlue': (0, 0, 128),
            'darkCyan': (0, 128, 128),
            'darkGreen': (0, 128, 0),
            'darkMagenta': (128, 0, 128),
            'darkRed': (128, 0, 0),
            'darkYellow': (128, 128, 0),
            'darkGray': (128, 128, 128),
            'lightGray': (192, 192, 192),
            'black': (0, 0, 0),
            'white': (255, 255, 255),
        }
        
        # Find the closest color
        closest_color = None
        min_distance = float('inf')
        
        for color_name, color_rgb in highlight_colors.items():
            # Calculate Euclidean distance
            distance = (
                (r - color_rgb[0]) ** 2 + 
                (g - color_rgb[1]) ** 2 + 
                (b - color_rgb[2]) ** 2
            ) ** 0.5
            
            if distance < min_distance:
                min_distance = distance
                closest_color = color_name
        
        return closest_color or 'yellow'  # Default to yellow if no match
    
    def _fill_table_with_data(self, table: Table, headers: List[str], rows: List[List[str]], style: str = "default") -> None:
        """
        Fill a PowerPoint table with data and apply formatting.
        
        Args:
            table: PowerPoint table to fill.
            headers: List of column headers.
            rows: List of rows, where each row is a list of cell values.
            style: Style preset to apply to the table.
        """
        # Get the style preset
        style_preset = self.TABLE_STYLES.get(style, self.TABLE_STYLES["default"])
        
        # Ensure the table has the correct number of rows and columns
        actual_rows = len(table.rows)
        actual_cols = len(table.columns)
        needed_rows = len(rows) + 1  # +1 for header row
        needed_cols = len(headers)
        
        # Validate dimensions
        if needed_rows > actual_rows or needed_cols > actual_cols:
            logger.warning(
                f"Table dimensions mismatch: needed {needed_rows}x{needed_cols}, "
                f"actual {actual_rows}x{actual_cols}. "
                f"Data may be truncated."
            )
        
        # Add headers
        for col_idx, header in enumerate(headers):
            if col_idx < actual_cols:
                cell = table.cell(0, col_idx)
                self._format_table_cell(cell, header, is_header=True, style_preset=style_preset)
        
        # Add rows
        for row_idx, row_data in enumerate(rows):
            if row_idx + 1 < actual_rows:  # +1 to skip header row
                for col_idx, cell_value in enumerate(row_data):
                    if col_idx < actual_cols and col_idx < len(headers):  # Ensure we don't exceed headers
                        cell = table.cell(row_idx + 1, col_idx)
                        is_alternate_row = style_preset.get("banded_rows", False) and row_idx % 2 == 1
                        is_alternate_col = style_preset.get("banded_cols", False) and col_idx % 2 == 1
                        
                        self._format_table_cell(
                            cell, 
                            cell_value, 
                            is_header=False, 
                            is_alternate_row=is_alternate_row,
                            is_alternate_col=is_alternate_col,
                            style_preset=style_preset
                        )
        
        # Apply table styling
        self._apply_table_style(table, style_preset)
    
    def _format_table_cell(self, cell: _Cell, text: str, is_header: bool = False, 
                          is_alternate_row: bool = False, is_alternate_col: bool = False,
                          style_preset: Dict[str, Any] = None) -> None:
        """
        Format a table cell with text and styling.
        
        Args:
            cell: PowerPoint table cell to format.
            text: Text to add to the cell.
            is_header: Whether this is a header cell.
            is_alternate_row: Whether this cell is in an alternate row.
            is_alternate_col: Whether this cell is in an alternate column.
            style_preset: Style preset to apply to the cell.
        """
        if style_preset is None:
            style_preset = self.TABLE_STYLES["default"]
        
        # Clear any existing text
        text_frame = cell.text_frame
        text_frame.clear()
        
        # Add the text with formatting support
        self._add_formatted_text(text_frame, text)
        
        # Set vertical alignment to middle
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # Set text alignment
        for paragraph in text_frame.paragraphs:
            if is_header:
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                paragraph.alignment = PP_ALIGN.LEFT
            
            # Ensure consistent font size
            for run in paragraph.runs:
                if not run.font.size:
                    run.font.size = Pt(10)
        
        # Apply background color
        if is_header and style_preset.get("header_bg"):
            color = style_preset["header_bg"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._hex_to_rgb(color)
        elif is_alternate_row and style_preset.get("accent_color"):
            color = style_preset["accent_color"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._hex_to_rgb(color)
        
        # Apply text color
        if is_header and style_preset.get("header_text"):
            color = style_preset["header_text"]
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = self._hex_to_rgb(color)
        elif not is_header and style_preset.get("body_text"):
            color = style_preset["body_text"]
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = self._hex_to_rgb(color)
    
    def _apply_table_style(self, table: Table, style_preset: Dict[str, Any]) -> None:
        """
        Apply styling to a PowerPoint table.
        
        Args:
            table: PowerPoint table to style.
            style_preset: Style preset to apply to the table.
        """
        # Get style properties
        border_color = style_preset.get("border_color")
        border_width = style_preset.get("border_width", Pt(1))
        
        # Check for empty table
        if len(table.columns) == 0 or len(table.rows) == 0:
            logger.warning("Cannot apply style to empty table (no rows or columns)")
            return
        
        # Apply consistent row heights and column widths
        total_width = sum(column.width for column in table.columns)
        total_height = sum(row.height for row in table.rows)
        
        # Distribute column widths evenly
        col_width = total_width / len(table.columns)
        for column in table.columns:
            column.width = col_width
        
        # Set first row height slightly taller
        row_height = total_height / len(table.rows)
        if len(table.rows) > 1:
            # Use row_height * 1.2 directly instead of max() to avoid type comparison issues
            table.rows[0].height = row_height * 1.2
            for i in range(1, len(table.rows)):
                table.rows[i].height = row_height
        
        # Apply borders if specified
        if border_color:
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    # Top border
                    cell.border_top.color.rgb = self._hex_to_rgb(border_color)
                    cell.border_top.width = border_width
                    
                    # Bottom border
                    cell.border_bottom.color.rgb = self._hex_to_rgb(border_color)
                    cell.border_bottom.width = border_width
                    
                    # Left border
                    cell.border_left.color.rgb = self._hex_to_rgb(border_color)
                    cell.border_left.width = border_width
                    
                    # Right border
                    cell.border_right.color.rgb = self._hex_to_rgb(border_color)
                    cell.border_right.width = border_width
    
    @staticmethod
    def _hex_to_rgb(hex_value: str) -> RGBColor:
        """
        Convert a hex color string to an RGBColor object.
        
        Args:
            hex_value: Hex color string (with or without #).
            
        Returns:
            RGBColor object.
        """
        # Remove '#' if present
        if hex_value.startswith('#'):
            hex_value = hex_value[1:]
        
        # Ensure 6 digits
        if len(hex_value) == 3:
            hex_value = ''.join(c + c for c in hex_value)
        
        # Convert to RGB
        try:
            r = int(hex_value[0:2], 16)
            g = int(hex_value[2:4], 16)
            b = int(hex_value[4:6], 16)
            return RGBColor(r, g, b)
        except (ValueError, IndexError):
            # Default to black if color is invalid
            return RGBColor(0, 0, 0)