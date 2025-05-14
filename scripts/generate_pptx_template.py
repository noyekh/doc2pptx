from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.dml import MSO_FILL
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
import uuid
import datetime

def analyze_pptx_layouts(pptx_path, verbose=True):

    """
    Analyzes a PPTX file, extracts and prints information about all layouts,
    and returns structured data about them.
    
    Args:
        pptx_path: Path to the PPTX file
        verbose: Whether to print detailed information
        
    Returns:
        Dictionary containing structured information about all layouts
    """
    prs = Presentation(pptx_path)
    
    if verbose:
        print(f"\n=== ANALYZING LAYOUTS IN: {pptx_path} ===")
        print(f"Presentation has {len(prs.slide_layouts)} slide layouts and {len(prs.slides)} slides")
    
    # Store layout information
    layouts_info = {
        "slide_layouts": [],
        "slides": []
    }
    
    # Analyze slide layouts
    for layout_idx, layout in enumerate(prs.slide_layouts):
        layout_info = {
            "idx": layout_idx,
            "name": layout.name,
            "placeholders": []
        }
        
        if verbose:
            print(f"\n--- Layout {layout_idx}: {layout.name} ---")
            print(f"  Number of placeholders: {len(layout.placeholders)}")
        
        # Analyze placeholders in the layout
        for ph_idx, placeholder in enumerate(layout.placeholders):
            ph_info = {
                "idx": ph_idx,
                "name": placeholder.name,
                "shape_type": placeholder.shape_type,
                "placeholder_type": placeholder.placeholder_format.type if hasattr(placeholder, "placeholder_format") else None,
                "width": placeholder.width.inches,
                "height": placeholder.height.inches,
                "left": placeholder.left.inches,
                "top": placeholder.top.inches
            }
            
            # Try to get text properties if available
            try:
                if hasattr(placeholder, "text") and placeholder.text:
                    ph_info["text"] = placeholder.text
                
                if hasattr(placeholder, "text_frame"):
                    ph_info["vertical_anchor"] = placeholder.text_frame.vertical_anchor
                    ph_info["word_wrap"] = placeholder.text_frame.word_wrap
                    ph_info["auto_size"] = placeholder.text_frame.auto_size
                    
                    # Check paragraph properties
                    if placeholder.text_frame.paragraphs:
                        para = placeholder.text_frame.paragraphs[0]
                        ph_info["alignment"] = para.alignment
                        
                        # Check font properties
                        if para.runs:
                            font = para.runs[0].font
                            ph_info["font_name"] = font.name
                            ph_info["font_size"] = font.size.pt if font.size else None
                            ph_info["font_bold"] = font.bold
                            ph_info["font_italic"] = font.italic
                            
                            # Get color information safely
                            if hasattr(font.color, "rgb") and font.color.rgb:
                                try:
                                    # Access RGB values safely
                                    rgb = font.color.rgb
                                    ph_info["font_color"] = (
                                        rgb.r if hasattr(rgb, "r") else 0,
                                        rgb.g if hasattr(rgb, "g") else 0,
                                        rgb.b if hasattr(rgb, "b") else 0
                                    )
                                except Exception as e:
                                    ph_info["font_color_error"] = str(e)
            except Exception as e:
                ph_info["text_error"] = str(e)
            
            # Get fill information
            try:
                if hasattr(placeholder, "fill"):
                    ph_info["fill_type"] = placeholder.fill.type
                    
                    if placeholder.fill.type == MSO_FILL.SOLID:
                        if hasattr(placeholder.fill.fore_color, "rgb") and placeholder.fill.fore_color.rgb:
                            try:
                                # Access RGB values safely
                                rgb = placeholder.fill.fore_color.rgb
                                ph_info["fill_color"] = (
                                    rgb.r if hasattr(rgb, "r") else 0,
                                    rgb.g if hasattr(rgb, "g") else 0,
                                    rgb.b if hasattr(rgb, "b") else 0
                                )
                            except Exception as e:
                                ph_info["fill_color_error"] = str(e)
            except Exception as e:
                ph_info["fill_error"] = str(e)
            
            layout_info["placeholders"].append(ph_info)
            
            if verbose:
                print(f"    Placeholder {ph_idx}:")
                print(f"      Name: {placeholder.name}")
                print(f"      Type: {placeholder.shape_type}")
                print(f"      Placeholder Type: {placeholder.placeholder_format.type if hasattr(placeholder, 'placeholder_format') else 'None'}")
                print(f"      Position: Left={placeholder.left.inches}in, Top={placeholder.top.inches}in")
                print(f"      Size: Width={placeholder.width.inches}in, Height={placeholder.height.inches}in")
                
                # Print text content if any
                if hasattr(placeholder, "text") and placeholder.text:
                    print(f"      Text: {placeholder.text[:50]}...")
                
                # Print fill information
                if hasattr(placeholder, "fill"):
                    print(f"      Fill Type: {placeholder.fill.type}")
                    # Safely print color information if available
                    if placeholder.fill.type == MSO_FILL.SOLID and hasattr(placeholder.fill.fore_color, "rgb"):
                        try:
                            rgb = placeholder.fill.fore_color.rgb
                            if rgb and hasattr(rgb, "r") and hasattr(rgb, "g") and hasattr(rgb, "b"):
                                print(f"      Fill Color: RGB({rgb.r}, {rgb.g}, {rgb.b})")
                        except:
                            print(f"      Fill Color: [Color information not accessible]")
        
        layouts_info["slide_layouts"].append(layout_info)
    
    # Analyze actual slides in the presentation
    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {
            "idx": slide_idx,
            "layout_name": slide.slide_layout.name,
            "shapes": []
        }
        
        if verbose:
            print(f"\n=== Slide {slide_idx+1} ===")
            print(f"  Uses Layout: {slide.slide_layout.name}")
            print(f"  Number of shapes: {len(slide.shapes)}")
        
        # Analyze shapes in the slide
        for shape_idx, shape in enumerate(slide.shapes):
            shape_info = {
                "idx": shape_idx,
                "name": shape.name,
                "shape_type": shape.shape_type,
                "width": shape.width.inches,
                "height": shape.height.inches,
                "left": shape.left.inches,
                "top": shape.top.inches
            }
            
            # Check if it's a placeholder
            if shape.is_placeholder:
                shape_info["is_placeholder"] = True
                shape_info["placeholder_type"] = shape.placeholder_format.type
                
                # Find which layout placeholder this corresponds to
                for layout_info in layouts_info["slide_layouts"]:
                    if layout_info["name"] == slide.slide_layout.name:
                        for ph in layout_info["placeholders"]:
                            if ph["placeholder_type"] == shape.placeholder_format.type:
                                shape_info["layout_placeholder_idx"] = ph["idx"]
                                break
                        break
            
            # Try to get text properties if available
            try:
                if hasattr(shape, "text") and shape.text:
                    shape_info["text"] = shape.text
                
                if hasattr(shape, "text_frame"):
                    shape_info["vertical_anchor"] = shape.text_frame.vertical_anchor
                    shape_info["word_wrap"] = shape.text_frame.word_wrap
                    shape_info["auto_size"] = shape.text_frame.auto_size
                    
                    # Check paragraph properties
                    if shape.text_frame.paragraphs:
                        para = shape.text_frame.paragraphs[0]
                        shape_info["alignment"] = para.alignment
                        
                        # Check font properties
                        if para.runs:
                            font = para.runs[0].font
                            shape_info["font_name"] = font.name
                            shape_info["font_size"] = font.size.pt if font.size else None
                            shape_info["font_bold"] = font.bold
                            shape_info["font_italic"] = font.italic
                            
                            # Get color information safely
                            if hasattr(font.color, "rgb") and font.color.rgb:
                                try:
                                    # Access RGB values safely
                                    rgb = font.color.rgb
                                    shape_info["font_color"] = (
                                        rgb.r if hasattr(rgb, "r") else 0,
                                        rgb.g if hasattr(rgb, "g") else 0,
                                        rgb.b if hasattr(rgb, "b") else 0
                                    )
                                except Exception as e:
                                    shape_info["font_color_error"] = str(e)
            except Exception as e:
                shape_info["text_error"] = str(e)
            
            # Get fill information
            try:
                if hasattr(shape, "fill"):
                    shape_info["fill_type"] = shape.fill.type
                    
                    if shape.fill.type == MSO_FILL.SOLID:
                        if hasattr(shape.fill.fore_color, "rgb") and shape.fill.fore_color.rgb:
                            try:
                                # Access RGB values safely
                                rgb = shape.fill.fore_color.rgb
                                shape_info["fill_color"] = (
                                    rgb.r if hasattr(rgb, "r") else 0,
                                    rgb.g if hasattr(rgb, "g") else 0,
                                    rgb.b if hasattr(rgb, "b") else 0
                                )
                            except Exception as e:
                                shape_info["fill_color_error"] = str(e)
            except Exception as e:
                shape_info["fill_error"] = str(e)
            
            slide_info["shapes"].append(shape_info)
            
            if verbose:
                print(f"    Shape {shape_idx}:")
                print(f"      Name: {shape.name}")
                print(f"      Type: {shape.shape_type}")
                
                if shape.is_placeholder:
                    print(f"      Is Placeholder: True")
                    print(f"      Placeholder Type: {shape.placeholder_format.type}")
                
                print(f"      Position: Left={shape.left.inches}in, Top={shape.top.inches}in")
                print(f"      Size: Width={shape.width.inches}in, Height={shape.height.inches}in")
                
                # Print text content if any
                if hasattr(shape, "text") and shape.text:
                    text_preview = shape.text[:50]
                    if len(shape.text) > 50:
                        text_preview += "..."
                    print(f"      Text: {text_preview}")
                
                # Print fill information
                if hasattr(shape, "fill"):
                    print(f"      Fill Type: {shape.fill.type}")
                    # Safely print color information if available
                    if shape.fill.type == MSO_FILL.SOLID and hasattr(shape.fill.fore_color, "rgb"):
                        try:
                            rgb = shape.fill.fore_color.rgb
                            if rgb and hasattr(rgb, "r") and hasattr(rgb, "g") and hasattr(rgb, "b"):
                                print(f"      Fill Color: RGB({rgb.r}, {rgb.g}, {rgb.b})")
                        except:
                            print(f"      Fill Color: [Color information not accessible]")
        
        layouts_info["slides"].append(slide_info)
    
    return layouts_info



def create_editable_template(source_pptx_path, output_path=None):
    """
    Creates an editable template based on a source PPTX file.
    Each slide is converted to a unique layout in the template.
    
    Args:
        source_pptx_path: Path to the source PPTX file
        output_path: Path where the template will be saved. If None, a default name will be used.
        
    Returns:
        Path to the created template and layout information dictionary
    """
    if output_path is None:
        # Generate a default output path
        basename = os.path.splitext(os.path.basename(source_pptx_path))[0]
        output_path = f"{basename}_template_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    
    # First, analyze the source presentation
    layouts_info = analyze_pptx_layouts(source_pptx_path, verbose=False)
    
    # Create a new presentation for the template
    template_prs = Presentation()
    
    # Dictionary to store layout IDs for the template
    template_layouts = {}
    
    # Process each slide in the source presentation
    for slide_idx, slide_info in enumerate(layouts_info["slides"]):
        # Create a new slide layout in the template
        layout_name = f"Custom Layout {slide_idx+1} (from {slide_info['layout_name']})"
        
        # Add a new slide layout
        slide_layout = template_prs.slide_layouts.add_slide_layout()
        
        # We'll keep track of shapes we've processed
        processed_shapes = []
        
        # Process each shape in the slide
        for shape_info in slide_info["shapes"]:
            # Skip placeholders from the original slide layout
            if "is_placeholder" in shape_info and shape_info["is_placeholder"]:
                continue
            
            # For non-placeholder shapes, create a placeholder in the template layout
            left = Inches(shape_info["left"])
            top = Inches(shape_info["top"])
            width = Inches(shape_info["width"])
            height = Inches(shape_info["height"])
            
            placeholder_type = PP_PLACEHOLDER.OBJECT  # Default placeholder type
            
            # Create a descriptive name for the placeholder
            if "text" in shape_info:
                # Use the first few words of the text as part of the name
                text_sample = shape_info["text"][:20].replace("\n", " ").strip()
                name = f"Placeholder {shape_info['idx']+1} ({text_sample}...)"
            else:
                name = f"Placeholder {shape_info['idx']+1}"
            
            # Create a placeholder in the layout
            try:
                placeholder = slide_layout.placeholders.add_placeholder(
                    ph_type=placeholder_type,
                    name=name,
                    left=left,
                    top=top,
                    width=width,
                    height=height
                )
                
                # Set placeholder properties based on the original shape
                if "text" in shape_info:
                    placeholder.text = shape_info["text"]
                
                # Set text properties if available
                if hasattr(placeholder, "text_frame"):
                    if "vertical_anchor" in shape_info:
                        placeholder.text_frame.vertical_anchor = shape_info["vertical_anchor"]
                    if "word_wrap" in shape_info:
                        placeholder.text_frame.word_wrap = shape_info["word_wrap"]
                    if "auto_size" in shape_info:
                        placeholder.text_frame.auto_size = shape_info["auto_size"]
                    
                    # Set paragraph properties
                    if placeholder.text_frame.paragraphs:
                        para = placeholder.text_frame.paragraphs[0]
                        if "alignment" in shape_info:
                            para.alignment = shape_info["alignment"]
                        
                        # Set font properties
                        if para.runs:
                            font = para.runs[0].font
                            if "font_name" in shape_info:
                                font.name = shape_info["font_name"]
                            if "font_size" in shape_info and shape_info["font_size"]:
                                font.size = Pt(shape_info["font_size"])
                            if "font_bold" in shape_info:
                                font.bold = shape_info["font_bold"]
                            if "font_italic" in shape_info:
                                font.italic = shape_info["font_italic"]
                            
                            # Set font color
                            if "font_color" in shape_info:
                                r, g, b = shape_info["font_color"]
                                font.color.rgb = RGBColor(r, g, b)
                
                # Set fill properties if available
                if hasattr(placeholder, "fill") and "fill_type" in shape_info:
                    placeholder.fill.solid()
                    if "fill_color" in shape_info:
                        r, g, b = shape_info["fill_color"]
                        placeholder.fill.fore_color.rgb = RGBColor(r, g, b)
                
                processed_shapes.append(placeholder)
            except Exception as e:
                print(f"Error creating placeholder for shape {shape_info['idx']} on slide {slide_idx+1}: {e}")
        
        # Generate a unique ID for this layout
        layout_id = str(uuid.uuid4())
        template_layouts[layout_id] = {
            "layout_idx": len(template_layouts),
            "layout_name": layout_name,
            "source_slide_idx": slide_idx,
            "placeholders": [
                {
                    "idx": idx,
                    "name": shape.name,
                    "original_shape_idx": shape_info["idx"],
                    "has_text": hasattr(shape, "text") and bool(shape.text)
                }
                for idx, (shape, shape_info) in enumerate(zip(processed_shapes, [
                    s for s in slide_info["shapes"] 
                    if not ("is_placeholder" in s and s["is_placeholder"])
                ]))
            ]
        }
        
        # Create a slide using this layout as a demonstration
        slide = template_prs.slides.add_slide(slide_layout)
        
        # Add a title to the slide indicating this is a template
        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.5)
        )
        title_shape.text = f"Template Slide {slide_idx+1} (from source slide {slide_idx+1})"
        
        # Make sure the text frame has paragraphs
        if title_shape.text_frame.paragraphs:
            # Add bold formatting
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(14)
        
        # Add instructions at the bottom of the slide
        instructions = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5), Inches(9), Inches(0.5)
        )
        instructions.text = f"This is a template generated from slide {slide_idx+1}. Use the update_layout_text() function to modify text."
        
        # Make sure the text frame has paragraphs
        if instructions.text_frame.paragraphs:
            # Add italic formatting
            for paragraph in instructions.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.italic = True
                    run.font.size = Pt(10)
    
    # Save the template
    template_prs.save(output_path)
    
    print(f"\n=== Template Creation Complete ===")
    print(f"Created template with {len(template_layouts)} custom layouts")
    print(f"Saved to: {output_path}")
    
    # Print information about the template layouts
    print("\n=== Template Layouts Information ===")
    for layout_id, layout_info in template_layouts.items():
        print(f"\nLayout ID: {layout_id}")
        print(f"  Name: {layout_info['layout_name']}")
        print(f"  Based on source slide: {layout_info['source_slide_idx'] + 1}")
        print(f"  Placeholders: {len(layout_info['placeholders'])}")
        
        for ph_info in layout_info['placeholders']:
            print(f"    Placeholder {ph_info['idx']}:")
            print(f"      Name: {ph_info['name']}")
            print(f"      Original Shape Index: {ph_info['original_shape_idx']}")
            print(f"      Has Text: {ph_info['has_text']}")
    
    return output_path, template_layouts

def update_layout_text(template_path, slide_idx, placeholder_idx, new_text, new_color=None, output_path=None):
    """
    Update the text and optionally the color in a placeholder in the template
    
    Args:
        template_path: Path to the template PPTX
        slide_idx: Index of the slide (1-based)
        placeholder_idx: Index of the placeholder (1-based)
        new_text: New text to set
        new_color: Optional RGBColor to set
        output_path: Path to save the modified template. If None, modifies the original.
    
    Returns:
        True if successful, False otherwise
    """
    if output_path is None:
        output_path = template_path
    
    prs = Presentation(template_path)
    
    # Convert to 0-based indices
    slide_idx = slide_idx - 1
    placeholder_idx = placeholder_idx - 1
    
    if slide_idx < 0 or slide_idx >= len(prs.slides):
        print(f"Error: Slide index {slide_idx+1} is out of range")
        return False
    
    slide = prs.slides[slide_idx]
    
    # Get all non-placeholder shapes or shapes that we added as actual content
    shapes = [shape for shape in slide.shapes]
    
    if placeholder_idx < 0 or placeholder_idx >= len(shapes):
        print(f"Error: Shape index {placeholder_idx+1} on slide {slide_idx+1} is out of range")
        print(f"  This slide has {len(shapes)} shapes")
        return False
    
    shape = shapes[placeholder_idx]
    
    if not hasattr(shape, "text_frame"):
        print(f"Error: Shape {placeholder_idx+1} on slide {slide_idx+1} does not have a text frame")
        return False
    
    # Update text
    if shape.has_text_frame:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = new_text
        
        # Update color if specified
        if new_color:
            run.font.color.rgb = new_color
    else:
        print(f"Warning: Shape {placeholder_idx+1} on slide {slide_idx+1} doesn't support text directly")
        return False
    
    prs.save(output_path)
    print(f"Updated text of shape {placeholder_idx+1} on slide {slide_idx+1}")
    return True

def create_presentation_from_template(template_path, layout_info, output_path, content_dict):
    """
    Create a presentation using a template and provided content
    
    Args:
        template_path: Path to the template PPTX
        layout_info: Layout information dictionary from create_editable_template
        output_path: Path to save the new presentation
        content_dict: Dictionary mapping {slide_idx: {placeholder_idx: {"text": text, "color": color}}}
    
    Returns:
        Path to the created presentation
    """
    prs = Presentation(template_path)
    
    for slide_idx, placeholders in content_dict.items():
        # Convert to 0-based index
        slide_0idx = slide_idx - 1
        
        if slide_0idx < 0 or slide_0idx >= len(prs.slides):
            print(f"Warning: Slide index {slide_idx} is out of range, skipping")
            continue
        
        slide = prs.slides[slide_0idx]
        
        for placeholder_idx, content in placeholders.items():
            # Convert to 0-based index
            ph_0idx = placeholder_idx - 1
            
            shapes = [shape for shape in slide.shapes]
            
            if ph_0idx < 0 or ph_0idx >= len(shapes):
                print(f"Warning: Shape index {placeholder_idx} on slide {slide_idx} is out of range, skipping")
                continue
            
            shape = shapes[ph_0idx]
            
            if not hasattr(shape, "text_frame"):
                print(f"Warning: Shape {placeholder_idx} on slide {slide_idx} does not have a text frame, skipping")
                continue
            
            # Update text
            if "text" in content:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = content["text"]
                
                # Update color if specified
                if "color" in content:
                    run.font.color.rgb = content["color"]
    
    prs.save(output_path)
    print(f"Created presentation from template: {output_path}")
    return output_path

def extract_slide_shapes(pptx_path, verbose=True):
    """
    Extract all shapes from all slides in a presentation and organize them for easy updating
    
    Args:
        pptx_path: Path to the PPTX file
        verbose: Whether to print detailed information
        
    Returns:
        Dictionary with shape information by slide and index
    """
    prs = Presentation(pptx_path)
    
    result = {}
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_key = slide_idx + 1  # 1-based indexing for user interface
        result[slide_key] = {}
        
        if verbose:
            print(f"\nSlide {slide_key}:")
            
        for shape_idx, shape in enumerate(slide.shapes):
            shape_key = shape_idx + 1  # 1-based indexing
            
            shape_info = {
                "type": shape.shape_type,
                "name": shape.name,
                "has_text": hasattr(shape, "text"),
                "is_placeholder": shape.is_placeholder if hasattr(shape, "is_placeholder") else False,
            }
            
            if hasattr(shape, "text"):
                shape_info["text"] = shape.text
                
            if hasattr(shape, "fill") and shape.fill.type:
                shape_info["fill_type"] = shape.fill.type
            
            result[slide_key][shape_key] = shape_info
            
            if verbose:
                print(f"  Shape {shape_key}: {shape.name}")
                if hasattr(shape, "text") and shape.text:
                    print(f"    Text: {shape.text[:50]}..." if len(shape.text) > 50 else f"    Text: {shape.text}")
                if hasattr(shape, "fill") and shape.fill.type:
                    print(f"    Fill Type: {shape.fill.type}")
    
    return result

# Usage example
if __name__ == "__main__":
    # Path to your PPTX file
    pptx_file = "slide6.pptx"
    
    # Step 1: Analyze the PPTX file
    analyze_pptx_layouts(pptx_file)
    
    # Step 2: Extract shapes for easier manipulation
    shapes_info = extract_slide_shapes(pptx_file)
    print("\nExtracted Shapes:")
    for slide_idx, shapes in shapes_info.items():
        print(f"Slide {slide_idx}:")
        for shape_idx, shape_info in shapes.items():
            if shape_info.get("has_text", False):
                print(f"  Shape {shape_idx}: {shape_info['name']}")
                if "text" in shape_info:
                    print(f"    Text: {shape_info['text'][:30]}...")
    
    # Step 3: Create an editable template
    template_path, layout_info = create_editable_template(pptx_file)
    
    # Step 4: Example of updating text in the template
    # This updates the text in the third shape of the first slide
    update_layout_text(
        template_path, 
        slide_idx=1,   # First slide
        placeholder_idx=3,  # Third shape - adjust based on actual index 
        new_text="This is updated text!",
        new_color=RGBColor(0, 0, 255)  # Blue
    )
    
    # Step 5: Example of creating a new presentation from the template
    # Mapping of slide_idx -> shape_idx -> content
    content_dict = {
        1: {  # Slide 1
            3: {"text": "New content for box 1", "color": RGBColor(255, 0, 0)},  # Red
            4: {"text": "New content for box 2", "color": RGBColor(0, 128, 0)}   # Green
        }
    }
    
    create_presentation_from_template(
        template_path,
        layout_info,
        "final_presentation.pptx",
        content_dict
    )