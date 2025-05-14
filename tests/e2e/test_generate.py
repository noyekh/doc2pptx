"""End-to-end tests : JSON → Builder → CLI → fichier PPTX.

Ces scénarios valident l’intégration complète sans dépendance réseau :
lecture du JSON, sélection des layouts, génération du deck et exécution
de la CLI Typer (pipeline « utilisateur final »).

La suite vise **> 90 %** de couverture pour `src/doc2pptx/cli.py` *et*
exerce les chemins principaux de `PPTBuilder`.
"""

from __future__ import annotations

import json
import logging
import subprocess
import sys
from pathlib import Path
from typing import List

import pytest
from pptx import Presentation as PptxPresentation
from typer.testing import CliRunner

from doc2pptx.cli import app as cli_app
from doc2pptx.core.models import Presentation
from doc2pptx.ingest.json_loader import load_sections
from doc2pptx.ppt.builder import PPTBuilder
from doc2pptx.ppt.overflow import OverflowHandler
from uuid import uuid4


# --------------------------------------------------------------------------- fixtures
@pytest.fixture(scope="module")
def runner() -> CliRunner:
    """Typer runner partagé entre les tests CLI."""
    return CliRunner(mix_stderr=False)


@pytest.fixture()
def presentation_model(sample_json: Path, base_template: Path) -> Presentation:
    """Charge le JSON d’exemple en modèle Pydantic et ajoute un template."""
    data = json.loads(sample_json.read_text(encoding="utf-8"))
    data["template_path"] = str(base_template)
    return Presentation.model_validate(data)


@pytest.fixture()
def builder(base_template: Path) -> PPTBuilder:
    """Instance prête à l’emploi du builder, basée sur le template de test."""
    return PPTBuilder(template_path=base_template)


# --------------------------------------------------------------------- helpers internes
def _open_pptx(path: Path) -> PptxPresentation:
    assert path.exists(), "Le fichier PPTX n’a pas été généré."
    return PptxPresentation(path)


def _collect_visible_text(slides) -> List[str]:
    texts: List[str] = []
    for slide in slides:
        for shape in slide.shapes:
            if getattr(shape, "text", None):
                txt = shape.text.strip()
                if txt:
                    texts.append(txt.lower())
    return texts


# ----------------------------------------------------------------------------- tests
def test_builder_generates_file(
    builder: PPTBuilder, presentation_model: Presentation, tmp_output_path: Path
) -> None:
    """Le builder seul produit un fichier valide et non vide."""
    built = builder.build(presentation_model, tmp_output_path)
    ppt = _open_pptx(built)

    # 1. au moins autant de slides dans la sortie que dans le modèle
    expected = sum(len(s.slides) for s in presentation_model.sections)
    assert len(ppt.slides) >= expected

    # 2. les titres de section apparaissent quelque part
    slide_texts = _collect_visible_text(ppt.slides)
    for section in presentation_model.sections:
        assert any(section.title.lower() in t for t in slide_texts), (
            f"Titre de section « {section.title} » absent du deck"
        )


def test_cli_generate_happy_path(
    runner: CliRunner,
    sample_json: Path,
    base_template: Path,
    tmp_output_path: Path,
) -> None:
    """Chemin « utilisateur final » : appel CLI avec options complètes."""
    res = runner.invoke(
        cli_app,
        [
            "generate",
            str(sample_json),
            "--template",
            str(base_template),
            "--output",
            str(tmp_output_path),
        ],
    )
    assert res.exit_code == 0, res.stdout

    ppt = _open_pptx(tmp_output_path)
    assert len(ppt.slides) > 0


def test_cli_generate_default_output(
    runner: CliRunner, sample_json: Path, base_template: Path, tmp_path: Path
) -> None:
    """Quand --output est omis, le fichier porte le nom du JSON."""
    with runner.isolated_filesystem(temp_dir=tmp_path) as cwd:
        res = runner.invoke(
            cli_app,
            ["generate", str(sample_json), "--template", str(base_template)],
        )
        assert res.exit_code == 0, res.stdout

        expected = Path(cwd) / f"{sample_json.stem}.pptx"
        ppt = _open_pptx(expected)
        assert len(ppt.slides) > 0


def test_cli_error_bad_json(
    runner: CliRunner, base_template: Path, tmp_path: Path
) -> None:
    """Le JSON mal formé déclenche une erreur gérée (exit-code 1)."""
    bad_json = tmp_path / "bad.json"
    bad_json.write_text('{"oops": }', encoding="utf-8")

    res = runner.invoke(
        cli_app,
        ["generate", str(bad_json), "--template", str(base_template)],
    )
    assert res.exit_code == 1
    assert "Erreur de parsing JSON" in res.stdout


def test_invalid_template_path(tmp_path: Path) -> None:
    """PPTBuilder lève FileNotFoundError si le template est manquant."""
    missing = tmp_path / "ghost.pptx"
    with pytest.raises(FileNotFoundError):
        PPTBuilder(template_path=missing)


def test_invalid_layout_names_are_fallbacked(
    presentation_model: Presentation,
    base_template: Path,
    tmp_output_path: Path,
) -> None:
    """Des layouts inexistants sont remplacés par le fallback interne."""
    # force un faux layout
    for section in presentation_model.sections:
        for slide in section.slides:
            slide.layout_name = "DoesNotExist"

    builder = PPTBuilder(template_path=base_template)
    built = builder.build(presentation_model, tmp_output_path)
    ppt = _open_pptx(built)
    assert len(ppt.slides) > 0  # génération quand même


def test_overflow_warning_logged(
    base_template: Path, tmp_output_path: Path, caplog: pytest.LogCaptureFixture
) -> None:
    """Texte très long ⇒ logger.warning d’OverflowHandler.will_text_overflow."""
    from doc2pptx.core.models import Section, Slide, SlideBlock, SlideContent, ContentType

    long_text = "A" * 8_000  # déraisonnablement long

    model = Presentation(
        id="ovf",
        title="overflow",
        template_path=base_template,
        sections=[
            Section(
                id="s1",
                title="Overflow",
                type="content",
                slides=[
                    Slide(
                        id="sl1",
                        title="t",
                        layout_name="Titre et texte",
                        blocks=[
                            SlideBlock(
                                id="b1",
                                title="t",
                                content=SlideContent(content_type=ContentType.TEXT, text=long_text),
                            )
                        ],
                    )
                ],
            )
        ],
    )

    caplog.set_level(logging.WARNING)
    builder = PPTBuilder(template_path=base_template)
    builder.build(model, tmp_output_path)

    # un avertissement OverflowHandler doit être présent
    msgs = [rec.getMessage() for rec in caplog.records]
    assert any("overflow" in m.lower() for m in msgs)


def test_subprocess_entrypoint(
    sample_json: Path, base_template: Path, tmp_path: Path
) -> None:
    """On frappe `python -m doc2pptx.cli …` comme un utilisateur shell."""
    out = tmp_path / "via_subprocess.pptx"
    cmd = [
        sys.executable,
        "-m",
        "doc2pptx.cli",
        "generate",
        str(sample_json),
        "--template",
        str(base_template),
        "--output",
        str(out),
    ]

    proc = subprocess.run(cmd, capture_output=True, text=True)
    assert proc.returncode == 0, proc.stderr
    _open_pptx(out)  # existence + slide(s) vérifiés plus haut

@pytest.mark.parametrize("verbose,expected", [(False, logging.INFO), (True, logging.DEBUG)])
def test_configure_logging_sets_level(verbose: bool, expected: int, caplog) -> None:
    _configure_logging(verbose)
    # root logger mis à jour
    assert logging.getLogger().level == expected


@pytest.mark.parametrize(
    "sub_cmd, extra",
    [
        ("edit", ["--presentation", "dummy.pptx", "--command", "noop"]),
        ("prompt", ["--presentation", "dummy.pptx", "--nl", "noop"]),
    ],
)
def test_stub_commands_exit_1(tmp_path, runner: CliRunner, sub_cmd: str, extra: list[str]) -> None:
    dummy = tmp_path / "dummy.pptx"
    dummy.write_bytes(b"x")  # besoin d’un fichier existant
    res = runner.invoke(cli_app, [sub_cmd, *extra])
    assert res.exit_code == 1
    assert "pas encore disponible" in res.stdout.lower()


def test_version_command(runner: CliRunner) -> None:
    res = runner.invoke(cli_app, ["version"])
    # le numéro peut ne pas être résolu si package non installé, mais sortie non vide
    assert res.exit_code == 0 and res.stdout.strip()