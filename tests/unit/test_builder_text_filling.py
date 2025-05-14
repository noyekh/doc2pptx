"""
Unit tests for slide filling functionality in PPTBuilder.
"""
import pytest
from unittest.mock import MagicMock, patch, PropertyMock

from pptx.presentation import Presentation as PptxPresentation
from pptx.slide import Slide as PptxSlide
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from doc2pptx.core.models import (
    ContentType, 
    Presentation, 
    Section, 
    Slide, 
    SlideBlock, 
    SlideContent,
    SectionType,
    TableData,
    MermaidDiagram,
    ImageSource,
    ChartData
)
from doc2pptx.ppt.builder_v3 import PPTBuilder


@pytest.fixture
def builder():
    """Create a PPTBuilder instance."""
    return PPTBuilder()


@pytest.fixture
def mock_pptx_slide():
    """Create a mock PowerPoint slide with basic placeholders."""
    slide = MagicMock(spec=PptxSlide)
    
    # Create title placeholder
    title_placeholder = MagicMock()
    title_placeholder.is_placeholder = True
    title_placeholder.placeholder_format = MagicMock()
    title_placeholder.placeholder_format.type = PP_PLACEHOLDER.TITLE
    title_placeholder.placeholder_format.idx = 0
    title_placeholder.text_frame = MagicMock()
    
    # Create content placeholder
    content_placeholder = MagicMock()
    content_placeholder.is_placeholder = True
    content_placeholder.placeholder_format = MagicMock()
    content_placeholder.placeholder_format.type = PP_PLACEHOLDER.BODY
    content_placeholder.placeholder_format.idx = 1
    content_placeholder.text_frame = MagicMock()
    content_placeholder.text_frame.paragraphs = [MagicMock()]
    
    # Setup shapes collection
    slide.shapes = MagicMock()
    placeholders = [title_placeholder, content_placeholder]
    slide.shapes.__iter__.return_value = iter(placeholders)
    slide.shapes.__getitem__.side_effect = lambda idx: placeholders[idx]
    slide.shapes.__len__.return_value = len(placeholders)
    
    # Setup notes slide
    slide.notes_slide = MagicMock()
    slide.notes_slide.notes_text_frame = MagicMock()
    
    return slide


def test_fill_slide_title(builder, mock_pptx_slide):
    """Test filling a slide title."""
    builder._fill_slide_title(mock_pptx_slide, "Test Title")
    
    # Verify title was added
    title_placeholder = mock_pptx_slide.shapes[0]
    assert title_placeholder.placeholder_format.type == PP_PLACEHOLDER.TITLE
    title_placeholder.text_frame.clear.assert_called_once()


def test_fill_slide_title_with_empty_title(builder, mock_pptx_slide):
    """Test filling a slide with an empty title."""
    builder._fill_slide_title(mock_pptx_slide, None)
    
    # Verify no title was added
    title_placeholder = mock_pptx_slide.shapes[0]
    title_placeholder.text_frame.clear.assert_not_called()


def test_fill_slide_title_no_placeholder(builder):
    """Test filling a slide title when no title placeholder exists."""
    # Create mock slide with no title placeholder
    slide = MagicMock(spec=PptxSlide)
    
    # Create content placeholder only
    content_placeholder = MagicMock()
    content_placeholder.is_placeholder = True
    content_placeholder.placeholder_format = MagicMock()
    content_placeholder.placeholder_format.type = PP_PLACEHOLDER.BODY
    content_placeholder.placeholder_format.idx = 1
    content_placeholder.text_frame = MagicMock()
    
    # Setup shapes collection
    slide.shapes = MagicMock()
    placeholders = [content_placeholder]
    slide.shapes.__iter__.return_value = iter(placeholders)
    slide.shapes.__getitem__.side_effect = lambda idx: placeholders[idx] if idx < len(placeholders) else None
    slide.shapes.__len__.return_value = len(placeholders)
    
    # Test filling title
    with patch('doc2pptx.ppt.builder_v3.logger.warning') as mock_warning:
        builder._fill_slide_title(slide, "Test Title")
        
        # Should log a warning
        mock_warning.assert_called_once()


def test_fill_slide_with_notes(builder, mock_pptx_slide):
    """Test filling a slide with speaker notes."""
    # Create a slide with notes
    slide = Slide(
        id="slide1",
        title="Slide with Notes",
        layout_name="Titre et texte",
        blocks=[
            SlideBlock(
                id="block1",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Text content"
                )
            )
        ],
        notes="These are speaker notes."
    )
    
    # Create a section
    section = Section(
        id="section1",
        title="Test Section",
        type=SectionType.CONTENT,
        slides=[slide]
    )
    
    # Mock methods to avoid side effects
    builder._fill_content_slide = MagicMock()
    builder._fill_slide_title = MagicMock()
    
    # Fill the slide
    builder._fill_slide(mock_pptx_slide, slide, section)
    
    # Verify notes were added
    assert mock_pptx_slide.notes_slide.notes_text_frame.text == "These are speaker notes."


def test_fill_title_slide(builder, mock_pptx_slide):
    """Test filling a title slide."""
    # Add a subtitle placeholder to the mock slide
    subtitle_placeholder = MagicMock()
    subtitle_placeholder.is_placeholder = True
    subtitle_placeholder.placeholder_format = MagicMock()
    subtitle_placeholder.placeholder_format.type = PP_PLACEHOLDER.SUBTITLE
    subtitle_placeholder.placeholder_format.idx = 2
    subtitle_placeholder.text_frame = MagicMock()
    
    # Add the subtitle placeholder to the shapes collection
    placeholders = list(mock_pptx_slide.shapes)
    placeholders.append(subtitle_placeholder)
    mock_pptx_slide.shapes.__iter__.return_value = iter(placeholders)
    mock_pptx_slide.shapes.__getitem__.side_effect = lambda idx: placeholders[idx] if idx < len(placeholders) else None
    mock_pptx_slide.shapes.__len__.return_value = len(placeholders)
    
    # Create a title slide
    slide = Slide(
        id="slide1",
        title="Presentation Title",
        layout_name="Diapositive de titre",
        blocks=[
            SlideBlock(
                id="block1",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Subtitle Text"
                )
            )
        ]
    )
    
    # Mock _add_formatted_text to avoid side effects
    builder._add_formatted_text = MagicMock()
    
    # Fill the slide
    builder._fill_title_slide(mock_pptx_slide, slide)
    
    # Verify subtitle was added
    builder._add_formatted_text.assert_called_with(subtitle_placeholder.text_frame, "Subtitle Text")


def test_fill_content_slide(builder, mock_pptx_slide):
    """Test filling a standard content slide."""
    # Create a content slide
    slide = Slide(
        id="slide1",
        title="Content Slide",
        layout_name="Titre et texte",
        blocks=[
            SlideBlock(
                id="block1",
                title="Block Title",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Text content"
                )
            )
        ]
    )
    
    # Mock methods to avoid side effects
    builder._add_text_content_to_placeholder = MagicMock()
    
    # Fill the slide
    builder._fill_content_slide(mock_pptx_slide, slide)
    
    # Verify content was added
    content_placeholder = mock_pptx_slide.shapes[1]
    content_placeholder.text_frame.clear.assert_called_once()
    builder._add_text_content_to_placeholder.assert_called_once()


def test_fill_content_slide_multiple_blocks(builder, mock_pptx_slide):
    """Test filling a content slide with multiple blocks."""
    # Create a content slide with multiple blocks
    slide = Slide(
        id="slide1",
        title="Content Slide",
        layout_name="Titre et texte",
        blocks=[
            SlideBlock(
                id="block1",
                title="Block 1",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Text 1"
                )
            ),
            SlideBlock(
                id="block2",
                title="Block 2",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Text 2"
                )
            )
        ]
    )
    
    # Mock methods to avoid side effects
    builder._add_text_content_to_placeholder = MagicMock()
    
    # Fill the slide
    builder._fill_content_slide(mock_pptx_slide, slide)
    
    # Verify content was added for both blocks
    content_placeholder = mock_pptx_slide.shapes[1]
    content_placeholder.text_frame.clear.assert_called_once()
    content_placeholder.text_frame.add_paragraph.assert_called()
    assert builder._add_text_content_to_placeholder.call_count == 2


def test_fill_content_slide_bullet_points(builder, mock_pptx_slide):
    """Test filling a content slide with bullet points."""
    # Create a content slide with bullet points
    slide = Slide(
        id="slide1",
        title="Bullet Points",
        layout_name="Titre et texte",
        blocks=[
            SlideBlock(
                id="block1",
                title="Bullet Points",
                content=SlideContent(
                    content_type=ContentType.BULLET_POINTS,
                    bullet_points=["Point 1", "Point 2", "Point 3"],
                    as_bullets=True
                )
            )
        ]
    )
    
    # Mock methods to avoid side effects
    builder._add_bullet_points_to_placeholder = MagicMock()
    
    # Fill the slide
    builder._fill_content_slide(mock_pptx_slide, slide)
    
    # Verify bullet points were added
    content_placeholder = mock_pptx_slide.shapes[1]
    content_placeholder.text_frame.clear.assert_called_once()
    builder._add_bullet_points_to_placeholder.assert_called_once_with(
        content_placeholder,
        ["Point 1", "Point 2", "Point 3"],
        True
    )


def test_fill_content_slide_unsupported_content(builder, mock_pptx_slide):
    """Test filling a content slide with unsupported content type."""
    # Create a content slide with unsupported content
    slide = Slide(
        id="slide1",
        title="Unsupported Content",
        layout_name="Titre et texte",
        blocks=[
            SlideBlock(
                id="block1",
                title="Chart",
                content=SlideContent(
                    content_type=ContentType.CHART,
                    chart=ChartData(
                        chart_type="bar",
                        categories=["Cat1", "Cat2", "Cat3"],
                        series=[{"name": "Series1", "data": [1, 2, 3]}]
                    )
                )
            )
        ]
    )
    
    # Fill the slide
    builder._fill_content_slide(mock_pptx_slide, slide)
    
    # Verify fallback text was added
    content_placeholder = mock_pptx_slide.shapes[1]
    content_placeholder.text_frame.clear.assert_called_once()
    content_placeholder.text_frame.add_paragraph.assert_called()
    paragraph = content_placeholder.text_frame.add_paragraph.return_value
    assert "chart content not shown" in paragraph.text.lower()


def test_fill_table_slide(builder, mock_pptx_slide):
    """Test filling a slide with a table."""
    # Create a table slide
    slide = Slide(
        id="slide1",
        title="Table Slide",
        layout_name="Titre et tableau",
        blocks=[
            SlideBlock(
                id="block1",
                content=SlideContent(
                    content_type=ContentType.TABLE,
                    table=TableData(
                        headers=["Col1", "Col2"],
                        rows=[["A", "B"], ["C", "D"]]
                    )
                )
            )
        ]
    )
    
    # Mock methods to avoid side effects
    with patch.object(builder, '_add_table_to_slide', return_value=MagicMock()) as mock_add_table:
        # Fill the slide
        builder._fill_table_slide(mock_pptx_slide, slide)
        
        # Verify _add_table_to_slide was called with processed headers
        mock_add_table.assert_called_once()
        args = mock_add_table.call_args[0]
        assert args[0] == mock_pptx_slide
        assert args[1] == ["Col1", "Col2"]
        assert args[2] == [["A", "B"], ["C", "D"]]
        assert args[3] == "default"


def test_fill_table_slide_with_style(builder, mock_pptx_slide):
    """Test filling a slide with a table that has style information."""
    # Create a table slide with style
    slide = Slide(
        id="slide1",
        title="Table Slide with Style",
        layout_name="Titre et tableau",
        blocks=[
            SlideBlock(
                id="block1",
                content=SlideContent(
                    content_type=ContentType.TABLE,
                    table=TableData(
                        headers=["Col1", "Col2", "style:accent1"],
                        rows=[
                            ["A", "B", "C"],  # Match the number of columns in headers
                            ["D", "E", "F"]   # Match the number of columns in headers
                        ]
                    )
                )
            )
        ]
    )
    
    # Create a mock for _add_table_to_slide
    mock_table = MagicMock()
    with patch.object(builder, '_add_table_to_slide', return_value=mock_table) as mock_add_table:
        # Mock _get_style_from_headers to return "accent1" and verify call
        with patch.object(builder, '_get_style_from_headers', return_value="accent1") as mock_get_style:
            # Execute the method
            builder._fill_table_slide(mock_pptx_slide, slide)
            
            # Verify _get_style_from_headers was called
            mock_get_style.assert_called_once()
            
            # Verify _add_table_to_slide was called with processed headers
            mock_add_table.assert_called_once()
            args = mock_add_table.call_args[0]
            assert args[0] == mock_pptx_slide
            assert "style:accent1" not in args[1]  # Style marker should be removed
            assert args[3] == "accent1"  # Style should be passed


def test_fill_table_slide(builder, mock_pptx_slide):
    """Test filling a slide with a table."""
    # Create a table slide
    slide = Slide(
        id="slide1",
        title="Table Slide",
        layout_name="Titre et tableau",
        blocks=[
            SlideBlock(
                id="block1",
                content=SlideContent(
                    content_type=ContentType.TABLE,
                    table=TableData(
                        headers=["Col1", "Col2"],
                        rows=[["A", "B"], ["C", "D"]]
                    )
                )
            )
        ]
    )
    
    # Mock methods to avoid side effects
    builder._add_table_to_slide = MagicMock(return_value=MagicMock())
    
    # Patch _get_style_from_headers to return None (default style)
    with patch.object(builder, '_get_style_from_headers', return_value=None):
        # Fill the slide
        builder._fill_table_slide(mock_pptx_slide, slide)
        
        # Verify table was added
        builder._add_table_to_slide.assert_called_once()
        # Check that we're passing the headers, rows, and style ("default")
        args = builder._add_table_to_slide.call_args[0]
        assert args[0] == mock_pptx_slide
        assert args[1] == ["Col1", "Col2"]
        assert args[2] == [["A", "B"], ["C", "D"]]
        assert args[3] == "default"


def test_fill_column_layout_slide(builder, mock_pptx_slide):
    """Test filling a slide with multiple columns."""
    # Add two more content placeholders for columns
    column2_placeholder = MagicMock()
    column2_placeholder.is_placeholder = True
    column2_placeholder.placeholder_format = MagicMock()
    column2_placeholder.placeholder_format.type = PP_PLACEHOLDER.BODY
    column2_placeholder.placeholder_format.idx = 2
    column2_placeholder.text_frame = MagicMock()
    # Set left as an integer, not a MagicMock
    type(column2_placeholder).left = PropertyMock(return_value=300)
    
    column3_placeholder = MagicMock()
    column3_placeholder.is_placeholder = True
    column3_placeholder.placeholder_format = MagicMock()
    column3_placeholder.placeholder_format.type = PP_PLACEHOLDER.BODY
    column3_placeholder.placeholder_format.idx = 3
    column3_placeholder.text_frame = MagicMock()
    # Set left as an integer, not a MagicMock
    type(column3_placeholder).left = PropertyMock(return_value=500)
    
    # For the existing content placeholder, also set left
    type(mock_pptx_slide.shapes[1]).left = PropertyMock(return_value=100)
    
    # Add the placeholders to the shapes collection
    placeholders = list(mock_pptx_slide.shapes)
    placeholders.append(column2_placeholder)
    placeholders.append(column3_placeholder)
    mock_pptx_slide.shapes.__iter__.return_value = iter(placeholders)
    mock_pptx_slide.shapes.__getitem__.side_effect = lambda idx: placeholders[idx] if idx < len(placeholders) else None
    mock_pptx_slide.shapes.__len__.return_value = len(placeholders)
    
    # Create a column layout slide
    slide = Slide(
        id="slide1",
        title="Column Layout",
        layout_name="Titre et 3 colonnes",
        blocks=[
            SlideBlock(
                id="block1",
                title="Column 1",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Content for column 1"
                )
            ),
            SlideBlock(
                id="block2",
                title="Column 2",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Content for column 2"
                )
            ),
            SlideBlock(
                id="block3",
                title="Column 3",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Content for column 3"
                )
            )
        ]
    )
    
    # Mock methods to avoid side effects
    builder._add_block_to_placeholder = MagicMock()
    
    # Fill the slide
    builder._fill_column_layout_slide(mock_pptx_slide, slide)
    
    # Verify blocks were added to column placeholders
    assert builder._add_block_to_placeholder.call_count == 3
    
    # Verify they were added in the right order (sorted by left position)
    column_placeholders = [mock_pptx_slide.shapes[1], column2_placeholder, column3_placeholder]
    column_placeholders.sort(key=lambda p: p.left)
    
    for i, block in enumerate(slide.blocks):
        builder._add_block_to_placeholder.assert_any_call(column_placeholders[i], block)


def test_fill_column_layout_slide_more_blocks_than_columns(builder, mock_pptx_slide):
    """Test filling a column layout slide with more blocks than columns."""
    # Add one more content placeholder for column 2
    column2_placeholder = MagicMock()
    column2_placeholder.is_placeholder = True
    column2_placeholder.placeholder_format = MagicMock()
    column2_placeholder.placeholder_format.type = PP_PLACEHOLDER.BODY
    column2_placeholder.placeholder_format.idx = 2
    column2_placeholder.text_frame = MagicMock()
    column2_placeholder.left = 300  # Use direct attribute instead of PropertyMock
    
    # For the existing content placeholder, also set left
    mock_pptx_slide.shapes[1].left = 100  # Use direct attribute
    
    # Add the placeholder to the shapes collection
    placeholders = list(mock_pptx_slide.shapes)
    placeholders.append(column2_placeholder)
    mock_pptx_slide.shapes.__iter__.return_value = iter(placeholders)
    mock_pptx_slide.shapes.__getitem__.side_effect = lambda idx: placeholders[idx] if idx < len(placeholders) else None
    mock_pptx_slide.shapes.__len__.return_value = len(placeholders)
    
    # Create a column layout slide with more blocks than columns
    slide = Slide(
        id="slide1",
        title="Column Layout",
        layout_name="Titre et 3 colonnes",
        blocks=[
            SlideBlock(
                id="block1",
                title="Block 1",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Content for block 1"
                )
            ),
            SlideBlock(
                id="block2",
                title="Block 2",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Content for block 2"
                )
            ),
            SlideBlock(
                id="block3",
                title="Block 3",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Content for block 3"
                )
            ),
            SlideBlock(
                id="block4",
                title="Block 4",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Content for block 4"
                )
            )
        ]
    )
    
    # Mock methods to avoid side effects
    with patch.object(builder, '_add_block_to_placeholder') as mock_add_block:
        # Fill the slide
        builder._fill_column_layout_slide(mock_pptx_slide, slide)
        
        # Verify blocks were distributed among available columns
        assert mock_add_block.call_count >= 4  # At least 4 calls for 4 blocks
        
        # Verify first block goes to first placeholder
        mock_add_block.assert_any_call(mock_pptx_slide.shapes[1], slide.blocks[0])
        
        # Verify second block goes to second placeholder
        mock_add_block.assert_any_call(column2_placeholder, slide.blocks[1])


def test_fill_column_layout_slide_no_placeholders(builder, mock_pptx_slide):
    """Test filling a column layout slide with no column placeholders."""
    # Remove all content placeholders, leaving only title
    title_placeholder = mock_pptx_slide.shapes[0]
    mock_pptx_slide.shapes = MagicMock()
    mock_pptx_slide.shapes.__iter__.return_value = iter([title_placeholder])
    mock_pptx_slide.shapes.__getitem__.side_effect = lambda idx: title_placeholder if idx == 0 else None
    mock_pptx_slide.shapes.__len__.return_value = 1
    
    # Create a column layout slide
    slide = Slide(
        id="slide1",
        title="Column Layout",
        layout_name="Titre et 3 colonnes",
        blocks=[
            SlideBlock(
                id="block1",
                title="Column 1",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Content for column 1"
                )
            )
        ]
    )
    
    # Mock methods to avoid side effects
    builder._add_block_to_placeholder = MagicMock()
    
    # Fill the slide
    with patch('doc2pptx.ppt.builder_v3.logger.warning') as mock_warning:
        builder._fill_column_layout_slide(mock_pptx_slide, slide)
        
        # Should log a warning
        mock_warning.assert_called_once()
        
        # Should not attempt to add blocks
        builder._add_block_to_placeholder.assert_not_called()


def test_fill_image_layout_slide(builder, mock_pptx_slide):
    """Test filling a slide with image and text."""
    # Add an image placeholder
    image_placeholder = MagicMock()
    image_placeholder.is_placeholder = True
    image_placeholder.placeholder_format = MagicMock()
    image_placeholder.placeholder_format.type = PP_PLACEHOLDER.PICTURE
    image_placeholder.placeholder_format.idx = 2
    image_placeholder.text_frame = MagicMock()
    
    # Add the placeholder to the shapes collection
    placeholders = list(mock_pptx_slide.shapes)
    placeholders.append(image_placeholder)
    mock_pptx_slide.shapes.__iter__.return_value = iter(placeholders)
    mock_pptx_slide.shapes.__getitem__.side_effect = lambda idx: placeholders[idx] if idx < len(placeholders) else None
    mock_pptx_slide.shapes.__len__.return_value = len(placeholders)
    
    # Create an image layout slide
    slide = Slide(
        id="slide1",
        title="Image Layout",
        layout_name="Titre et texte 1 visuel gauche",
        blocks=[
            SlideBlock(
                id="block1",
                content=SlideContent(
                    content_type=ContentType.IMAGE,
                    image=ImageSource(
                        query="test image",
                        alt_text="Test image description"
                    )
                )
            ),
            SlideBlock(
                id="block2",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Text content to go with the image"
                )
            )
        ]
    )
    
    # Mock methods to avoid side effects
    builder._add_text_content_to_placeholder = MagicMock()
    
    # Fill the slide
    builder._fill_image_layout_slide(mock_pptx_slide, slide)
    
    # Verify text was added to content placeholder
    builder._add_text_content_to_placeholder.assert_called_once_with(
        mock_pptx_slide.shapes[1],
        "Text content to go with the image"
    )
    
    # Verify image description was added to image placeholder
    assert "[Image: Query: test image]" in image_placeholder.text_frame.text


def test_fill_chart_layout_slide(builder, mock_pptx_slide):
    """Test filling a slide with chart and text."""
    # Add a chart placeholder
    chart_placeholder = MagicMock()
    chart_placeholder.is_placeholder = True
    chart_placeholder.placeholder_format = MagicMock()
    chart_placeholder.placeholder_format.type = PP_PLACEHOLDER.CHART
    chart_placeholder.placeholder_format.idx = 2
    chart_placeholder.text_frame = MagicMock()
    
    # Add the placeholder to the shapes collection
    placeholders = list(mock_pptx_slide.shapes)
    placeholders.append(chart_placeholder)
    mock_pptx_slide.shapes.__iter__.return_value = iter(placeholders)
    mock_pptx_slide.shapes.__getitem__.side_effect = lambda idx: placeholders[idx] if idx < len(placeholders) else None
    mock_pptx_slide.shapes.__len__.return_value = len(placeholders)
    
    # Create a chart layout slide
    slide = Slide(
        id="slide1",
        title="Chart Layout",
        layout_name="Titre et texte 1 histogramme",
        blocks=[
            SlideBlock(
                id="block1",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Text content to go with the chart"
                )
            ),
            SlideBlock(
                id="block2",
                content=SlideContent(
                    content_type=ContentType.CHART,
                    chart=ChartData(
                        chart_type="bar",
                        categories=["Cat1", "Cat2", "Cat3"],
                        series=[{"name": "Series1", "data": [1, 2, 3]}],
                        title="Test Chart"
                    )
                )
            )
        ]
    )
    
    # Mock methods to avoid side effects
    builder._add_text_content_to_placeholder = MagicMock()
    
    # Fill the slide
    builder._fill_chart_layout_slide(mock_pptx_slide, slide)
    
    # Verify text was added to content placeholder
    builder._add_text_content_to_placeholder.assert_called_once_with(
        mock_pptx_slide.shapes[1],
        "Text content to go with the chart"
    )
    
    # Verify chart description was added to chart placeholder
    assert "[Chart: bar" in chart_placeholder.text_frame.text
    assert "Title: Test Chart" in chart_placeholder.text_frame.text


def test_fill_chart_layout_slide_with_mermaid(builder, mock_pptx_slide):
    """Test filling a chart layout slide with a mermaid diagram."""
    # Add a chart placeholder
    chart_placeholder = MagicMock()
    chart_placeholder.is_placeholder = True
    chart_placeholder.placeholder_format = MagicMock()
    chart_placeholder.placeholder_format.type = PP_PLACEHOLDER.CHART
    chart_placeholder.placeholder_format.idx = 2
    chart_placeholder.text_frame = MagicMock()
    
    # Add the placeholder to the shapes collection
    placeholders = list(mock_pptx_slide.shapes)
    placeholders.append(chart_placeholder)
    mock_pptx_slide.shapes.__iter__.return_value = iter(placeholders)
    mock_pptx_slide.shapes.__getitem__.side_effect = lambda idx: placeholders[idx] if idx < len(placeholders) else None
    mock_pptx_slide.shapes.__len__.return_value = len(placeholders)
    
    # Create a chart layout slide with mermaid content
    slide = Slide(
        id="slide1",
        title="Mermaid Diagram",
        layout_name="Titre et texte 1 histogramme",
        blocks=[
            SlideBlock(
                id="block1",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Text content to go with the diagram"
                )
            ),
            SlideBlock(
                id="block2",
                content=SlideContent(
                    content_type=ContentType.MERMAID,
                    mermaid=MermaidDiagram(
                        code="graph TD; A-->B; B-->C;",
                        caption="Flow Diagram"
                    )
                )
            )
        ]
    )
    
    # Mock methods to avoid side effects
    builder._add_text_content_to_placeholder = MagicMock()
    
    # Fill the slide
    builder._fill_chart_layout_slide(mock_pptx_slide, slide)
    
    # Verify text was added to content placeholder
    builder._add_text_content_to_placeholder.assert_called_once_with(
        mock_pptx_slide.shapes[1],
        "Text content to go with the diagram"
    )
    
    # Verify mermaid description was added to chart placeholder
    assert "[Mermaid diagram" in chart_placeholder.text_frame.text
    assert "Flow Diagram" in chart_placeholder.text_frame.text


def test_add_block_to_placeholder(builder):
    """Test adding a block to a placeholder."""
    # Create a mock placeholder
    placeholder = MagicMock()
    placeholder.text_frame = MagicMock()
    placeholder.text_frame.paragraphs = [MagicMock()]
    placeholder.text_frame.add_paragraph = MagicMock(return_value=MagicMock())
    
    # Create a block
    block = SlideBlock(
        id="block1",
        title="Block Title",
        content=SlideContent(
            content_type=ContentType.TEXT,
            text="Block content"
        )
    )
    
    # Mock methods to avoid side effects
    builder._add_text_content_to_placeholder = MagicMock()
    
    # Add block to placeholder
    builder._add_block_to_placeholder(placeholder, block)
    
    # Verify placeholder was cleared
    placeholder.text_frame.clear.assert_called_once()
    
    # Verify block title was added
    placeholder.text_frame.add_paragraph.assert_called_once()
    para = placeholder.text_frame.add_paragraph.return_value
    assert para.text == "Block Title"
    assert para.font.bold is True
    
    # Verify block content was added
    builder._add_text_content_to_placeholder.assert_called_once_with(placeholder, "Block content")


def test_add_block_to_placeholder_with_bullet_points(builder):
    """Test adding a block with bullet points to a placeholder."""
    # Create a mock placeholder
    placeholder = MagicMock()
    placeholder.text_frame = MagicMock()
    placeholder.text_frame.paragraphs = [MagicMock()]
    placeholder.text_frame.add_paragraph = MagicMock(return_value=MagicMock())
    
    # Create a block with bullet points
    block = SlideBlock(
        id="block1",
        title="Bullet Points",
        content=SlideContent(
            content_type=ContentType.BULLET_POINTS,
            bullet_points=["Point 1", "Point 2", "Point 3"],
            as_bullets=True
        )
    )
    
    # Mock methods to avoid side effects
    builder._add_bullet_points_to_placeholder = MagicMock()
    
    # Add block to placeholder
    builder._add_block_to_placeholder(placeholder, block)
    
    # Verify placeholder was cleared
    placeholder.text_frame.clear.assert_called_once()
    
    # Verify block title was added
    placeholder.text_frame.add_paragraph.assert_called_once()
    
    # Verify bullet points were added
    builder._add_bullet_points_to_placeholder.assert_called_once_with(
        placeholder,
        ["Point 1", "Point 2", "Point 3"],
        True
    )


def test_add_block_to_placeholder_unsupported_content(builder):
    """Test adding a block with unsupported content to a placeholder."""
    # Create a mock placeholder
    placeholder = MagicMock()
    placeholder.text_frame = MagicMock()
    placeholder.text_frame.paragraphs = [MagicMock()]
    placeholder.text_frame.add_paragraph = MagicMock(return_value=MagicMock())
    
    # Create a block with unsupported content
    block = SlideBlock(
        id="block1",
        title="Unsupported Content",
        content=SlideContent(
            content_type=ContentType.TABLE,
            table=TableData(
                headers=["Col1", "Col2"],
                rows=[["A", "B"], ["C", "D"]]
            )
        )
    )
    
    # Add block to placeholder
    builder._add_block_to_placeholder(placeholder, block)
    
    # Verify placeholder was cleared
    placeholder.text_frame.clear.assert_called_once()
    
    # Verify block title was added
    placeholder.text_frame.add_paragraph.assert_called()
    
    # Verify fallback text was added
    para = placeholder.text_frame.add_paragraph.return_value
    assert "table content not shown" in para.text.lower()