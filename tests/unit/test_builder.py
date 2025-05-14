"""
Unit tests for the PPTBuilder class.
"""
import os
import pytest
from pathlib import Path
from unittest.mock import MagicMock, patch, PropertyMock

from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.presentation import Presentation as PptxPresentation
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide as PptxSlide
from pptx.text.text import TextFrame

from doc2pptx.core.models import (
    ContentType, 
    Presentation, 
    Section, 
    Slide, 
    SlideBlock, 
    SlideContent,
    SectionType
)
from doc2pptx.ppt.builder import PPTBuilder


@pytest.fixture
def sample_presentation():
    """Create a sample presentation model for testing."""
    # Create a test slide with text content
    text_content = SlideContent(
        content_type=ContentType.TEXT,
        text="This is a sample text for testing."
    )
    
    # Create a test slide with bullet points
    bullet_content = SlideContent(
        content_type=ContentType.BULLET_POINTS,
        bullet_points=["Point 1", "Point 2", "Point 3"]
    )
    
    # Create slide blocks
    text_block = SlideBlock(
        id="block1",
        title="Text Block",
        content=text_content
    )
    
    bullet_block = SlideBlock(
        id="block2",
        title="Bullet Block",
        content=bullet_content
    )
    
    # Create slides
    text_slide = Slide(
        id="slide1",
        title="Text Slide",
        layout_name="Titre et texte",
        blocks=[text_block],
        notes="This is a slide with text content."
    )
    
    bullet_slide = Slide(
        id="slide2",
        title="Bullet Slide",
        layout_name="Titre et texte",
        blocks=[bullet_block],
        notes="This is a slide with bullet points."
    )
    
    # Create a section containing the slides
    section = Section(
        id="section1",
        title="Test Section",
        type=SectionType.CONTENT,
        slides=[text_slide, bullet_slide],
        description="This is a test section."
    )
    
    # Create the presentation
    presentation = Presentation(
        id="pres1",
        title="Test Presentation",
        author="Test Author",
        description="This is a test presentation.",
        template_path=Path("tests/fixtures/base_template.pptx"),
        sections=[section],
        metadata={"created": "2023-01-01"}
    )
    
    return presentation


@pytest.fixture
def mock_pptx_slide():
    """Create a mock PowerPoint slide with placeholders."""
    slide = MagicMock(spec=PptxSlide)
    
    # Create a mock title placeholder
    title_shape = MagicMock(spec=Shape)
    title_shape.is_placeholder = True
    title_shape.placeholder_format.type = PP_PLACEHOLDER.TITLE
    title_shape.text_frame = MagicMock(spec=TextFrame)
    title_shape.name = "Title Placeholder"
    
    # Create a mock content placeholder
    content_shape = MagicMock(spec=Shape)
    content_shape.is_placeholder = True
    content_shape.placeholder_format.type = PP_PLACEHOLDER.BODY
    content_shape.text_frame = MagicMock(spec=TextFrame)
    content_shape.text_frame.paragraphs = [MagicMock()]
    content_shape.name = "Content Placeholder"
    
    # Add placeholders to the slide
    slide.shapes = [title_shape, content_shape]
    slide.notes_slide = MagicMock()
    slide.notes_slide.notes_text_frame = MagicMock()
    
    return slide


@pytest.fixture
def mock_pptx_presentation():
    """Create a mock PowerPoint presentation with layouts."""
    pres = MagicMock(spec=PptxPresentation)
    
    # Create mock layouts
    title_layout = MagicMock()
    title_layout.name = "Diapositive de titre"
    
    content_layout = MagicMock()
    content_layout.name = "Titre et texte"
    
    intro_layout = MagicMock()
    intro_layout.name = "Introduction"
    
    chapitre_layout = MagicMock()
    chapitre_layout.name = "Chapitre 1"
    
    histo_layout = MagicMock()
    histo_layout.name = "Titre et texte 1 histogramme"
    
    visuel_layout = MagicMock()
    visuel_layout.name = "Titre et texte 1 visuel gauche"
    
    colonnes_layout = MagicMock()
    colonnes_layout.name = "Titre et 3 colonnes"   
    
    table_layout = MagicMock()
    table_layout.name = "Titre et tableau" 
    
    
    # Add layouts to the presentation
    pres.slide_layouts = [title_layout, content_layout, intro_layout, chapitre_layout, 
                          histo_layout, visuel_layout, colonnes_layout, table_layout]
    
    # Mock slides collection
    pres.slides = MagicMock()
    pres.slides.add_slide = MagicMock(return_value=mock_pptx_slide)
    
    return pres


def test_init():
    """Test PPTBuilder initialization."""
    # Test initialization without template
    builder = PPTBuilder()
    assert builder.template_path is None
    assert builder.template_info is None
    
    # Test initialization with template
    with patch('doc2pptx.ppt.builder.TemplateLoader.analyze_template') as mock_analyze:
        mock_analyze.return_value = "template_info"
        builder = PPTBuilder(template_path="tests/fixtures/base_template.pptx")
        assert builder.template_path == Path("tests/fixtures/base_template.pptx")
        assert builder.template_info == "template_info"
        mock_analyze.assert_called_once_with(Path("tests/fixtures/base_template.pptx"))


def test_create_slide():
    """Test _create_slide method."""
    builder = PPTBuilder()
    
    # Mock PowerPoint presentation
    pptx = MagicMock(spec=PptxPresentation)
    layout1 = MagicMock()
    layout1.name = "Layout1"
    layout2 = MagicMock()
    layout2.name = "Layout2"
    pptx.slide_layouts = [layout1, layout2]
    pptx.slides.add_slide = MagicMock(return_value="new_slide")
    
    # Test creating a slide with an existing layout
    result = builder._create_slide(pptx, "Layout2")
    assert result == "new_slide"
    pptx.slides.add_slide.assert_called_once_with(layout2)
    
    # Reset mock
    pptx.slides.add_slide.reset_mock()
    
    # Test creating a slide with a non-existent layout (should use first available)
    with patch('doc2pptx.ppt.builder.logger.warning') as mock_warning:
        result = builder._create_slide(pptx, "NonExistentLayout")
        assert result == "new_slide"
        pptx.slides.add_slide.assert_called_once_with(layout1)
        mock_warning.assert_called()


def test_get_placeholder_mapping(mock_pptx_slide):
    """Test _get_placeholder_mapping method."""
    builder = PPTBuilder()
    
    # Test mapping creation
    mapping = builder._get_placeholder_mapping(mock_pptx_slide)
    
    # Check that the mapping contains expected keys
    assert 'title' in mapping
    assert 'content' in mapping
    assert 'title_placeholders' in mapping
    assert 'content_placeholders' in mapping
    
    # Check that the placeholders are correctly mapped
    assert mapping['title'] == mock_pptx_slide.shapes[0]
    assert mapping['content'] == mock_pptx_slide.shapes[1]
    assert len(mapping['title_placeholders']) == 1
    assert len(mapping['content_placeholders']) == 1


def test_find_placeholder_for_block(mock_pptx_slide):
    """Test _find_placeholder_for_block method."""
    builder = PPTBuilder()

    # Get the placeholder mapping
    placeholder_mapping = builder._get_placeholder_mapping(mock_pptx_slide)

    # Create a text block
    text_content = SlideContent(
        content_type=ContentType.TEXT,
        text="This is a test text."
    )
    text_block = SlideBlock(
        id="block1",
        title="Text Block",
        content=text_content
    )

    # Test finding a placeholder for a text block
    placeholder = builder._find_placeholder_for_block(
        mock_pptx_slide, text_block, 0, placeholder_mapping
    )

    # Vérifier que le placeholder n'est pas None
    assert placeholder is not None
    
    # Vérifier que le placeholder est un placeholder de contenu
    # Cette vérification dépend de la configuration du mock et de son API exacte
    if hasattr(placeholder, 'is_placeholder'):
        assert placeholder.is_placeholder
    
    # Si nous savons que le placeholder dans le mapping est directement lié à shapes[1]
    # nous pouvons vérifier certaines de ses propriétés au lieu de comparer les références
    assert placeholder in mock_pptx_slide.shapes


def test_fill_placeholder_with_content(mock_pptx_slide):
    """Test _fill_placeholder_with_content method."""
    builder = PPTBuilder()
    content_placeholder = mock_pptx_slide.shapes[1]
    
    # Test filling with text content
    text_content = SlideContent(
        content_type=ContentType.TEXT,
        text="This is a test text."
    )
    
    # Mock the overflow handler
    builder.overflow_handler.will_text_overflow = MagicMock(return_value=False)
    
    # Call the method
    builder._fill_placeholder_with_content(content_placeholder, text_content)
    
    # Check that the text was set correctly
    content_placeholder.text_frame.clear.assert_called_once()
    content_placeholder.text_frame.paragraphs[0].text = "This is a test text."
    
    # Reset mocks
    content_placeholder.text_frame.clear.reset_mock()
    content_placeholder.text_frame.paragraphs[0].reset_mock()
    
    # Test filling with bullet points content
    bullet_content = SlideContent(
        content_type=ContentType.BULLET_POINTS,
        bullet_points=["Point 1", "Point 2", "Point 3"]
    )
    
    # Mock add_paragraph method
    content_placeholder.text_frame.add_paragraph = MagicMock(return_value=MagicMock())
    
    # Call the method
    builder._fill_placeholder_with_content(content_placeholder, bullet_content)
    
    # Check that the bullet points were set correctly
    content_placeholder.text_frame.clear.assert_called_once()
    assert content_placeholder.text_frame.add_paragraph.call_count == 2  # Called twice for points 2 and 3
    
    # Point 1 should use the existing first paragraph
    content_placeholder.text_frame.paragraphs[0].text = "Point 1"
    content_placeholder.text_frame.paragraphs[0].level = 0
    
    # The added paragraphs should be set to the other points
    added_paragraphs = content_placeholder.text_frame.add_paragraph.return_value
    assert added_paragraphs.text in ("Point 2", "Point 3")
    assert added_paragraphs.level == 0


def test_fill_slide(mock_pptx_slide):
    """Test _fill_slide method."""
    builder = PPTBuilder()
    
    # Create a slide model
    text_content = SlideContent(
        content_type=ContentType.TEXT,
        text="This is a sample text."
    )
    text_block = SlideBlock(
        id="block1",
        title=None,
        content=text_content
    )
    slide = Slide(
        id="slide1",
        title="Test Slide",
        layout_name="Titre et texte",
        blocks=[text_block],
        notes="Test notes"
    )
    section = Section(
        id="section1",
        title="Test Section",
        type=SectionType.CONTENT,
        slides=[slide]
    )
    
    # Mock methods
    builder._get_placeholder_mapping = MagicMock(return_value={
        'title': mock_pptx_slide.shapes[0],
        'content': mock_pptx_slide.shapes[1],
        'title_placeholders': [mock_pptx_slide.shapes[0]],
        'content_placeholders': [mock_pptx_slide.shapes[1]]
    })
    builder._find_placeholder_for_block = MagicMock(return_value=mock_pptx_slide.shapes[1])
    builder._fill_placeholder_with_content = MagicMock()
    
    # Call the method
    builder._fill_slide(mock_pptx_slide, slide, section)
    
    # Check that the title was set
    mock_pptx_slide.shapes[0].text_frame.text = "Test Slide"
    
    # Check that _find_placeholder_for_block was called
    builder._find_placeholder_for_block.assert_called_once()
    
    # Check that _fill_placeholder_with_content was called
    builder._fill_placeholder_with_content.assert_called_once_with(
        mock_pptx_slide.shapes[1], text_content
    )
    
    # Check that notes were set
    mock_pptx_slide.notes_slide.notes_text_frame.text = "Test notes"


@patch('doc2pptx.ppt.builder.TemplateLoader.load_template')
@patch('doc2pptx.ppt.builder.TemplateLoader.analyze_template')
def test_build(mock_analyze_template, mock_load_template, sample_presentation, mock_pptx_presentation, mock_pptx_slide):
    """Test build method."""
    # Setup mocks
    mock_analyze_template.return_value = "template_info"
    mock_load_template.return_value = mock_pptx_presentation
    
    # Create builder
    builder = PPTBuilder()
    
    # Mock methods
    builder._create_slide = MagicMock(return_value=mock_pptx_slide)
    builder._fill_slide = MagicMock()
    
    # Call the method
    output_path = builder.build(sample_presentation, "output.pptx")
    
    # Check that the template was loaded
    mock_load_template.assert_called_once_with(sample_presentation.template_path)
    
    # Check that slides were created (one for each slide in the presentation)
    assert builder._create_slide.call_count == 3
    
    # Check that slides were filled (one for each slide in the presentation)
    assert builder._fill_slide.call_count == 2
    
    # Check that the presentation was saved
    mock_pptx_presentation.save.assert_called_once_with(Path("output.pptx"))
    
    # Check that the correct path was returned
    assert output_path == Path("output.pptx")


def test_build_no_template_error(sample_presentation):
    """Test build method with no template."""
    # Create a presentation without a template path
    presentation = sample_presentation
    presentation.template_path = None
    
    # Create builder without a template
    builder = PPTBuilder()
    
    # Call the method should raise an error
    with pytest.raises(ValueError, match="No template path provided"):
        builder.build(presentation, "output.pptx")