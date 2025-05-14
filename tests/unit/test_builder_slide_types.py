"""
Unit tests for different slide types in PPTBuilder.
"""
import pytest
from unittest.mock import MagicMock, patch

from pptx.slide import Slide as PptxSlide
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from doc2pptx.core.models import (
    ContentType, 
    Slide, 
    SlideBlock, 
    SlideContent,
    TableData,
    ImageSource,
    ChartData,
    MermaidDiagram,
    SectionType,
    Section
)
from doc2pptx.ppt.builder_v3 import PPTBuilder


@pytest.fixture
def builder():
    """Create a PPTBuilder instance."""
    return PPTBuilder()


@pytest.fixture
def title_slide():
    """Create a title slide."""
    return Slide(
        id="title_slide",
        title="Presentation Title",
        layout_name="Diapositive de titre",
        blocks=[
            SlideBlock(
                id="subtitle_block",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Subtitle: Important Presentation"
                )
            )
        ]
    )


@pytest.fixture
def content_slide():
    """Create a standard content slide."""
    return Slide(
        id="content_slide",
        title="Content Slide",
        layout_name="Titre et texte",
        blocks=[
            SlideBlock(
                id="content_block",
                title="Block Title",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="This is the main content of the slide.\nIt has multiple paragraphs.\nAnd formats like **bold** and *italic*."
                )
            )
        ]
    )


@pytest.fixture
def bullet_points_slide():
    """Create a slide with bullet points."""
    return Slide(
        id="bullet_points_slide",
        title="Bullet Points Slide",
        layout_name="Titre et texte",
        blocks=[
            SlideBlock(
                id="bullet_points_block",
                title="Important Points",
                content=SlideContent(
                    content_type=ContentType.BULLET_POINTS,
                    bullet_points=[
                        "First important point",
                        "Second important point",
                        "Third important point with **bold** formatting"
                    ],
                    as_bullets=True
                )
            )
        ]
    )


@pytest.fixture
def multi_block_slide():
    """Create a slide with multiple content blocks."""
    return Slide(
        id="multi_block_slide",
        title="Multiple Blocks Slide",
        layout_name="Titre et 3 colonnes",
        blocks=[
            SlideBlock(
                id="block1",
                title="First Column",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Content for the first column"
                )
            ),
            SlideBlock(
                id="block2",
                title="Second Column",
                content=SlideContent(
                    content_type=ContentType.BULLET_POINTS,
                    bullet_points=["Point 1", "Point 2", "Point 3"],
                    as_bullets=True
                )
            ),
            SlideBlock(
                id="block3",
                title="Third Column",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Content for the third column"
                )
            )
        ]
    )


@pytest.fixture
def table_slide():
    """Create a slide with a table."""
    return Slide(
        id="table_slide",
        title="Table Slide",
        layout_name="Titre et tableau",
        blocks=[
            SlideBlock(
                id="table_block",
                content=SlideContent(
                    content_type=ContentType.TABLE,
                    table=TableData(
                        headers=["Name", "Age", "Occupation"],  # Retirer "style:accent1"
                        rows=[
                            ["John Doe", "32", "Software Engineer"],
                            ["Jane Smith", "28", "Data Scientist"],
                            ["Robert Johnson", "45", "Project Manager"]
                        ]
                    )
                )
            )
        ]
    )


@pytest.fixture
def image_slide():
    """Create a slide with image and text."""
    return Slide(
        id="image_slide",
        title="Image Slide",
        layout_name="Titre et texte 1 visuel gauche",
        blocks=[
            SlideBlock(
                id="image_block",
                content=SlideContent(
                    content_type=ContentType.IMAGE,
                    image=ImageSource(
                        query="scenic mountain landscape",
                        alt_text="Beautiful mountain landscape"
                    )
                )
            ),
            SlideBlock(
                id="text_block",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="This text describes the image of a beautiful mountain landscape.\nThe image is used to illustrate natural beauty."
                )
            )
        ]
    )


@pytest.fixture
def chart_slide():
    """Create a slide with chart and text."""
    return Slide(
        id="chart_slide",
        title="Chart Slide",
        layout_name="Titre et texte 1 histogramme",
        blocks=[
            SlideBlock(
                id="text_block",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="This text explains the chart data.\nIt shows sales performance over multiple quarters."
                )
            ),
            SlideBlock(
                id="chart_block",
                content=SlideContent(
                    content_type=ContentType.CHART,
                    chart=ChartData(
                        chart_type="bar",
                        categories=["Q1", "Q2", "Q3", "Q4"],
                        series=[
                            {"name": "2023", "data": [100, 150, 200, 180]},
                            {"name": "2024", "data": [120, 180, 210, 220]}
                        ],
                        title="Quarterly Sales Comparison"
                    )
                )
            )
        ]
    )


@pytest.fixture
def mermaid_slide():
    """Create a slide with mermaid diagram and text."""
    return Slide(
        id="mermaid_slide",
        title="Diagram Slide",
        layout_name="Titre et texte 1 histogramme",
        blocks=[
            SlideBlock(
                id="text_block",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="This text explains the flowchart.\nIt illustrates a simple process flow."
                )
            ),
            SlideBlock(
                id="mermaid_block",
                content=SlideContent(
                    content_type=ContentType.MERMAID,
                    mermaid=MermaidDiagram(
                        code="graph TD; A[Start] --> B[Process]; B --> C[End];",
                        caption="Simple Process Flow"
                    )
                )
            )
        ]
    )


@pytest.fixture
def section_header_slide():
    """Create a section header slide."""
    return Slide(
        id="section_header",
        title="Section Title",
        layout_name="Chapitre 1",
        blocks=[]
    )


@pytest.fixture
def mock_pptx_slide():
    """Create a mock PowerPoint slide with common placeholders."""
    mock_slide = MagicMock(spec=PptxSlide)
    
    # Create title placeholder
    title_placeholder = MagicMock()
    title_placeholder.is_placeholder = True
    title_placeholder.placeholder_format = MagicMock()
    title_placeholder.placeholder_format.type = PP_PLACEHOLDER.TITLE
    title_placeholder.placeholder_format.idx = 0
    title_placeholder.text_frame = MagicMock()
    
    # Create content placeholder
    content_placeholder = MagicMock()
    content_placeholder.is_placeholder = True
    content_placeholder.placeholder_format = MagicMock()
    content_placeholder.placeholder_format.type = PP_PLACEHOLDER.BODY
    content_placeholder.placeholder_format.idx = 1
    content_placeholder.text_frame = MagicMock()
    content_placeholder.text_frame.paragraphs = [MagicMock()]
    
    # Setup shapes collection
    mock_slide.shapes = MagicMock()
    placeholders = [title_placeholder, content_placeholder]
    mock_slide.shapes.__iter__.return_value = iter(placeholders)
    mock_slide.shapes.__getitem__.side_effect = lambda idx: placeholders[idx] if idx < len(placeholders) else None
    mock_slide.shapes.__len__.return_value = len(placeholders)
    
    # Add shapes.add_table method
    mock_table_shape = MagicMock()
    mock_table = MagicMock()
    mock_table_shape.table = mock_table
    mock_slide.shapes.add_table = MagicMock(return_value=mock_table_shape)
    
    # Setup notes slide
    mock_slide.notes_slide = MagicMock()
    mock_slide.notes_slide.notes_text_frame = MagicMock()
    
    return mock_slide


def test_process_different_slide_types(builder, mock_pptx_slide, title_slide, content_slide, 
                                      bullet_points_slide, multi_block_slide, table_slide, 
                                      image_slide, chart_slide, mermaid_slide, section_header_slide):
    """Test processing different types of slides."""
    # Create a section for context
    section = Section(
        id="section1",
        title="Test Section",
        type=SectionType.CONTENT,
        slides=[]
    )
    
    # Test each slide type
    slides = [
        title_slide,
        content_slide,
        bullet_points_slide,
        multi_block_slide,
        table_slide, 
        image_slide,
        chart_slide,
        mermaid_slide,
        section_header_slide
    ]
    
    # Mock specific processing methods
    builder._fill_title_slide = MagicMock()
    builder._fill_content_slide = MagicMock()
    builder._fill_column_layout_slide = MagicMock()
    builder._fill_table_slide = MagicMock()
    builder._fill_image_layout_slide = MagicMock()
    builder._fill_chart_layout_slide = MagicMock()
    builder._fill_slide_title = MagicMock()
    
    # Process each slide
    for slide in slides:
        builder._fill_slide(mock_pptx_slide, slide, section)
        
        # Verify title was set
        builder._fill_slide_title.assert_called_with(mock_pptx_slide, slide.title)
        builder._fill_slide_title.reset_mock()
        
        # Verify appropriate method was called based on layout
        if slide.layout_name == "Diapositive de titre":
            builder._fill_title_slide.assert_called_once_with(mock_pptx_slide, slide)
            builder._fill_title_slide.reset_mock()
        elif slide.layout_name == "Titre et texte":
            builder._fill_content_slide.assert_called_once_with(mock_pptx_slide, slide)
            builder._fill_content_slide.reset_mock()
        elif slide.layout_name == "Titre et 3 colonnes":
            builder._fill_column_layout_slide.assert_called_once_with(mock_pptx_slide, slide)
            builder._fill_column_layout_slide.reset_mock()
        elif slide.layout_name == "Titre et tableau":
            builder._fill_table_slide.assert_called_once_with(mock_pptx_slide, slide)
            builder._fill_table_slide.reset_mock()
        elif slide.layout_name == "Titre et texte 1 visuel gauche":
            builder._fill_image_layout_slide.assert_called_once_with(mock_pptx_slide, slide)
            builder._fill_image_layout_slide.reset_mock()
        elif slide.layout_name == "Titre et texte 1 histogramme":
            builder._fill_chart_layout_slide.assert_called_once_with(mock_pptx_slide, slide)
            builder._fill_chart_layout_slide.reset_mock()


def test_fill_slide_with_notes(builder, mock_pptx_slide, content_slide):
    """Test filling a slide with speaker notes."""
    # Add notes to the slide
    content_slide.notes = "These are important speaker notes for the presenter."
    
    # Create a section for context
    section = Section(
        id="section1",
        title="Test Section",
        type=SectionType.CONTENT,
        slides=[content_slide]
    )
    
    # Mock the fill_content_slide method
    builder._fill_content_slide = MagicMock()
    builder._fill_slide_title = MagicMock()
    
    # Fill the slide
    builder._fill_slide(mock_pptx_slide, content_slide, section)
    
    # Verify notes were added
    assert mock_pptx_slide.notes_slide.notes_text_frame.text == "These are important speaker notes for the presenter."


def test_process_slide_with_unknown_layout(builder, mock_pptx_slide, content_slide):
    """Test processing a slide with an unknown layout."""
    # Change the layout name to something unknown
    content_slide.layout_name = "Unknown Layout"
    
    # Create a section for context
    section = Section(
        id="section1",
        title="Test Section",
        type=SectionType.CONTENT,
        slides=[content_slide]
    )
    
    # Mock methods to avoid side effects
    builder._fill_content_slide = MagicMock()
    builder._fill_slide_title = MagicMock()
    
    # Process the slide
    builder._fill_slide(mock_pptx_slide, content_slide, section)
    
    # Verify it fell back to the content slide method
    builder._fill_content_slide.assert_called_once_with(mock_pptx_slide, content_slide)


def test_special_slide_types(builder, mock_pptx_slide):
    """Test handling of special slide types."""
    # Create slides with special layouts
    intro_slide = Slide(
        id="intro_slide",
        title="Introduction",
        layout_name="Introduction",
        blocks=[
            SlideBlock(
                id="intro_block",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Introduction content here."
                )
            )
        ]
    )
    
    chapter_slide = Slide(
        id="chapter_slide",
        title="Chapter Title",
        layout_name="Chapitre 1",
        blocks=[]
    )
    
    # Create a section for context
    section = Section(
        id="section1",
        title="Test Section",
        type=SectionType.CONTENT,
        slides=[intro_slide, chapter_slide]
    )
    
    # Mock methods to avoid side effects
    builder._fill_content_slide = MagicMock()
    builder._fill_slide_title = MagicMock()
    
    # Process the intro slide - should use content slide handler
    builder._fill_slide(mock_pptx_slide, intro_slide, section)
    builder._fill_content_slide.assert_called_once_with(mock_pptx_slide, intro_slide)
    builder._fill_content_slide.reset_mock()
    
    # Reset le mock de _fill_slide_title
    builder._fill_slide_title.reset_mock()
    
    # Process the chapter slide - should only fill title
    builder._fill_slide(mock_pptx_slide, chapter_slide, section)
    builder._fill_slide_title.assert_called_once_with(mock_pptx_slide, chapter_slide.title)


def test_custom_formatting_in_text(builder):
    """Test that custom text formatting is properly applied."""
    # Create a paragraph mock
    paragraph = MagicMock()
    run = MagicMock()
    paragraph.add_run.return_value = run
    
    # Test with various formatting
    formatted_text = "Normal **bold** *italic* {color:red}colored{/color} {size:16pt}sized{/size} text"
    
    # Mock _apply_pattern to avoid complex formatting logic
    original_apply_pattern = builder._apply_pattern
    
    def mock_apply_pattern(segments, pattern, formatter):
        if pattern == builder.BOLD_PATTERN:
            # Simulate finding bold text
            return [
                {'text': 'Normal '},
                {'text': 'bold', 'bold': True},
                {'text': ' *italic* {color:red}colored{/color} {size:16pt}sized{/size} text'}
            ]
        elif pattern == builder.ITALIC_PATTERN:
            # Simulate finding italic text
            return [
                {'text': 'Normal **bold** '},
                {'text': 'italic', 'italic': True},
                {'text': ' {color:red}colored{/color} {size:16pt}sized{/size} text'}
            ]
        elif pattern == builder.COLOR_PATTERN:
            # Simulate finding colored text
            return [
                {'text': 'Normal **bold** *italic* '},
                {'text': 'colored', 'color': 'red'},
                {'text': ' {size:16pt}sized{/size} text'}
            ]
        elif pattern == builder.FONT_SIZE_PATTERN:
            # Simulate finding sized text
            return [
                {'text': 'Normal **bold** *italic* {color:red}colored{/color} '},
                {'text': 'sized', 'size': '16'},
                {'text': ' text'}
            ]
        # Default pass-through for other patterns
        return segments
    
    # Replace the _apply_pattern method
    builder._apply_pattern = mock_apply_pattern
    
    # Also mock _hex_to_rgb
    builder._hex_to_rgb = MagicMock()
    builder._closest_highlight_color = MagicMock()
    
    try:
        # Call the method
        builder._add_formatted_text_to_paragraph(paragraph, formatted_text)
        
        # Verify paragraph.add_run was called
        assert paragraph.add_run.called
        
        # Test is successful if no exceptions were raised
    finally:
        # Restore original method
        builder._apply_pattern = original_apply_pattern


def test_complex_slide_layout_validation(builder):
    """Test complex validation of slide layouts for different content."""
    # Create a slide with mixed content that doesn't match the layout
    complex_slide = Slide(
        id="complex_slide",
        title="Complex Slide",
        layout_name="Titre et texte",  # Text layout
        blocks=[
            SlideBlock(
                id="block1",
                content=SlideContent(
                    content_type=ContentType.TABLE,  # Table content in text layout
                    table=TableData(
                        headers=["Col1", "Col2"],
                        rows=[["A", "B"], ["C", "D"]]
                    )
                )
            ),
            SlideBlock(
                id="block2",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Additional text content"
                )
            )
        ]
    )
    
    # Validate the layout
    with patch('doc2pptx.ppt.builder_v3.logger.warning') as mock_warning:
        new_layout = builder._validate_layout_for_content(complex_slide)
        
        # Should recommend a table layout due to table content
        assert new_layout == "Titre et tableau"
        assert mock_warning.called


def test_add_block_to_placeholder_with_each_content_type(builder):
    """Test adding a block with each content type to a placeholder."""
    # Create a mock placeholder
    placeholder = MagicMock()
    placeholder.text_frame = MagicMock()
    placeholder.text_frame.paragraphs = [MagicMock()]
    placeholder.text_frame.add_paragraph = MagicMock(return_value=MagicMock())
    
    # Mock methods to avoid side effects
    builder._add_text_content_to_placeholder = MagicMock()
    builder._add_bullet_points_to_placeholder = MagicMock()
    
    # Test each content type
    content_types = [
        (ContentType.TEXT, SlideContent(content_type=ContentType.TEXT, text="Text content")),
        (ContentType.BULLET_POINTS, SlideContent(
            content_type=ContentType.BULLET_POINTS, 
            bullet_points=["Point 1", "Point 2"],
            as_bullets=True
        )),
        (ContentType.TABLE, SlideContent(
            content_type=ContentType.TABLE,
            table=TableData(headers=["Col1", "Col2"], rows=[["A", "B"]])
        )),
        (ContentType.IMAGE, SlideContent(
            content_type=ContentType.IMAGE,
            image=ImageSource(query="test image", alt_text="Test image")
        )),
        (ContentType.CHART, SlideContent(
            content_type=ContentType.CHART,
            chart=ChartData(
                chart_type="bar",
                categories=["A", "B"],
                series=[{"name": "Series1", "data": [1, 2]}]
            )
        )),
        (ContentType.MERMAID, SlideContent(
            content_type=ContentType.MERMAID,
            mermaid=MermaidDiagram(code="graph TD; A-->B;", caption="Test diagram")
        ))
    ]
    
    for content_type, content in content_types:
        # Create a block with the content type
        block = SlideBlock(
            id=f"block_{content_type.value}",
            title=f"{content_type.value.title()} Block",
            content=content
        )
        
        # Reset mocks
        placeholder.text_frame.clear.reset_mock()
        placeholder.text_frame.add_paragraph.reset_mock()
        builder._add_text_content_to_placeholder.reset_mock()
        builder._add_bullet_points_to_placeholder.reset_mock()
        
        # Add block to placeholder
        builder._add_block_to_placeholder(placeholder, block)
        
        # Verify placeholder was cleared
        placeholder.text_frame.clear.assert_called_once()
        
        # Verify block title was added
        placeholder.text_frame.add_paragraph.assert_called()
        
        # Verify appropriate method was called based on content type
        if content_type == ContentType.TEXT:
            builder._add_text_content_to_placeholder.assert_called_once_with(placeholder, "Text content")
        elif content_type == ContentType.BULLET_POINTS:
            builder._add_bullet_points_to_placeholder.assert_called_once_with(placeholder, ["Point 1", "Point 2"], True)
        else:
            # For other content types, a fallback text should be added
            assert placeholder.text_frame.add_paragraph.call_count > 1
            last_paragraph = placeholder.text_frame.add_paragraph.return_value
            assert f"{content_type.value} content not shown" in last_paragraph.text.lower()