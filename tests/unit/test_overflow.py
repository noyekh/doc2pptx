"""
Unit tests for the OverflowHandler class.
"""
import pytest
from unittest.mock import MagicMock, patch, PropertyMock

from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.text.text import TextFrame
from pptx.util import Pt, Inches

from doc2pptx.ppt.overflow import OverflowHandler


@pytest.fixture
def mock_shape():
    """Create a mock shape with text frame for testing."""
    shape = MagicMock(spec=BaseShape)
    shape.text_frame = MagicMock(spec=TextFrame)
    
    # Mock width and height properties
    type(shape).width = PropertyMock(return_value=Inches(10))
    type(shape).height = PropertyMock(return_value=Inches(5))
    
    return shape


def test_init():
    """Test OverflowHandler initialization."""
    handler = OverflowHandler()
    assert handler.max_chars_per_line == 90
    assert handler.max_lines_per_text_box == 15
    assert handler.avg_char_width_pt == 7.0
    assert handler.avg_line_height_pt == 18.0
    
    # Test with custom parameters
    handler = OverflowHandler(
        max_chars_per_line=80,
        max_lines_per_text_box=10,
        avg_char_width_pt=8.0,
        avg_line_height_pt=20.0
    )
    assert handler.max_chars_per_line == 80
    assert handler.max_lines_per_text_box == 10
    assert handler.avg_char_width_pt == 8.0
    assert handler.avg_line_height_pt == 20.0


def test_estimate_lines_needed():
    """Test _estimate_lines_needed method."""
    handler = OverflowHandler()
    
    # Test with empty text
    assert handler._estimate_lines_needed("", 80) == 0
    
    # Test with single line
    assert handler._estimate_lines_needed("This is a short line.", 80) == 1
    
    # Test with multiple lines due to line breaks
    text = "Line 1\nLine 2\nLine 3"
    assert handler._estimate_lines_needed(text, 80) == 3
    
    # Test with line wrapping
    long_line = "This is a very long line that should wrap to multiple lines because it exceeds the maximum characters per line."
    # With 30 chars per line, this should wrap to multiple lines
    assert handler._estimate_lines_needed(long_line, 30) > 1
    
    # Test with empty paragraphs
    text_with_empty = "Line 1\n\nLine 3"
    assert handler._estimate_lines_needed(text_with_empty, 80) == 3


def test_will_text_overflow(mock_shape):
    """Test will_text_overflow method."""
    handler = OverflowHandler()
    
    # Test with short text (should not overflow)
    short_text = "This is a short text that should fit in the placeholder."
    assert not handler.will_text_overflow(mock_shape, short_text)
    
    # Test with long text (should overflow)
    # Generate a long text that will definitely overflow
    long_text = "This is a test line. " * 100
    assert handler.will_text_overflow(mock_shape, long_text)
    
    # Test with a shape that doesn't have a text_frame
    shape_without_text_frame = MagicMock(spec=BaseShape)
    # Remove text_frame attribute
    del shape_without_text_frame.text_frame
    
    with patch('doc2pptx.ppt.overflow.logger.warning') as mock_warning:
        result = handler.will_text_overflow(shape_without_text_frame, "Some text")
        assert not result  # Should return False
        mock_warning.assert_called()  # Should log a warning


def test_will_text_overflow_no_dimensions():
    """Test will_text_overflow when shape dimensions cannot be determined."""
    handler = OverflowHandler()
    
    # Shape minimaliste dont width / height lèvent AttributeError
    class BrokenShape:
        @property
        def text_frame(self):
            return MagicMock(spec=TextFrame)

        @property
        def width(self):
            raise AttributeError("no width")

        @property
        def height(self):
            raise AttributeError("no height")

    shape = BrokenShape()
    
    # Mock _estimate_overflow to test the fallback
    handler._estimate_overflow = MagicMock(return_value=True)
    
    with patch('doc2pptx.ppt.overflow.logger.warning') as mock_warning:
        result = handler.will_text_overflow(shape, "Some text")
        assert result  # Should return result from _estimate_overflow
        mock_warning.assert_called()  # Should log a warning
        handler._estimate_overflow.assert_called_once_with("Some text")


def test_estimate_overflow():
    """Test _estimate_overflow method."""
    handler = OverflowHandler()
    
    # Test with short text (should not overflow)
    short_text = "This is a short text that should fit in the placeholder."
    assert not handler._estimate_overflow(short_text)
    
    # Test with text that exceeds character limit
    # Create a text that's longer than max_chars_per_line * max_lines_per_text_box
    long_text = "a" * (handler.max_chars_per_line * handler.max_lines_per_text_box + 1)
    assert handler._estimate_overflow(long_text)
    
    # Test with text that has too many lines
    many_lines = "\n".join(["Line " + str(i) for i in range(handler.max_lines_per_text_box + 5)])
    assert handler._estimate_overflow(many_lines)


def test_split_text_for_overflow():
    """Test split_text_for_overflow method."""
    handler = OverflowHandler()
    
    # Test with empty text
    empty_chunks = handler.split_text_for_overflow("")
    assert len(empty_chunks) == 1
    assert empty_chunks[0] == ""
    
    # Test with short text
    short_text = "This is a short text."
    short_chunks = handler.split_text_for_overflow(short_text)
    assert len(short_chunks) == 1
    assert short_chunks[0] == short_text
    
    # Test with long text that needs splitting
    # Create paragraphs that will exceed max_chars_per_slide
    paragraphs = ["Paragraph " + str(i) + ": " + ("content " * 50) for i in range(10)]
    long_text = "\n".join(paragraphs)
    
    # Split with custom max_chars_per_slide
    max_chars = 200
    chunks = handler.split_text_for_overflow(long_text, max_chars_per_slide=max_chars)
    
    # Verify that each chunk is no longer than max_chars
    for chunk in chunks:
        assert len(chunk) <= max_chars


def test_split_bullet_points_for_overflow():
    """Test split_bullet_points_for_overflow method."""
    handler = OverflowHandler()
    
    # Test with empty list
    empty_chunks = handler.split_bullet_points_for_overflow([])
    assert len(empty_chunks) == 1
    assert empty_chunks[0] == []
    
    # Test with short list
    bullet_points = ["Point 1", "Point 2", "Point 3"]
    short_chunks = handler.split_bullet_points_for_overflow(bullet_points)
    assert len(short_chunks) == 1
    assert short_chunks[0] == bullet_points
    
    # Test with long list that needs splitting
    long_bullet_points = ["Point " + str(i) for i in range(25)]
    
    # Split with default max_points_per_slide (10)
    chunks = handler.split_bullet_points_for_overflow(long_bullet_points)
    
    # Should have 3 chunks (25 points with 10 per chunk = 3 chunks)
    assert len(chunks) == 3
    
    # First two chunks should have 10 points, last chunk should have 5 points
    assert len(chunks[0]) == 10
    assert len(chunks[1]) == 10
    assert len(chunks[2]) == 5
    
    # Verify content of first and last chunks
    assert chunks[0][0] == "Point 0"
    assert chunks[0][9] == "Point 9"
    assert chunks[2][0] == "Point 20"
    assert chunks[2][4] == "Point 24"
    
    # Test with custom max_points_per_slide
    custom_chunks = handler.split_bullet_points_for_overflow(long_bullet_points, max_points_per_slide=5)
    
    # Should have 5 chunks (25 points with 5 per chunk = 5 chunks)
    assert len(custom_chunks) == 5
    assert all(len(chunk) == 5 for chunk in custom_chunks)


def test_truncate_text_for_placeholder(mock_shape):
    """Test truncate_text_for_placeholder method."""
    handler = OverflowHandler()
    
    # Test with short text (should not be truncated)
    short_text = "This is a short text that should fit in the placeholder."
    truncated = handler.truncate_text_for_placeholder(mock_shape, short_text)
    assert truncated == short_text
    
    # Test with long text (should be truncated)
    # Generate a long text that will definitely be truncated
    long_text = "This is a test line. " * 100
    truncated = handler.truncate_text_for_placeholder(mock_shape, long_text)
    
    # Truncated text should be shorter than original
    assert len(truncated) < len(long_text)
    
    # Truncated text should end with ellipsis
    assert truncated.endswith(" ...")
    
    # Test without ellipsis
    truncated_no_ellipsis = handler.truncate_text_for_placeholder(
        mock_shape, long_text, add_ellipsis=False
    )
    assert not truncated_no_ellipsis.endswith(" ...")
    
    # Test with a shape that doesn't have a text_frame
    shape_without_text_frame = MagicMock(spec=BaseShape)
    # Remove text_frame attribute
    del shape_without_text_frame.text_frame
    
    with patch('doc2pptx.ppt.overflow.logger.warning') as mock_warning:
        result = handler.truncate_text_for_placeholder(shape_without_text_frame, "Some text")
        assert result == "Some text"  # Should return original text
        mock_warning.assert_called()  # Should log a warning


def test_truncate_text_for_placeholder_no_dimensions():
    """Test truncate_text_for_placeholder when shape dimensions cannot be determined."""
    handler = OverflowHandler()
    
    # Create a shape that will raise AttributeError when accessing width/height
    shape = MagicMock(spec=BaseShape)
    shape.text_frame = MagicMock(spec=TextFrame)
    
    # Make width/height properties raise AttributeError
    type(shape).width = PropertyMock(side_effect=AttributeError)
    type(shape).height = PropertyMock(side_effect=AttributeError)
    
    # Mock _truncate_text_by_chars to test the fallback
    handler._truncate_text_by_chars = MagicMock(return_value="Truncated text")
    
    with patch('doc2pptx.ppt.overflow.logger.warning') as mock_warning:
        result = handler.truncate_text_for_placeholder(shape, "Some text")
        assert result == "Truncated text"  # Should return result from _truncate_text_by_chars
        mock_warning.assert_called()  # Should log a warning
        handler._truncate_text_by_chars.assert_called_once_with("Some text", True)


def test_truncate_text_by_chars():
    """Test _truncate_text_by_chars method."""
    handler = OverflowHandler()
    
    # Test with short text (should not be truncated)
    short_text = "This is a short text."
    truncated = handler._truncate_text_by_chars(short_text)
    assert truncated == short_text
    
    # Test with long text (should be truncated)
    max_chars = handler.max_chars_per_line * handler.max_lines_per_text_box
    long_text = "word " * (max_chars // 5)  # Each "word " is 5 chars
    truncated = handler._truncate_text_by_chars(long_text)
    
    # Should be truncated and end with ellipsis
    assert len(truncated) < len(long_text)
    assert truncated.endswith(" ...")
    
    # Test word boundary truncation
    # Create a text with a few words that will exceed the max_chars
    boundary_text = "a " * (max_chars - 10) + "supercalifragilisticexpialidocious"
    truncated = handler._truncate_text_by_chars(boundary_text)
    
    # Should be truncated at the last space before the long word
    assert len(truncated) < len(boundary_text)
    assert not "supercal" in truncated  # The long word should not be included
    assert truncated.endswith(" ...")
    
    # Test with no ellipsis
    truncated_no_ellipsis = handler._truncate_text_by_chars(long_text, add_ellipsis=False)
    assert len(truncated_no_ellipsis) < len(long_text)
    assert not truncated_no_ellipsis.endswith(" ...")


def test_integration():
    """Test the integration of multiple OverflowHandler methods."""
    handler = OverflowHandler()
    
    # Create a large collection of bullet points
    bullet_points = [f"Bullet point {i}: This is a detailed bullet point with more information than necessary." for i in range(20)]
    
    # Split the bullet points into chunks
    chunks = handler.split_bullet_points_for_overflow(bullet_points, max_points_per_slide=8)
    
    # Verify that we get the right number of chunks
    assert len(chunks) == 3  # 20 points / 8 per slide = 3 slides (8 + 8 + 4)
    
    # Verify each chunk has correct number of points
    assert len(chunks[0]) == 8
    assert len(chunks[1]) == 8
    assert len(chunks[2]) == 4
    
    # Create a mock shape
    shape = MagicMock(spec=BaseShape)
    shape.text_frame = MagicMock(spec=TextFrame)
    type(shape).width = PropertyMock(return_value=Inches(8))
    type(shape).height = PropertyMock(return_value=Inches(3))  # Smaller height to force overflow
    
    # Test if a chunk of bullet points would overflow
    # Convert bullet points to text for overflow check
    bullet_text = "\n".join(chunks[0])
    
    # Check for overflow (with our small shape, it should overflow)
    assert handler.will_text_overflow(shape, bullet_text)
    
    # Truncate the text to fit the shape
    truncated = handler.truncate_text_for_placeholder(shape, bullet_text)
    
    # Verify the truncated text is shorter
    assert len(truncated) < len(bullet_text)
    
    # Verify ellipsis
    assert truncated.endswith(" ...")


def test_real_world_scenarios():
    """Test OverflowHandler with realistic content scenarios."""
    handler = OverflowHandler()
    
    # Scenario 1: Presentation slide with a paragraph of text
    paragraph = (
        "This is a typical slide with a paragraph of text that explains a concept. "
        "The text should be clear and concise, but sometimes we need to include "
        "more details than can comfortably fit on a slide. In those cases, "
        "we need to detect overflow and handle it appropriately to ensure "
        "the presentation looks professional."
    )
    
    # Create a mock shape with typical slide dimensions
    text_shape = MagicMock(spec=BaseShape)
    text_shape.text_frame = MagicMock(spec=TextFrame)
    type(text_shape).width = PropertyMock(return_value=Inches(9))  # Standard slide width
    type(text_shape).height = PropertyMock(return_value=Inches(4))  # Content area height
    
    # This paragraph should fit in a standard slide
    assert not handler.will_text_overflow(text_shape, paragraph)
    
    # Scenario 2: Slide with too many bullet points
    bullet_points = [
        "Introduction to the topic",
        "Background information and context",
        "Key point 1 with supporting details",
        "Key point 2 with examples and clarification",
        "Key point 3 with technical specifications",
        "Comparison with alternative approaches",
        "Benefits and advantages of our solution",
        "Potential challenges and limitations",
        "Implementation timeline and requirements",
        "Budget considerations and resource allocation",
        "Expected outcomes and success metrics",
        "Next steps and action items"
    ]
    
    # Convert to text
    bullet_text = "\n".join(bullet_points)
    
    # Create a mock shape for bullet points
    bullet_shape = MagicMock(spec=BaseShape)
    bullet_shape.text_frame = MagicMock(spec=TextFrame)
    type(bullet_shape).width = PropertyMock(return_value=Inches(9))
    type(bullet_shape).height = PropertyMock(return_value=Inches(5))
    
    # Check if bullets would overflow typical slide
    overflow_result = handler.will_text_overflow(bullet_shape, bullet_text)
    
    # Since this is somewhat dependent on the exact estimations, we don't assert true/false
    # Instead, we just run the code to ensure it doesn't crash
    
    # If overflow detected, split the bullets
    if overflow_result:
        chunks = handler.split_bullet_points_for_overflow(bullet_points, max_points_per_slide=6)
        assert len(chunks) == 2  # 12 points / 6 per slide = 2 slides