"""
Unit tests for the improved table handling and layout capabilities in PPTBuilder.
"""
import os
import pytest
from pathlib import Path
from unittest.mock import MagicMock, patch, PropertyMock
from pptx.util import Emu

from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
#from pptx.presentation import Presentation as PptxPresentation
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide as PptxSlide
from pptx.text.text import TextFrame, _Paragraph, _Run
from pptx.table import Table, _Cell, _Row, _Column
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

from doc2pptx.core.models import (
    ContentType, 
    Presentation, 
    Section, 
    Slide, 
    SlideBlock, 
    SlideContent,
    SectionType,
    TableData
)
from doc2pptx.ppt.builder_v3 import PPTBuilder, LAYOUT_CAPABILITIES, LAYOUT_PLACEHOLDER_MAP
from doc2pptx.ingest.json_loader import load_presentation


# ===== FIXTURES =====

@pytest.fixture
def sample_table_data():
    """Create sample table data for testing."""
    return TableData(
        headers=["Name", "Age", "Occupation"],
        rows=[
            ["John Doe", "32", "Software Engineer"],
            ["Jane Smith", "28", "Data Scientist"],
            ["Robert Johnson", "45", "Project Manager"]
        ]
    )


@pytest.fixture
def styled_table_data():
    """Create sample table data with style for testing."""
    return TableData(
        headers=["Product", "Price", "Stock", "style:accent1"],
        rows=[
            ["Laptop", "$999", "15"],
            ["Smartphone", "$499", "42"],
            ["Tablet", "$299", "23"]
        ]
    )


@pytest.fixture
def table_slide():
    """Create a slide with a table."""
    table_content = SlideContent(
        content_type=ContentType.TABLE,
        table=TableData(
            headers=["Name", "Age", "Occupation"],
            rows=[
                ["John Doe", "32", "Software Engineer"],
                ["Jane Smith", "28", "Data Scientist"],
                ["Robert Johnson", "45", "Project Manager"]
            ]
        )
    )
    
    return Slide(
        id="slide1",
        title="Simple Table",
        layout_name="Titre et tableau",
        blocks=[
            SlideBlock(
                id="block1",
                title=None,
                content=table_content
            )
        ]
    )


@pytest.fixture
def wrong_layout_table_slide():
    """Create a slide with a table but wrong layout."""
    table_content = SlideContent(
        content_type=ContentType.TABLE,
        table=TableData(
            headers=["Name", "Age", "Occupation"],
            rows=[
                ["John Doe", "32", "Software Engineer"],
                ["Jane Smith", "28", "Data Scientist"],
                ["Robert Johnson", "45", "Project Manager"]
            ]
        )
    )
    
    return Slide(
        id="slide1",
        title="Table in Wrong Layout",
        layout_name="Titre et texte",  # Wrong layout for a table
        blocks=[
            SlideBlock(
                id="block1",
                title=None,
                content=table_content
            )
        ]
    )


@pytest.fixture
def column_layout_slide():
    """Create a slide with multiple columns."""
    return Slide(
        id="slide_columns",
        title="Multiple Columns",
        layout_name="Titre et 3 colonnes",
        blocks=[
            SlideBlock(
                id="block1",
                title="Column 1",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Content for column 1."
                )
            ),
            SlideBlock(
                id="block2",
                title="Column 2",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Content for column 2."
                )
            ),
            SlideBlock(
                id="block3",
                title="Column 3",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Content for column 3."
                )
            )
        ]
    )


@pytest.fixture
def image_layout_slide():
    """Create a slide with image and text."""
    return Slide(
        id="slide_image",
        title="Image and Text",
        layout_name="Titre et texte 1 visuel gauche",
        blocks=[
            SlideBlock(
                id="image_block",
                content=SlideContent(
                    content_type=ContentType.IMAGE,
                    text="Description of the image",
                    image={
                        "query": "test image",
                        "alt_text": "Test image"
                    }
                )
            ),
            SlideBlock(
                id="text_block",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="This is the text that goes with the image."
                )
            )
        ]
    )


@pytest.fixture
def chart_layout_slide():
    """Create a slide with chart and text."""
    return Slide(
        id="slide_chart",
        title="Chart and Text",
        layout_name="Titre et texte 1 histogramme",
        blocks=[
            SlideBlock(
                id="chart_block",
                content=SlideContent(
                    content_type=ContentType.MERMAID,
                    mermaid={
                        "code": "graph TD; A-->B; B-->C;",
                        "caption": "Test diagram"
                    }
                )
            ),
            SlideBlock(
                id="text_block",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="This is the text that goes with the chart."
                )
            )
        ]
    )


@pytest.fixture
def mock_pptx_slide():
    """Create a mock PowerPoint slide."""
    slide = MagicMock(spec=PptxSlide)
    
    # Create different placeholder types
    title_placeholder = MagicMock()
    title_placeholder.is_placeholder = True
    title_placeholder.placeholder_format = MagicMock()
    title_placeholder.placeholder_format.type = PP_PLACEHOLDER.TITLE
    title_placeholder.placeholder_format.idx = 0
    title_placeholder.text_frame = MagicMock()
    
    content_placeholder = MagicMock()
    content_placeholder.is_placeholder = True
    content_placeholder.placeholder_format = MagicMock()
    content_placeholder.placeholder_format.type = PP_PLACEHOLDER.BODY
    content_placeholder.placeholder_format.idx = 1
    content_placeholder.text_frame = MagicMock()
    content_placeholder.text_frame.paragraphs = [MagicMock()]
    
    picture_placeholder = MagicMock()
    picture_placeholder.is_placeholder = True
    picture_placeholder.placeholder_format = MagicMock()
    picture_placeholder.placeholder_format.type = PP_PLACEHOLDER.PICTURE
    picture_placeholder.placeholder_format.idx = 2
    picture_placeholder.text_frame = MagicMock()
    
    chart_placeholder = MagicMock()
    chart_placeholder.is_placeholder = True
    chart_placeholder.placeholder_format = MagicMock()
    chart_placeholder.placeholder_format.type = PP_PLACEHOLDER.CHART
    chart_placeholder.placeholder_format.idx = 2
    chart_placeholder.text_frame = MagicMock()
    
    # Configurer shapes correctement avec un mock au lieu d'une liste
    slide.shapes = MagicMock()
    
    # Configurer le comportement d'itération et d'indexation
    placeholders = [title_placeholder, content_placeholder, picture_placeholder, chart_placeholder]
    slide.shapes.__iter__.return_value = iter(placeholders)
    slide.shapes.__getitem__.side_effect = lambda idx: placeholders[idx]
    slide.shapes.__len__.return_value = len(placeholders)
    
    # Mock shapes collection methods
    mock_table = MagicMock()
    mock_table_shape = MagicMock()
    mock_table_shape.table = mock_table
    slide.shapes.add_table = MagicMock(return_value=mock_table_shape)
    
    # Mock text box
    slide.shapes.add_textbox = MagicMock()
    
    # Mock notes slide
    slide.notes_slide = MagicMock()
    slide.notes_slide.notes_text_frame = MagicMock()
    
    return slide

@pytest.fixture
def mock_pptx_column_slide():
    """Create a mock PowerPoint slide with multiple column placeholders."""
    slide = MagicMock(spec=PptxSlide)
    
    # Create title placeholder
    title_placeholder = MagicMock()
    title_placeholder.is_placeholder = True
    title_placeholder.placeholder_format = MagicMock()
    title_placeholder.placeholder_format.type = PP_PLACEHOLDER.TITLE
    title_placeholder.placeholder_format.idx = 0
    title_placeholder.text_frame = MagicMock()
    
    # Create three column placeholders
    column1_placeholder = MagicMock()
    column1_placeholder.is_placeholder = True
    column1_placeholder.placeholder_format = MagicMock()
    column1_placeholder.placeholder_format.type = PP_PLACEHOLDER.BODY
    column1_placeholder.placeholder_format.idx = 1
    column1_placeholder.text_frame = MagicMock()
    column1_placeholder.text_frame.paragraphs = [MagicMock()]
    column1_placeholder.left = 100  # Position for sorting
    
    column2_placeholder = MagicMock()
    column2_placeholder.is_placeholder = True
    column2_placeholder.placeholder_format = MagicMock()
    column2_placeholder.placeholder_format.type = PP_PLACEHOLDER.BODY
    column2_placeholder.placeholder_format.idx = 2
    column2_placeholder.text_frame = MagicMock()
    column2_placeholder.text_frame.paragraphs = [MagicMock()]
    column2_placeholder.left = 300  # Position for sorting
    
    column3_placeholder = MagicMock()
    column3_placeholder.is_placeholder = True
    column3_placeholder.placeholder_format = MagicMock()
    column3_placeholder.placeholder_format.type = PP_PLACEHOLDER.BODY
    column3_placeholder.placeholder_format.idx = 3
    column3_placeholder.text_frame = MagicMock()
    column3_placeholder.text_frame.paragraphs = [MagicMock()]
    column3_placeholder.left = 500  # Position for sorting
    
    # Configurer shapes correctement avec un mock au lieu d'une liste
    slide.shapes = MagicMock()
    
    # Configurer le comportement d'itération et d'indexation
    placeholders = [title_placeholder, column1_placeholder, column2_placeholder, column3_placeholder]
    slide.shapes.__iter__.return_value = iter(placeholders)
    slide.shapes.__getitem__.side_effect = lambda idx: placeholders[idx]
    slide.shapes.__len__.return_value = len(placeholders)
    
    # Mock shapes collection methods
    mock_table = MagicMock()
    mock_table_shape = MagicMock()
    mock_table_shape.table = mock_table
    slide.shapes.add_table = MagicMock(return_value=mock_table_shape)
    
    # Mock notes slide
    slide.notes_slide = MagicMock()
    slide.notes_slide.notes_text_frame = MagicMock()
    
    return slide

@pytest.fixture
def mock_pptx_image_slide():
    """Create a mock PowerPoint slide with image and text placeholders."""
    slide = MagicMock(spec=PptxSlide)
    
    # Create title placeholder
    title_placeholder = MagicMock()
    title_placeholder.is_placeholder = True
    title_placeholder.placeholder_format = MagicMock()
    title_placeholder.placeholder_format.type = PP_PLACEHOLDER.TITLE
    title_placeholder.placeholder_format.idx = 0
    title_placeholder.text_frame = MagicMock()
    
    # Create content placeholder
    content_placeholder = MagicMock()
    content_placeholder.is_placeholder = True
    content_placeholder.placeholder_format = MagicMock()
    content_placeholder.placeholder_format.type = PP_PLACEHOLDER.BODY
    content_placeholder.placeholder_format.idx = 1
    content_placeholder.text_frame = MagicMock()
    content_placeholder.text_frame.paragraphs = [MagicMock()]
    
    # Create image placeholder
    image_placeholder = MagicMock()
    image_placeholder.is_placeholder = True
    image_placeholder.placeholder_format = MagicMock()
    image_placeholder.placeholder_format.type = PP_PLACEHOLDER.PICTURE
    image_placeholder.placeholder_format.idx = 2
    image_placeholder.text_frame = MagicMock()
    
    # Configurer shapes correctement avec un mock au lieu d'une liste
    slide.shapes = MagicMock()
    
    # Configurer le comportement d'itération et d'indexation
    placeholders = [title_placeholder, content_placeholder, image_placeholder]
    slide.shapes.__iter__.return_value = iter(placeholders)
    slide.shapes.__getitem__.side_effect = lambda idx: placeholders[idx]
    slide.shapes.__len__.return_value = len(placeholders)
    
    # Mock shapes collection methods
    mock_table = MagicMock()
    mock_table_shape = MagicMock()
    mock_table_shape.table = mock_table
    slide.shapes.add_table = MagicMock(return_value=mock_table_shape)
    
    # Mock notes slide
    slide.notes_slide = MagicMock()
    slide.notes_slide.notes_text_frame = MagicMock()
    
    return slide

@pytest.fixture
def mock_pptx_chart_slide():
    """Create a mock PowerPoint slide with chart and text placeholders."""
    slide = MagicMock(spec=PptxSlide)
    
    # Create title placeholder
    title_placeholder = MagicMock()
    title_placeholder.is_placeholder = True
    title_placeholder.placeholder_format = MagicMock()
    title_placeholder.placeholder_format.type = PP_PLACEHOLDER.TITLE
    title_placeholder.placeholder_format.idx = 0
    title_placeholder.text_frame = MagicMock()
    
    # Create content placeholder
    content_placeholder = MagicMock()
    content_placeholder.is_placeholder = True
    content_placeholder.placeholder_format = MagicMock()
    content_placeholder.placeholder_format.type = PP_PLACEHOLDER.BODY
    content_placeholder.placeholder_format.idx = 1
    content_placeholder.text_frame = MagicMock()
    content_placeholder.text_frame.paragraphs = [MagicMock()]
    
    # Create chart placeholder
    chart_placeholder = MagicMock()
    chart_placeholder.is_placeholder = True
    chart_placeholder.placeholder_format = MagicMock()
    chart_placeholder.placeholder_format.type = PP_PLACEHOLDER.CHART
    chart_placeholder.placeholder_format.idx = 2
    chart_placeholder.text_frame = MagicMock()
    
    # Configurer shapes correctement avec un mock au lieu d'une liste
    slide.shapes = MagicMock()
    
    # Configurer le comportement d'itération et d'indexation
    placeholders = [title_placeholder, content_placeholder, chart_placeholder]
    slide.shapes.__iter__.return_value = iter(placeholders)
    slide.shapes.__getitem__.side_effect = lambda idx: placeholders[idx]
    slide.shapes.__len__.return_value = len(placeholders)
    
    # Mock shapes collection methods
    mock_table = MagicMock()
    mock_table_shape = MagicMock()
    mock_table_shape.table = mock_table
    slide.shapes.add_table = MagicMock(return_value=mock_table_shape)
    
    # Mock notes slide
    slide.notes_slide = MagicMock()
    slide.notes_slide.notes_text_frame = MagicMock()
    
    return slide


@pytest.fixture
def mock_table():
    """Create a mock PowerPoint table."""
    table = MagicMock(spec=Table)
    
    # Create rows and columns
    rows = []
    for i in range(4):  # 1 header + 3 data rows
        row = MagicMock(spec=_Row)
        row.height = 100
        rows.append(row)
    table.rows = rows
    
    columns = []
    for i in range(3):  # 3 columns
        column = MagicMock(spec=_Column)
        column.width = 100
        columns.append(column)
    table.columns = columns
    
    # Create cells
    cells = []
    for i in range(4):
        row_cells = []
        for j in range(3):
            cell = MagicMock(spec=_Cell)
            cell.text_frame = MagicMock(spec=TextFrame)
            cell.text_frame.paragraphs = [MagicMock(spec=_Paragraph)]
            cell.text_frame.paragraphs[0].runs = []
            cell.fill = MagicMock()
            cell.fill.solid = MagicMock()
            cell.fill.fore_color = MagicMock()
            
            # Setup borders
            cell.border_top = MagicMock()
            cell.border_top.color = MagicMock()
            cell.border_top.color.rgb = None
            cell.border_bottom = MagicMock()
            cell.border_bottom.color = MagicMock()
            cell.border_bottom.color.rgb = None
            cell.border_left = MagicMock()
            cell.border_left.color = MagicMock()
            cell.border_left.color.rgb = None
            cell.border_right = MagicMock()
            cell.border_right.color = MagicMock()
            cell.border_right.color.rgb = None
            
            row_cells.append(cell)
        cells.append(row_cells)
    
    # Setup cell accessor
    def cell_accessor(row_idx, col_idx):
        return cells[row_idx][col_idx]
    
    table.cell = MagicMock(side_effect=cell_accessor)
    
    # Setup cells for each row
    for i, row in enumerate(table.rows):
        row.cells = cells[i]
    
    return table


@pytest.fixture
def mock_pptx_presentation():
    """Create a mock PowerPoint presentation with layouts."""
    presentation = MagicMock(spec=PptxPresentation)
    
    # Create mock layouts
    layouts = []
    for layout_name in LAYOUT_CAPABILITIES.keys():
        layout = MagicMock()
        layout.name = layout_name
        layouts.append(layout)
    
    presentation.slide_layouts = layouts
    
    # Mock slides collection
    presentation.slides = MagicMock()
    presentation.slides.add_slide = MagicMock(return_value=MagicMock(spec=PptxSlide))
    presentation.slides._sldIdLst = []  # For _clear_template_slides
    
    return presentation


@pytest.fixture
def builder():
    """Create a PPTBuilder instance."""
    return PPTBuilder()


# ===== TESTS FOR NEW BUILDER FUNCTIONALITY =====

def test_layout_capabilities_defined():
    """Test that layout capabilities are defined properly."""
    # Verify that all required layouts are defined
    required_layouts = [
        "Diapositive de titre", 
        "Introduction", 
        "Titre et texte", 
        "Titre et tableau",
        "Titre et texte 1 visuel gauche", 
        "Titre et texte 1 histogramme",
        "Titre et 3 colonnes", 
        "Chapitre 1"
    ]
    
    for layout in required_layouts:
        assert layout in LAYOUT_CAPABILITIES, f"Layout '{layout}' should be defined in LAYOUT_CAPABILITIES"
        
    # Verify structure of capabilities
    for layout, capabilities in LAYOUT_CAPABILITIES.items():
        assert "title" in capabilities
        assert "content" in capabilities
        assert "table" in capabilities
        assert "image" in capabilities
        assert "chart" in capabilities
        assert "max_blocks" in capabilities
        assert "description" in capabilities


def test_layout_placeholder_map_defined():
    """Test that layout placeholder map is defined properly."""
    # Verify that all layouts have a placeholder map
    for layout in LAYOUT_CAPABILITIES.keys():
        assert layout in LAYOUT_PLACEHOLDER_MAP, f"Layout '{layout}' should have a placeholder map"
        
    # Verify structure of placeholder maps
    for layout, placeholders in LAYOUT_PLACEHOLDER_MAP.items():
        if "Titre" in layout:  # All layouts with "Titre" should have a title placeholder
            assert "title" in placeholders


def test_validate_layout_for_content(builder, table_slide, wrong_layout_table_slide):
    """Test layout validation and correction."""
    # Test correct layout is kept
    validated_layout = builder._validate_layout_for_content(table_slide)
    assert validated_layout == "Titre et tableau", "Correct layout should not be changed"
    
    # Test wrong layout is corrected
    validated_layout = builder._validate_layout_for_content(wrong_layout_table_slide)
    assert validated_layout == "Titre et tableau", "Layout should be corrected for table content"


def test_validate_layout_for_content_block_count(builder):
    """Test layout validation based on block count."""
    # Create a slide with too many blocks for its layout
    slide = Slide(
        id="many_blocks",
        title="Too Many Blocks",
        layout_name="Titre et texte",  # Only supports 1 block
        blocks=[
            SlideBlock(id="block1", content=SlideContent(content_type=ContentType.TEXT, text="Text 1")),
            SlideBlock(id="block2", content=SlideContent(content_type=ContentType.TEXT, text="Text 2")),
            SlideBlock(id="block3", content=SlideContent(content_type=ContentType.TEXT, text="Text 3"))
        ]
    )
    
    # Validate layout
    validated_layout = builder._validate_layout_for_content(slide)
    assert validated_layout == "Titre et 3 colonnes", "Layout should be changed to support multiple blocks"


def test_fill_slide_title(builder, mock_pptx_slide):
    """Test filling slide title."""
    # Call the method
    builder._fill_slide_title(mock_pptx_slide, "Test Title")
    
    # Verify title was added
    title_placeholder = mock_pptx_slide.shapes[0]
    title_placeholder.text_frame.clear.assert_called_once()


def test_fill_title_slide(builder, mock_pptx_slide):
    """Test filling a title slide."""
    # Create a title slide
    slide = Slide(
        id="title_slide",
        title="Presentation Title",
        layout_name="Diapositive de titre",
        blocks=[
            SlideBlock(
                id="subtitle_block",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Subtitle Text"
                )
            )
        ]
    )
    
    # Mock methods
    builder._add_formatted_text = MagicMock()
    
    # Add a subtitle placeholder to the mock slide
    subtitle_placeholder = MagicMock()
    subtitle_placeholder.is_placeholder = True
    subtitle_placeholder.placeholder_format = MagicMock()
    subtitle_placeholder.placeholder_format.type = PP_PLACEHOLDER.SUBTITLE
    subtitle_placeholder.placeholder_format.idx = 1
    subtitle_placeholder.text_frame = MagicMock()
    
    # Add the subtitle placeholder to the shapes collection
    placeholders = list(mock_pptx_slide.shapes)
    placeholders.append(subtitle_placeholder)
    mock_pptx_slide.shapes.__iter__.return_value = iter(placeholders)
    mock_pptx_slide.shapes.__getitem__.side_effect = lambda idx: placeholders[idx]
    mock_pptx_slide.shapes.__len__.return_value = len(placeholders)
    
    # Call the method
    builder._fill_title_slide(mock_pptx_slide, slide)
    
    # Verify that _add_formatted_text was called on the subtitle placeholder
    assert builder._add_formatted_text.called, "The method should attempt to add formatted text"


def test_fill_content_slide(builder, mock_pptx_slide):
    """Test filling a standard content slide."""
    # Create a content slide
    slide = Slide(
        id="content_slide",
        title="Content Slide",
        layout_name="Titre et texte",
        blocks=[
            SlideBlock(
                id="text_block",
                title="Block Title",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Content text goes here."
                )
            )
        ]
    )
    
    # Mock methods
    builder._add_text_content_to_placeholder = MagicMock()
    
    # Call the method
    builder._fill_content_slide(mock_pptx_slide, slide)
    
    # Verify content was added
    assert builder._add_text_content_to_placeholder.called, "Text content should be added to the placeholder"


def test_fill_table_slide(builder, mock_pptx_slide, table_slide):
    """Test filling a slide with a table."""
    # Call the method
    builder._fill_table_slide(mock_pptx_slide, table_slide)
    
    # Verify that add_table is called on the slide's shapes
    mock_pptx_slide.shapes.add_table.assert_called_once()
    
    # Alternatively, we can mock _fill_table_with_data and verify it's called
    builder._fill_table_with_data = MagicMock()
    builder._fill_table_slide(mock_pptx_slide, table_slide)
    assert builder._fill_table_with_data.called, "Table should be filled with data"

def test_fill_column_layout_slide(builder, mock_pptx_column_slide, column_layout_slide):
    """Test filling a slide with multiple columns."""
    # Mock methods
    builder._add_block_to_placeholder = MagicMock()
    
    # Call the method
    builder._fill_column_layout_slide(mock_pptx_column_slide, column_layout_slide)
    
    # Verify blocks were added to column placeholders
    assert builder._add_block_to_placeholder.call_count == 3, "Each block should be added to a column placeholder"


def test_fill_image_layout_slide(builder, mock_pptx_image_slide, image_layout_slide):
    """Test filling a slide with image and text."""
    # Mock methods
    builder._add_text_content_to_placeholder = MagicMock()
    
    # Call the method
    builder._fill_image_layout_slide(mock_pptx_image_slide, image_layout_slide)
    
    # Verify content was added
    # Check that the _add_text_content_to_placeholder method was called
    assert builder._add_text_content_to_placeholder.called, "Text content should be added to the text placeholder"
    
    # Check that some text was added to the image placeholder
    placeholder_0 = mock_pptx_image_slide.shapes[0]  # Title
    placeholder_1 = mock_pptx_image_slide.shapes[1]  # Content
    placeholder_2 = mock_pptx_image_slide.shapes[2]  # Image
    
    # Ensure that we're working with the picture placeholder
    assert placeholder_2.placeholder_format.type == PP_PLACEHOLDER.PICTURE
    assert hasattr(placeholder_2, 'text_frame'), "Image placeholder should have a text_frame"


def test_fill_chart_layout_slide(builder, mock_pptx_chart_slide, chart_layout_slide):
    """Test filling a slide with chart and text."""
    # Mock methods
    builder._add_text_content_to_placeholder = MagicMock()
    
    # Call the method
    builder._fill_chart_layout_slide(mock_pptx_chart_slide, chart_layout_slide)
    
    # Verify content was added
    # Check that the _add_text_content_to_placeholder method was called
    assert builder._add_text_content_to_placeholder.called, "Text content should be added to the text placeholder"
    
    # Check that some text was added to the chart placeholder
    placeholder_0 = mock_pptx_chart_slide.shapes[0]  # Title
    placeholder_1 = mock_pptx_chart_slide.shapes[1]  # Content
    placeholder_2 = mock_pptx_chart_slide.shapes[2]  # Chart
    
    # Ensure that we're working with the chart placeholder
    assert placeholder_2.placeholder_format.type == PP_PLACEHOLDER.CHART
    assert hasattr(placeholder_2, 'text_frame'), "Chart placeholder should have a text_frame"


def test_add_block_to_placeholder(builder, mock_pptx_slide):
    """Test adding a content block to a placeholder."""
    # Create a content block
    block = SlideBlock(
        id="text_block",
        title="Block Title",
        content=SlideContent(
            content_type=ContentType.TEXT,
            text="Block content text."
        )
    )
    
    # Mock methods
    builder._add_text_content_to_placeholder = MagicMock()
    
    # Call the method
    builder._add_block_to_placeholder(mock_pptx_slide.shapes[1], block)
    
    # Verify content was added
    mock_pptx_slide.shapes[1].text_frame.clear.assert_called_once()
    assert builder._add_text_content_to_placeholder.called, "Text content should be added to the placeholder"


def test_add_text_content_to_placeholder(builder, mock_pptx_slide):
    """Test adding text content to a placeholder."""
    # Mock methods
    builder._add_formatted_text_to_paragraph = MagicMock()
    
    # Call the method
    builder._add_text_content_to_placeholder(mock_pptx_slide.shapes[1], "Line 1\nLine 2\nLine 3")
    
    # Verify content was added
    assert mock_pptx_slide.shapes[1].text_frame.add_paragraph.call_count == 3, "Three paragraphs should be added"
    assert builder._add_formatted_text_to_paragraph.call_count == 3, "Formatted text should be added to each paragraph"


def test_add_bullet_points_to_placeholder(builder, mock_pptx_slide):
    """Test adding bullet points to a placeholder."""
    # Mock methods
    builder._add_formatted_text_to_paragraph = MagicMock()
    
    # Call the method
    builder._add_bullet_points_to_placeholder(
        mock_pptx_slide.shapes[1], 
        ["Point 1", "Point 2", "Point 3"],
        True  # as_bullets=True
    )
    
    # Verify content was added
    assert mock_pptx_slide.shapes[1].text_frame.add_paragraph.call_count == 3, "Three paragraphs should be added"
    assert builder._add_formatted_text_to_paragraph.call_count == 3, "Formatted text should be added to each paragraph"


def test_add_table_to_slide(builder, mock_pptx_slide, sample_table_data):
    """Test adding a table to a slide."""
    # Mock methods
    builder._fill_table_with_data = MagicMock()
    
    # Call the method
    result = builder._add_table_to_slide(
        mock_pptx_slide,
        sample_table_data.headers,
        sample_table_data.rows
    )
    
    # Verify table was added
    mock_pptx_slide.shapes.add_table.assert_called_once()
    assert builder._fill_table_with_data.called, "Table should be filled with data"
    assert result is not None, "Should return the created table"


def test_add_table_to_slide_empty_data(builder, mock_pptx_slide):
    """Test adding a table with empty data."""
    # Call the method with empty data
    result = builder._add_table_to_slide(
        mock_pptx_slide,
        [],  # Empty headers
        []   # Empty rows
    )
    
    # Verify no table was added
    mock_pptx_slide.shapes.add_table.assert_not_called()
    assert result is None, "Should return None for empty data"


def test_get_style_from_headers(builder):
    """Test extracting style from table headers."""
    # Test with style marker
    headers_with_style = ["Column 1", "Column 2", "style:accent1"]
    style = builder._get_style_from_headers(headers_with_style)
    assert style == "accent1", "Style should be extracted from headers"
    
    # Test without style marker
    headers_without_style = ["Column 1", "Column 2", "Column 3"]
    style = builder._get_style_from_headers(headers_without_style)
    assert style is None, "Style should be None when no style marker is present"
    
    # Test with empty headers
    style = builder._get_style_from_headers([])
    assert style is None, "Style should be None for empty headers"


def test_fill_table_with_data(builder, mock_table, sample_table_data):
    """Test filling a table with data."""
    # Mock methods
    builder._format_table_cell = MagicMock()
    builder._apply_table_style = MagicMock()
    
    # Call the method
    builder._fill_table_with_data(
        mock_table,
        sample_table_data.headers,
        sample_table_data.rows
    )
    
    # Verify cells were formatted
    # Headers (1 row) + data rows (3 rows) * columns (3) = 12 cells
    assert builder._format_table_cell.call_count == 12, "All cells should be formatted"
    
    # Verify style was applied
    builder._apply_table_style.assert_called_once()


def test_format_table_cell(builder):
    """Test formatting a table cell."""
    # Create a mock cell
    cell = MagicMock(spec=_Cell)
    cell.text_frame = MagicMock(spec=TextFrame)
    cell.text_frame.paragraphs = [MagicMock(spec=_Paragraph)]
    cell.text_frame.paragraphs[0].runs = []
    cell.fill = MagicMock()
    cell.fill.solid = MagicMock()
    cell.fill.fore_color = MagicMock()
    
    # Add border attributes
    cell.border_top = MagicMock()
    cell.border_top.color = MagicMock()
    cell.border_top.color.rgb = None
    cell.border_bottom = MagicMock()
    cell.border_bottom.color = MagicMock()
    cell.border_bottom.color.rgb = None
    cell.border_left = MagicMock()
    cell.border_left.color = MagicMock()
    cell.border_left.color.rgb = None
    cell.border_right = MagicMock()
    cell.border_right.color = MagicMock()
    cell.border_right.color.rgb = None
    
    # Mock methods
    builder._add_formatted_text = MagicMock()
    
    # Format header cell
    style_preset = builder.TABLE_STYLES["default"]
    builder._format_table_cell(cell, "Header Text", is_header=True, style_preset=style_preset)
    
    # Verify formatting was applied
    builder._add_formatted_text.assert_called_once_with(cell.text_frame, "Header Text")
    assert cell.text_frame.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE
    
    # Header cells should have center alignment
    assert cell.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER
    
    # Header cells should use the header background color
    assert cell.fill.solid.called
    
    # Reset mocks
    builder._add_formatted_text.reset_mock()
    cell.fill.solid.reset_mock()
    
    # Format data cell
    builder._format_table_cell(cell, "Data Text", is_header=False, style_preset=style_preset)
    
    # Verify formatting was applied
    builder._add_formatted_text.assert_called_once_with(cell.text_frame, "Data Text")
    
    # Data cells should have left alignment
    assert cell.text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT


def test_apply_table_style(builder, mock_table):
    """Test applying a style to a table."""
    # Mock methods
    builder._hex_to_rgb = MagicMock(return_value=RGBColor(0, 0, 0))
    
    # Apply style
    style_preset = builder.TABLE_STYLES["default"]
    builder._apply_table_style(mock_table, style_preset)
    
    # Check if column widths were set
    for column in mock_table.columns:
        assert hasattr(column, 'width')
    
    # Check if row heights were set
    for row in mock_table.rows:
        assert hasattr(row, 'height')
    
    # Check if borders were set on all cells
    for row in mock_table.rows:
        for cell in row.cells:
            # Top border
            assert cell.border_top.color.rgb is not None
            # Bottom border
            assert cell.border_bottom.color.rgb is not None
            # Left border
            assert cell.border_left.color.rgb is not None
            # Right border
            assert cell.border_right.color.rgb is not None


def test_hex_to_rgb(builder):
    """Test converting hex color to RGB."""
    
    def extract_rgb_values(rgb_color):
        """Extract the RGB values from any RGBColor-like object."""
        try:
            # Try accessing as a tuple or sequence
            return tuple(rgb_color)
        except TypeError:
            try:
                # Try accessing as r, g, b attributes
                return (rgb_color.r, rgb_color.g, rgb_color.b)
            except AttributeError:
                try:
                    # Try accessing as red, green, blue attributes
                    return (rgb_color.red, rgb_color.green, rgb_color.blue)
                except AttributeError:
                    try:
                        # Try accessing as rgb attribute
                        return rgb_color.rgb
                    except AttributeError:
                        raise ValueError(f"Cannot extract RGB values from {rgb_color}")
    
    # Test with standard hex code (6 digits)
    rgb = builder._hex_to_rgb("FF0000")
    rgb_values = extract_rgb_values(rgb)
    assert rgb_values == (255, 0, 0), f"Expected (255, 0, 0), got {rgb_values}"
    
    # Test with hex code including # prefix
    rgb = builder._hex_to_rgb("#00FF00")
    rgb_values = extract_rgb_values(rgb)
    assert rgb_values == (0, 255, 0), f"Expected (0, 255, 0), got {rgb_values}"
    
    # Test with short hex code (3 digits)
    rgb = builder._hex_to_rgb("00F")
    rgb_values = extract_rgb_values(rgb)
    assert rgb_values == (0, 0, 255), f"Expected (0, 0, 255), got {rgb_values}"
    
    # Test with invalid hex code
    rgb = builder._hex_to_rgb("invalid")
    rgb_values = extract_rgb_values(rgb)
    assert rgb_values == (0, 0, 0), f"Expected (0, 0, 0), got {rgb_values}"


def test_validate_layout_for_unknown_layout(builder):
    """Test layout validation with unknown layout."""
    # Create a slide with an unknown layout
    slide = Slide(
        id="unknown_layout",
        title="Unknown Layout",
        layout_name="NonExistentLayout",
        blocks=[
            SlideBlock(
                id="block1",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Text content"
                )
            )
        ]
    )
    
    # Validate layout
    with patch('doc2pptx.ppt.builder_v3.logger.warning') as mock_warning:
        validated_layout = builder._validate_layout_for_content(slide)
        
        # Should fall back to default layout
        assert validated_layout == "Titre et texte", "Unknown layout should be replaced with default layout"
        
        # Should log a warning
        assert mock_warning.called, "Warning should be logged for unknown layout"


@patch('doc2pptx.ppt.builder_v3.logger.warning')
def test_fill_slide_with_different_layouts(mock_warning, builder, mock_pptx_presentation, mock_pptx_slide):
    """Test filling slides with different layouts."""
    # Create different slide types
    title_slide = Slide(
        id="title",
        title="Title Slide",
        layout_name="Diapositive de titre",
        blocks=[
            SlideBlock(
                id="block1",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Subtitle text"
                )
            )
        ]
    )
    
    content_slide = Slide(
        id="content",
        title="Content Slide",
        layout_name="Titre et texte",
        blocks=[
            SlideBlock(
                id="block2",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Content text"
                )
            )
        ]
    )
    
    table_slide = Slide(
        id="table",
        title="Table Slide",
        layout_name="Titre et tableau",
        blocks=[
            SlideBlock(
                id="block3",
                content=SlideContent(
                    content_type=ContentType.TABLE,
                    table=TableData(
                        headers=["Col1", "Col2"],
                        rows=[["A", "B"], ["C", "D"]]
                    )
                )
            )
        ]
    )
    
    # Mock methods to avoid actual implementation
    builder._fill_title_slide = MagicMock()
    builder._fill_content_slide = MagicMock()
    builder._fill_table_slide = MagicMock()
    builder._fill_column_layout_slide = MagicMock()
    builder._fill_image_layout_slide = MagicMock()
    builder._fill_chart_layout_slide = MagicMock()
    builder._fill_slide_title = MagicMock()
    
    # Create a dummy section
    section = Section(
        id="section1",
        title="Test Section",
        type=SectionType.CONTENT,
        slides=[]
    )
    
    # Test filling each slide type
    builder._fill_slide(mock_pptx_slide, title_slide, section)
    assert builder._fill_title_slide.called, "Title slide filling method should be called"
    builder._fill_title_slide.reset_mock()
    
    builder._fill_slide(mock_pptx_slide, content_slide, section)
    assert builder._fill_content_slide.called, "Content slide filling method should be called"
    builder._fill_content_slide.reset_mock()
    
    builder._fill_slide(mock_pptx_slide, table_slide, section)
    assert builder._fill_table_slide.called, "Table slide filling method should be called"


@patch('doc2pptx.ppt.builder_v3.logger.warning')
def test_fill_table_with_data_error_handling(mock_warning, builder, mock_table):
    """Test error handling in fill_table_with_data."""
    # Create inconsistent table data (more headers than columns)
    headers = ["Col1", "Col2", "Col3", "Col4"]  # 4 headers
    rows = [
        ["A", "B", "C"],  # 3 values
        ["D", "E", "F"]   # 3 values
    ]
    
    # Mock methods
    builder._format_table_cell = MagicMock()
    builder._apply_table_style = MagicMock()
    
    # Force mock_table to have only 3 columns
    mock_table.columns = [MagicMock(), MagicMock(), MagicMock()]
    
    # Call the method
    builder._fill_table_with_data(mock_table, headers, rows)
    
    # Should log a warning (not error)
    assert mock_warning.called, "Warning should be logged for dimension mismatch"
    
    # Should still format cells up to the available columns
    assert builder._format_table_cell.called, "Table cells should still be formatted"


@patch('doc2pptx.ppt.builder_v3.TemplateLoader.load_template')
@patch('doc2pptx.ppt.builder_v3.TemplateLoader.analyze_template')
def test_build_integration(mock_analyze_template, mock_load_template, builder, mock_pptx_presentation):
    """Integration test for the build method."""
    # Setup mocks
    mock_analyze_template.return_value = "template_info"
    mock_load_template.return_value = mock_pptx_presentation
    
    # Create a simple presentation
    slide1 = Slide(
        id="text_slide",
        title="Text Slide",
        layout_name="Titre et texte",
        blocks=[
            SlideBlock(
                id="block1",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Text content"
                )
            )
        ]
    )
    
    slide2 = Slide(
        id="table_slide",
        title="Table Slide",
        layout_name="Titre et tableau",
        blocks=[
            SlideBlock(
                id="block2",
                content=SlideContent(
                    content_type=ContentType.TABLE,
                    table=TableData(
                        headers=["Col1", "Col2"],
                        rows=[["A", "B"], ["C", "D"]]
                    )
                )
            )
        ]
    )
    
    section = Section(
        id="section1",
        title="Test Section",
        type=SectionType.CONTENT,
        slides=[slide1, slide2]
    )
    
    presentation = Presentation(
        id="test_pres",
        title="Test Presentation",
        template_path=Path("test_template.pptx"),
        sections=[section]
    )
    
    # Mock methods
    builder._create_slide = MagicMock(return_value=mock_pptx_slide)
    builder._fill_slide = MagicMock()
    builder._clear_template_slides = MagicMock()
    
    # Call the method
    output_path = builder.build(presentation, "output.pptx")
    
    # Check that slides were created and filled
    assert builder._create_slide.call_count == 2, "Two slides should be created"
    assert builder._fill_slide.call_count == 2, "Two slides should be filled"
    assert mock_pptx_presentation.save.called, "Presentation should be saved"
    
    # Check output path
    assert output_path == Path("output.pptx"), "Output path should be correct"


def test_build_with_slide_reordering(builder, mock_pptx_presentation):
    """Test building a presentation with slides that need layout correction."""
    # Setup mocks
    with patch('doc2pptx.ppt.builder_v3.TemplateLoader.load_template', return_value=mock_pptx_presentation), \
         patch('doc2pptx.ppt.builder_v3.TemplateLoader.analyze_template', return_value="template_info"), \
         patch.object(builder, '_create_slide', return_value=MagicMock()), \
         patch.object(builder, '_fill_slide'), \
         patch.object(builder, '_validate_layout_for_content') as mock_validate, \
         patch.object(builder, '_clear_template_slides'):
        
        # Make validate_layout_for_content return "Titre et tableau" for all slides
        mock_validate.return_value = "Titre et tableau"
        
        # Create a presentation with slides that need correction
        slide1 = Slide(
            id="slide1",
            title="Slide 1",
            layout_name="Titre et texte",  # Will be corrected to "Titre et tableau"
            blocks=[
                SlideBlock(
                    id="block1",
                    content=SlideContent(
                        content_type=ContentType.TABLE,
                        table=TableData(
                            headers=["Col1", "Col2"],
                            rows=[["A", "B"], ["C", "D"]]
                        )
                    )
                )
            ]
        )
        
        slide2 = Slide(
            id="slide2",
            title="Slide 2",
            layout_name="Titre et texte",  # Will be corrected to "Titre et tableau"
            blocks=[
                SlideBlock(
                    id="block2",
                    content=SlideContent(
                        content_type=ContentType.TABLE,
                        table=TableData(
                            headers=["Col1", "Col2"],
                            rows=[["E", "F"], ["G", "H"]]
                        )
                    )
                )
            ]
        )
        
        section = Section(
            id="section1",
            title="Test Section",
            type=SectionType.CONTENT,
            slides=[slide1, slide2]
        )
        
        presentation = Presentation(
            id="test_pres",
            title="Test Presentation",
            template_path=Path("test_template.pptx"),
            sections=[section]
        )
        
        # Call the method
        builder.build(presentation, "output.pptx")
        
        # Check that validate_layout_for_content was called for both slides
        assert mock_validate.call_count == 2, "Layout validation should be called for both slides"
        
        # Check that _create_slide was called with the corrected layout
        builder._create_slide.assert_any_call(mock_pptx_presentation, "Titre et tableau")


def test_fill_slide_with_notes(builder, mock_pptx_slide):
    """Test filling slide with speaker notes."""
    # Create a slide with notes
    slide = Slide(
        id="slide_with_notes",
        title="Slide with Notes",
        layout_name="Titre et texte",
        blocks=[
            SlideBlock(
                id="block1",
                content=SlideContent(
                    content_type=ContentType.TEXT,
                    text="Text content"
                )
            )
        ],
        notes="These are speaker notes for the slide."
    )
    
    section = Section(
        id="section1",
        title="Test Section",
        type=SectionType.CONTENT,
        slides=[slide]
    )
    
    # Mock methods
    builder._fill_content_slide = MagicMock()
    builder._fill_slide_title = MagicMock()
    
    # Call the method
    builder._fill_slide(mock_pptx_slide, slide, section)
    
    # Check that notes were added
    assert mock_pptx_slide.notes_slide.notes_text_frame.text == "These are speaker notes for the slide."


def test_build_with_template_path_from_presentation(builder):
    """Test building a presentation using the template path from the presentation."""
    # Create a presentation with a template path
    presentation = Presentation(
        id="test_pres",
        title="Test Presentation",
        template_path=Path("pres_template.pptx"),
        sections=[]
    )
    
    # Mock methods
    with patch('doc2pptx.ppt.builder_v3.TemplateLoader.load_template') as mock_load, \
         patch('doc2pptx.ppt.builder_v3.TemplateLoader.analyze_template') as mock_analyze, \
         patch.object(builder, '_clear_template_slides'), \
         patch('pathlib.Path.mkdir'), \
         patch('pptx.presentation.Presentation.save'):
        
        mock_load.return_value = MagicMock()
        mock_analyze.return_value = "template_info"
        
        # Call the method
        builder.build(presentation, "output.pptx")
        
        # Check that the template was loaded from the presentation path
        mock_load.assert_called_once_with(Path("pres_template.pptx"))


def test_build_with_no_template_path(builder):
    """Test building a presentation with no template path."""
    # Create a presentation without a template path
    presentation = Presentation(
        id="test_pres",
        title="Test Presentation",
        template_path=None,
        sections=[]
    )
    
    # Builder also has no template path
    builder.template_path = None
    
    # Call the method - should raise ValueError
    with pytest.raises(ValueError, match="No template path provided"):
        builder.build(presentation, "output.pptx")


def test_needs_section_header(builder):
    """Test the _needs_section_header method."""
    # Create a section
    section = Section(
        id="section1",
        title="Test Section",
        type=SectionType.CONTENT,
        slides=[
            Slide(
                id="slide1",
                title="Slide 1",
                layout_name="Titre et texte",
                blocks=[]
            )
        ]
    )
    
    # Call the method - should return False according to the current implementation
    result = builder._needs_section_header(section)
    assert result is False, "Current implementation should return False"


# Additional tests for edge cases and error handling

def test_add_table_to_slide_error(builder, mock_pptx_slide):
    """Test error handling in _add_table_to_slide."""
    # Mock add_table to raise an exception
    mock_pptx_slide.shapes.add_table.side_effect = Exception("Test error")
    
    # Call the method with valid data
    with patch('doc2pptx.ppt.builder_v3.logger.error') as mock_error:
        result = builder._add_table_to_slide(
            mock_pptx_slide,
            ["Col1", "Col2"],
            [["A", "B"], ["C", "D"]]
        )
        
        # Should log an error
        assert mock_error.called, "Error should be logged"
        
        # Should return None
        assert result is None, "Should return None on error"


def test_closest_highlight_color(builder):
    """Test _closest_highlight_color method."""
    # Test finding closest highlight color
    color = builder._closest_highlight_color(255, 0, 0)
    assert color == "red", "Should find red as closest color to (255,0,0)"
    
    color = builder._closest_highlight_color(0, 255, 0)
    assert color == "green", "Should find green as closest color to (0,255,0)"
    
    color = builder._closest_highlight_color(0, 0, 255)
    assert color == "blue", "Should find blue as closest color to (0,0,255)"
    
    # Test with in-between color
    color = builder._closest_highlight_color(128, 128, 0)
    assert color == "darkYellow", "Should find darkYellow as closest color to (128,128,0)"
    
    # Test with color not exactly matching any preset
    color = builder._closest_highlight_color(100, 100, 100)
    assert color in ["darkGray", "lightGray"], "Should find gray as closest color to (100,100,100)"


# tests/unit/test_tabledata.py
import pytest
from doc2pptx.core.models import TableData

def test_tabledata_accepts_style_header():
    data = TableData(
        headers=["Col1", "Col2", "style:minimal"],
        rows=[["A", "B"], ["C", "D"]],
    )
    assert len(data.headers) == 3  # le header style est bien présent

def test_tabledata_rejects_incorrect_rows():
    with pytest.raises(ValueError):
        TableData(
            headers=["H1", "H2"],
            rows=[["only one cell"]],
        )

def test_bullets_always_visible(prs_template, builder):
    # 1️⃣ Corps avec style → la puce doit déjà exister
    slide = prs_template.slides.add_slide(prs_template.slide_layouts[1])
    body_ph = next(sh for sh in slide.shapes if sh.is_placeholder and hasattr(sh, "text_frame"))
    builder._add_bullet_points_to_placeholder(body_ph, ["alpha", "beta"])
    assert "•" in body_ph.text_frame.text

    # 2️⃣ TextBox sans style → la fonction doit injecter la puce
    slide2 = prs_template.slides.add_slide(prs_template.slide_layouts[6])
    tb = slide2.shapes.add_textbox(0, 0, Emu(9144000), Emu(4572000))
    builder._add_bullet_points_to_placeholder(tb, ["one", "two"])
    assert "•" in tb.text_frame.text

def test_apply_table_style_has_only_int_emu(tmp_path, base_template):
    prs = PptxPresentation(str(base_template))
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    table_shape = slide.shapes.add_table(2, 2, Emu(0), Emu(0), Emu(9144000), Emu(6858000))  # 10x7,5 in
    table = table_shape.table

    builder = PPTBuilder()
    builder._apply_table_style(table, builder.TABLE_STYLES["default"])

    # Tous les widths/heights doivent être int/Emu
    for col in table.columns:
        assert isinstance(col.width, int)
    for row in table.rows:
        assert isinstance(row.height, int)


def test_bullet_points_no_attribute_error(tmp_path, base_template):
    prs = PptxPresentation(str(base_template))
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # layout avec BODY
    placeholder = next(sh for sh in slide.shapes if sh.is_placeholder and hasattr(sh, "text_frame"))
    builder = PPTBuilder()

    builder._add_bullet_points_to_placeholder(
        placeholder,
        ["First", "Second"],
        as_bullets=True,
    )

    # Il doit y avoir deux paragraphes et le premier doit être niveau 0
    assert len(placeholder.text_frame.paragraphs) == 2
    assert placeholder.text_frame.paragraphs[0].level == 0


if __name__ == "__main__":
    pytest.main(["-v"])