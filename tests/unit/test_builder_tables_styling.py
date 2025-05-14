"""
Unit tests for table styling functionality in PPTBuilder.
"""
import os
import pytest
from pathlib import Path
from unittest.mock import MagicMock, patch, PropertyMock

from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.table import Table, _Cell, _Row, _Column

from doc2pptx.core.models import TableData
from doc2pptx.ppt.builder_v3 import PPTBuilder


@pytest.fixture
def mock_cell():
    """Create a mock table cell."""
    cell = MagicMock(spec=_Cell)
    cell.text_frame = MagicMock()
    cell.text_frame.paragraphs = [MagicMock()]
    cell.text_frame.paragraphs[0].runs = []
    cell.fill = MagicMock()
    cell.fill.solid = MagicMock()
    cell.fill.fore_color = MagicMock()
    
    # Setup borders
    cell.border_top = MagicMock()
    cell.border_top.color = MagicMock()
    cell.border_top.color.rgb = None
    cell.border_bottom = MagicMock()
    cell.border_bottom.color = MagicMock()
    cell.border_bottom.color.rgb = None
    cell.border_left = MagicMock()
    cell.border_left.color = MagicMock()
    cell.border_left.color.rgb = None
    cell.border_right = MagicMock()
    cell.border_right.color = MagicMock()
    cell.border_right.color.rgb = None
    
    return cell


def test_hex_to_rgb():
    """Test the hex to RGB color conversion with various inputs."""
    builder = PPTBuilder()
    
    # Test with standard 6-digit hex
    rgb = builder._hex_to_rgb("FF0000")
    # RGBColor in python-pptx stores values as a tuple, not as separate attributes
    assert rgb[0] == 255
    assert rgb[1] == 0
    assert rgb[2] == 0
    
    # Test with 6-digit hex with # prefix
    rgb = builder._hex_to_rgb("#00FF00")
    assert rgb[0] == 0
    assert rgb[1] == 255
    assert rgb[2] == 0
    
    # Test with 3-digit hex
    rgb = builder._hex_to_rgb("00F")
    assert rgb[0] == 0
    assert rgb[1] == 0
    assert rgb[2] == 255
    
    # Test with uppercase hex
    rgb = builder._hex_to_rgb("FFFFFF")
    assert rgb[0] == 255
    assert rgb[1] == 255
    assert rgb[2] == 255
    
    # Test with lowercase hex
    rgb = builder._hex_to_rgb("ffffff")
    assert rgb[0] == 255
    assert rgb[1] == 255
    assert rgb[2] == 255
    
    # Test with invalid hex (should return black)
    rgb = builder._hex_to_rgb("XYZ")
    assert rgb[0] == 0
    assert rgb[1] == 0
    assert rgb[2] == 0
    
    # Test with empty string (should return black)
    rgb = builder._hex_to_rgb("")
    assert rgb[0] == 0
    assert rgb[1] == 0
    assert rgb[2] == 0


def test_format_table_cell_header(mock_cell):
    """Test formatting a header cell in a table."""
    builder = PPTBuilder()
    builder._add_formatted_text = MagicMock()
    
    # Test header cell formatting
    style_preset = builder.TABLE_STYLES["default"]
    builder._format_table_cell(mock_cell, "Header Text", is_header=True, style_preset=style_preset)
    
    # Verify header-specific formatting was applied
    builder._add_formatted_text.assert_called_once_with(mock_cell.text_frame, "Header Text")
    assert mock_cell.text_frame.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE
    assert mock_cell.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER
    mock_cell.fill.solid.assert_called_once()


def test_format_table_cell_data(mock_cell):
    """Test formatting a data cell in a table."""
    builder = PPTBuilder()
    builder._add_formatted_text = MagicMock()
    
    # Test data cell formatting
    style_preset = builder.TABLE_STYLES["default"]
    builder._format_table_cell(
        mock_cell,
        "Data Text",
        is_header=False,
        is_alternate_row=True,
        style_preset=style_preset
    )
    
    # Verify data-specific formatting was applied
    builder._add_formatted_text.assert_called_once_with(mock_cell.text_frame, "Data Text")
    assert mock_cell.text_frame.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE
    assert mock_cell.text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT


def test_format_table_cell_with_different_styles():
    """Test formatting cells with different style presets."""
    builder = PPTBuilder()
    builder._add_formatted_text = MagicMock()
    
    # Test with different style presets
    style_presets = ["default", "minimal", "grid", "accent1", "accent2", "accent3"]
    
    for style_name in style_presets:
        mock_cell = MagicMock(spec=_Cell)
        mock_cell.text_frame = MagicMock()
        mock_cell.text_frame.paragraphs = [MagicMock()]
        mock_cell.text_frame.paragraphs[0].runs = []
        mock_cell.fill = MagicMock()
        mock_cell.fill.solid = MagicMock()
        mock_cell.fill.fore_color = MagicMock()
        
        # Setup borders
        mock_cell.border_top = MagicMock()
        mock_cell.border_top.color = MagicMock()
        mock_cell.border_bottom = MagicMock()
        mock_cell.border_bottom.color = MagicMock()
        mock_cell.border_left = MagicMock()
        mock_cell.border_left.color = MagicMock()
        mock_cell.border_right = MagicMock()
        mock_cell.border_right.color = MagicMock()
        
        style_preset = builder.TABLE_STYLES[style_name]
        builder._format_table_cell(mock_cell, f"Cell with {style_name} style", is_header=True, style_preset=style_preset)
        
        # Verify basic formatting was applied
        builder._add_formatted_text.assert_called_with(mock_cell.text_frame, f"Cell with {style_name} style")
        assert mock_cell.text_frame.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE


def test_apply_table_style(mock_cell):
    """Test applying table styles to a table."""
    builder = PPTBuilder()
    builder._hex_to_rgb = MagicMock(return_value=RGBColor(255, 0, 0))
    
    # Create a mock table
    mock_table = MagicMock(spec=Table)
    
    # Create mock rows and columns
    rows = [MagicMock(spec=_Row) for _ in range(3)]
    columns = [MagicMock(spec=_Column) for _ in range(3)]
    
    # Add cells to rows
    for row in rows:
        cells = [MagicMock(spec=_Cell) for _ in range(3)]
        for cell in cells:
            # Setup borders
            cell.border_top = MagicMock()
            cell.border_top.color = MagicMock()
            cell.border_bottom = MagicMock()
            cell.border_bottom.color = MagicMock()
            cell.border_left = MagicMock()
            cell.border_left.color = MagicMock()
            cell.border_right = MagicMock()
            cell.border_right.color = MagicMock()
        row.cells = cells
    
    mock_table.rows = rows
    mock_table.columns = columns
    
    # Test with default style
    style_preset = builder.TABLE_STYLES["default"]
    builder._apply_table_style(mock_table, style_preset)
    
    # Verify border colors were set
    builder._hex_to_rgb.assert_called()
    
    # Test with empty table
    mock_empty_table = MagicMock(spec=Table)
    mock_empty_table.rows = []
    mock_empty_table.columns = []
    
    with patch('doc2pptx.ppt.builder_v3.logger.warning') as mock_warning:
        builder._apply_table_style(mock_empty_table, style_preset)
        mock_warning.assert_called_once()


def test_closest_highlight_color():
    """Test finding the closest PowerPoint highlight color."""
    builder = PPTBuilder()
    
    # Test primary colors
    assert builder._closest_highlight_color(255, 0, 0) == "red"
    assert builder._closest_highlight_color(0, 255, 0) == "green"
    assert builder._closest_highlight_color(0, 0, 255) == "blue"
    
    # Test secondary colors
    assert builder._closest_highlight_color(255, 255, 0) == "yellow"
    assert builder._closest_highlight_color(0, 255, 255) == "cyan"
    assert builder._closest_highlight_color(255, 0, 255) == "magenta"
    
    # Test grayscale
    assert builder._closest_highlight_color(0, 0, 0) == "black"
    assert builder._closest_highlight_color(255, 255, 255) == "white"
    assert builder._closest_highlight_color(128, 128, 128) == "darkGray"
    assert builder._closest_highlight_color(192, 192, 192) == "lightGray"
    
    # Test dark variations
    assert builder._closest_highlight_color(128, 0, 0) == "darkRed"
    assert builder._closest_highlight_color(0, 128, 0) == "darkGreen"
    assert builder._closest_highlight_color(0, 0, 128) == "darkBlue"


def test_get_style_from_headers():
    """Test extracting style information from table headers."""
    builder = PPTBuilder()
    
    # Test with style marker
    headers = ["Column 1", "Column 2", "style:accent1"]
    style = builder._get_style_from_headers(headers)
    assert style == "accent1"
    
    # Test with style in different position
    headers = ["style:accent2", "Column 1", "Column 2"]
    style = builder._get_style_from_headers(headers)
    assert style is None  # Should only look at the last header
    
    # Test with no style marker
    headers = ["Column 1", "Column 2", "Column 3"]
    style = builder._get_style_from_headers(headers)
    assert style is None
    
    # Test with empty headers
    headers = []
    style = builder._get_style_from_headers(headers)
    assert style is None
    
    # Test with non-string header
    headers = ["Column 1", "Column 2", 123]
    style = builder._get_style_from_headers(headers)
    assert style is None


def test_fill_table_with_data_dimension_mismatch():
    """Test handling of dimension mismatches when filling table data."""
    builder = PPTBuilder()
    builder._format_table_cell = MagicMock()
    builder._apply_table_style = MagicMock()
    
    # Create a mock table with specific dimensions
    mock_table = MagicMock(spec=Table)
    mock_table.rows = [MagicMock(spec=_Row) for _ in range(2)]  # Only 2 rows
    mock_table.columns = [MagicMock(spec=_Column) for _ in range(2)]  # Only 2 columns
    
    # Setup cell accessor
    def mock_cell(row_idx, col_idx):
        mock_cell = MagicMock(spec=_Cell)
        mock_cell.text_frame = MagicMock()
        mock_cell.text_frame.paragraphs = [MagicMock()]
        return mock_cell
    
    mock_table.cell = MagicMock(side_effect=mock_cell)
    
    # Data with more rows and columns than the table
    headers = ["Col1", "Col2", "Col3"]  # 3 columns, but table has 2
    rows = [
        ["A", "B", "C"],  # 3 columns, but table has 2
        ["D", "E", "F"],  # 3 columns, but table has 2
        ["G", "H", "I"]   # 3rd row, but table has 2 rows
    ]
    
    # Test with dimension mismatch
    with patch('doc2pptx.ppt.builder_v3.logger.warning') as mock_warning:
        builder._fill_table_with_data(mock_table, headers, rows)
        mock_warning.assert_called_once()
    
    # Verify that _format_table_cell was called for all available cells
    assert builder._format_table_cell.call_count == 4  # 2 rows * 2 cols


def test_table_with_different_style_presets():
    """Test creating tables with different style presets."""
    builder = PPTBuilder()
    
    # Test with all available style presets
    style_presets = list(builder.TABLE_STYLES.keys())
    
    mock_pptx_slide = MagicMock()
    mock_table_shape = MagicMock()
    mock_table = MagicMock(spec=Table)
    mock_table_shape.table = mock_table
    mock_pptx_slide.shapes.add_table.return_value = mock_table_shape
    
    # Setup mock table
    mock_table.rows = [MagicMock(spec=_Row) for _ in range(2)]
    mock_table.columns = [MagicMock(spec=_Column) for _ in range(2)]
    
    # Setup cell accessor
    def mock_cell(row_idx, col_idx):
        mock_cell = MagicMock(spec=_Cell)
        mock_cell.text_frame = MagicMock()
        mock_cell.text_frame.paragraphs = [MagicMock()]
        mock_cell.fill = MagicMock()
        mock_cell.fill.solid = MagicMock()
        mock_cell.fill.fore_color = MagicMock()
        
        # Setup borders
        mock_cell.border_top = MagicMock()
        mock_cell.border_top.color = MagicMock()
        mock_cell.border_bottom = MagicMock()
        mock_cell.border_bottom.color = MagicMock()
        mock_cell.border_left = MagicMock()
        mock_cell.border_left.color = MagicMock()
        mock_cell.border_right = MagicMock()
        mock_cell.border_right.color = MagicMock()
        return mock_cell
    
    mock_table.cell = MagicMock(side_effect=mock_cell)
    
    # Simple data
    headers = ["Col1", "Col2"]
    rows = [["A", "B"]]
    
    # Patch methods to avoid side effects
    with patch.object(builder, '_fill_table_with_data'), patch.object(builder, '_hex_to_rgb', return_value=RGBColor(0, 0, 0)):
        for style in style_presets:
            result = builder._add_table_to_slide(mock_pptx_slide, headers, rows, style)
            assert result is not None, f"Table creation with style '{style}' failed"
            mock_pptx_slide.shapes.add_table.assert_called()
            builder._fill_table_with_data.assert_called_with(mock_table, headers, rows, style)