"""
Integration tests for the build method in PPTBuilder.
"""
import os
import tempfile
import pytest
from pathlib import Path
from unittest.mock import MagicMock, patch

from pptx.presentation import Presentation as PptxPresentation
from pptx.slide import Slide as PptxSlide

from doc2pptx.core.models import (
    ContentType, 
    Presentation, 
    Section, 
    Slide, 
    SlideBlock, 
    SlideContent,
    SectionType,
    TableData
)
from doc2pptx.ppt.builder_v3 import PPTBuilder
from doc2pptx.ppt.template_loader import TemplateLoader, TemplateInfo


@pytest.fixture
def builder():
    """Create a PPTBuilder instance."""
    return PPTBuilder()


@pytest.fixture
def mock_template_loader():
    """Mock the TemplateLoader class."""
    with patch('doc2pptx.ppt.builder_v3.TemplateLoader') as MockTemplateLoader:
        # Mock the analyze_template method
        MockTemplateLoader.return_value.analyze_template.return_value = "mock_template_info"
        
        # Mock the load_template method
        mock_pptx = MagicMock(spec=PptxPresentation)
        MockTemplateLoader.return_value.load_template.return_value = mock_pptx
        
        # Setup slide layouts
        layout1 = MagicMock()
        layout1.name = "Diapositive de titre"
        layout2 = MagicMock()
        layout2.name = "Titre et texte"
        layout3 = MagicMock()
        layout3.name = "Titre et tableau"
        mock_pptx.slide_layouts = [layout1, layout2, layout3]
        
        # Create a mock slide
        mock_slide = MagicMock(spec=PptxSlide)
        mock_pptx.slides.add_slide.return_value = mock_slide
        
        # Setup slides collection
        mock_pptx.slides._sldIdLst = []
        
        yield MockTemplateLoader, mock_pptx


def test_build_with_template_path(builder, mock_template_loader):
    """Test building a presentation with a template path."""
    MockTemplateLoader, mock_pptx = mock_template_loader
    
    # Create a simple presentation
    presentation = Presentation(
        id="test_pres",
        title="Test Presentation",
        template_path=Path("test_template.pptx"),
        sections=[
            Section(
                id="section1",
                title="Test Section",
                type=SectionType.CONTENT,
                slides=[
                    Slide(
                        id="slide1",
                        title="Test Slide",
                        layout_name="Titre et texte",
                        blocks=[
                            SlideBlock(
                                id="block1",
                                content=SlideContent(
                                    content_type=ContentType.TEXT,
                                    text="Test content"
                                )
                            )
                        ]
                    )
                ]
            )
        ]
    )
    
    # Mock methods to avoid side effects
    with patch.object(Path, 'exists', return_value=True), \
         patch.object(builder, '_create_slide', return_value=mock_pptx.slides.add_slide.return_value), \
         patch.object(builder, '_fill_slide'), \
         patch.object(builder, '_clear_template_slides'), \
         patch.object(builder, '_validate_layout_for_content', return_value="Titre et texte"), \
         patch.object(LayoutSelector, '_validate_rules_against_template'), \
         patch('pathlib.Path.mkdir'):
        
        # Build the presentation
        output_path = builder.build(presentation, "output.pptx")
        
        # Verify template was loaded
        MockTemplateLoader.return_value.load_template.assert_called_once_with(Path("test_template.pptx"))
        
        # Verify slide was created
        builder._create_slide.assert_called_once_with(mock_pptx, "Titre et texte")
        
        # Verify slide was filled
        builder._fill_slide.assert_called_once()
        
        # Verify presentation was saved
        mock_pptx.save.assert_called_once_with(Path("output.pptx"))
        
        # Verify output path
        assert output_path == Path("output.pptx")


def test_build_with_non_existent_output_directory(builder, mock_template_loader):
    """Test building a presentation with a non-existent output directory."""
    MockTemplateLoader, mock_pptx = mock_template_loader
    
    # Create a simple presentation
    presentation = Presentation(
        id="test_pres",
        title="Test Presentation",
        template_path=Path("test_template.pptx"),
        sections=[
            Section(
                id="section1",
                title="Test Section",
                type=SectionType.CONTENT,
                slides=[
                    Slide(
                        id="slide1",
                        title="Test Slide",
                        layout_name="Titre et texte",
                        blocks=[
                            SlideBlock(
                                id="block1",
                                content=SlideContent(
                                    content_type=ContentType.TEXT,
                                    text="Test content"
                                )
                            )
                        ]
                    )
                ]
            )
        ]
    )
    
    # Mock methods to avoid side effects
    with patch.object(builder, '_create_slide', return_value=mock_pptx.slides.add_slide.return_value), \
         patch.object(builder, '_fill_slide'), \
         patch.object(builder, '_clear_template_slides'), \
         patch.object(builder, '_validate_layout_for_content', return_value="Titre et texte"), \
         patch('pathlib.Path.mkdir') as mock_mkdir:
        
        # Build the presentation
        output_path = builder.build(presentation, "non_existent_dir/output.pptx")
        
        # Verify directory was created
        mock_mkdir.assert_called_once_with(parents=True, exist_ok=True)
        
        # Verify presentation was saved
        mock_pptx.save.assert_called_once_with(Path("non_existent_dir/output.pptx"))


def test_build_with_no_template_path(builder):
    """Test building a presentation with no template path."""
    # Create a presentation without a template path
    presentation = Presentation(
        id="test_pres",
        title="Test Presentation",
        template_path=None,
        sections=[]
    )
    
    # Should raise ValueError
    with pytest.raises(ValueError, match="No template path provided"):
        builder.build(presentation, "output.pptx")


def test_build_with_constructor_template_path(mock_template_loader):
    """Test building a presentation using the template path from the constructor."""
    MockTemplateLoader, mock_pptx = mock_template_loader
    
    # Create a builder with a template path
    builder = PPTBuilder(template_path="constructor_template.pptx")
    
    # Create a presentation without a template path
    presentation = Presentation(
        id="test_pres",
        title="Test Presentation",
        template_path=None,
        sections=[
            Section(
                id="section1",
                title="Test Section",
                type=SectionType.CONTENT,
                slides=[
                    Slide(
                        id="slide1",
                        title="Test Slide",
                        layout_name="Titre et texte",
                        blocks=[
                            SlideBlock(
                                id="block1",
                                content=SlideContent(
                                    content_type=ContentType.TEXT,
                                    text="Test content"
                                )
                            )
                        ]
                    )
                ]
            )
        ]
    )
    
    # Mock methods to avoid side effects
    with patch.object(builder, '_create_slide', return_value=mock_pptx.slides.add_slide.return_value), \
         patch.object(builder, '_fill_slide'), \
         patch.object(builder, '_clear_template_slides'), \
         patch.object(builder, '_validate_layout_for_content', return_value="Titre et texte"), \
         patch('pathlib.Path.mkdir'):
        
        # Build the presentation
        output_path = builder.build(presentation, "output.pptx")
        
        # Verify template was loaded from constructor template path
        MockTemplateLoader.return_value.load_template.assert_called_once_with(Path("constructor_template.pptx"))
        
        # Verify slide was created
        builder._create_slide.assert_called_once_with(mock_pptx, "Titre et texte")
        
        # Verify slide was filled
        builder._fill_slide.assert_called_once()
        
        # Verify presentation was saved
        mock_pptx.save.assert_called_once_with(Path("output.pptx"))


def test_build_with_different_template_path(mock_template_loader):
    """Test building a presentation with a different template path than the constructor."""
    MockTemplateLoader, mock_pptx = mock_template_loader
    
    # Create a builder with a template path
    builder = PPTBuilder(template_path="constructor_template.pptx")
    
    # Create a presentation with a different template path
    presentation = Presentation(
        id="test_pres",
        title="Test Presentation",
        template_path=Path("presentation_template.pptx"),
        sections=[
            Section(
                id="section1",
                title="Test Section",
                type=SectionType.CONTENT,
                slides=[
                    Slide(
                        id="slide1",
                        title="Test Slide",
                        layout_name="Titre et texte",
                        blocks=[
                            SlideBlock(
                                id="block1",
                                content=SlideContent(
                                    content_type=ContentType.TEXT,
                                    text="Test content"
                                )
                            )
                        ]
                    )
                ]
            )
        ]
    )
    
    # Mock methods to avoid side effects
    with patch.object(builder, '_create_slide', return_value=mock_pptx.slides.add_slide.return_value), \
         patch.object(builder, '_fill_slide'), \
         patch.object(builder, '_clear_template_slides'), \
         patch.object(builder, '_validate_layout_for_content', return_value="Titre et texte"), \
         patch('pathlib.Path.mkdir'):
        
        # Build the presentation
        output_path = builder.build(presentation, "output.pptx")
        
        # Verify template was loaded from presentation template path, not constructor
        MockTemplateLoader.return_value.load_template.assert_called_once_with(Path("presentation_template.pptx"))
        
        # Verify template_info was updated for the new template
        MockTemplateLoader.return_value.analyze_template.assert_called_once_with(Path("presentation_template.pptx"))
        
        # Verify slide was created
        builder._create_slide.assert_called_once_with(mock_pptx, "Titre et texte")


def test_build_with_section_header(mock_template_loader):
    """Test building a presentation with a section header."""
    MockTemplateLoader, mock_pptx = mock_template_loader
    
    # Create a presentation
    presentation = Presentation(
        id="test_pres",
        title="Test Presentation",
        template_path=Path("test_template.pptx"),
        sections=[
            Section(
                id="section1",
                title="Test Section",
                type=SectionType.CONTENT,
                slides=[
                    Slide(
                        id="slide1",
                        title="Test Slide",
                        layout_name="Titre et texte",
                        blocks=[
                            SlideBlock(
                                id="block1",
                                content=SlideContent(
                                    content_type=ContentType.TEXT,
                                    text="Test content"
                                )
                            )
                        ]
                    )
                ]
            )
        ]
    )
    
    # Mock methods to avoid side effects
    with patch.object(builder, '_create_slide', return_value=mock_pptx.slides.add_slide.return_value), \
         patch.object(builder, '_fill_slide'), \
         patch.object(builder, '_clear_template_slides'), \
         patch.object(builder, '_validate_layout_for_content', return_value="Titre et texte"), \
         patch.object(builder, '_needs_section_header', return_value=True), \
         patch.object(builder, '_fill_slide_title'), \
         patch('pathlib.Path.mkdir'):
        
        # Build the presentation
        output_path = builder.build(presentation, "output.pptx")
        
        # Verify slide was created twice (section header + content slide)
        assert builder._create_slide.call_count == 2
        
        # Verify slide title was filled for the section header
        builder._fill_slide_title.assert_called_once_with(mock_pptx.slides.add_slide.return_value, "Test Section")


def test_build_with_layout_validation(mock_template_loader):
    """Test building a presentation with layout validation."""
    MockTemplateLoader, mock_pptx = mock_template_loader
    
    # Create a presentation with slides that need layout correction
    presentation = Presentation(
        id="test_pres",
        title="Test Presentation",
        template_path=Path("test_template.pptx"),
        sections=[
            Section(
                id="section1",
                title="Test Section",
                type=SectionType.CONTENT,
                slides=[
                    Slide(
                        id="slide1",
                        title="Table Slide",
                        layout_name="Titre et texte",  # Incorrect layout for a table
                        blocks=[
                            SlideBlock(
                                id="block1",
                                content=SlideContent(
                                    content_type=ContentType.TABLE,
                                    table=TableData(
                                        headers=["Col1", "Col2"],
                                        rows=[["A", "B"], ["C", "D"]]
                                    )
                                )
                            )
                        ]
                    )
                ]
            )
        ]
    )
    
    # Create a builder
    builder = PPTBuilder()
    
    # Mock methods to avoid side effects
    with patch.object(builder, '_create_slide', return_value=mock_pptx.slides.add_slide.return_value), \
         patch.object(builder, '_fill_slide'), \
         patch.object(builder, '_clear_template_slides'), \
         patch.object(builder, '_validate_layout_for_content', return_value="Titre et tableau"), \
         patch('pathlib.Path.mkdir'):
        
        # Build the presentation
        output_path = builder.build(presentation, "output.pptx")
        
        # Verify layout was validated
        builder._validate_layout_for_content.assert_called_once_with(presentation.sections[0].slides[0])
        
        # Verify slide was created with the corrected layout
        builder._create_slide.assert_called_once_with(mock_pptx, "Titre et tableau")


def test_build_with_multiple_sections_and_slides(mock_template_loader):
    """Test building a presentation with multiple sections and slides."""
    MockTemplateLoader, mock_pptx = mock_template_loader
    
    # Create a presentation with multiple sections and slides
    presentation = Presentation(
        id="test_pres",
        title="Test Presentation",
        template_path=Path("test_template.pptx"),
        sections=[
            Section(
                id="section1",
                title="Section 1",
                type=SectionType.TITLE,
                slides=[
                    Slide(
                        id="slide1",
                        title="Title Slide",
                        layout_name="Diapositive de titre",
                        blocks=[
                            SlideBlock(
                                id="block1",
                                content=SlideContent(
                                    content_type=ContentType.TEXT,
                                    text="Subtitle text"
                                )
                            )
                        ]
                    )
                ]
            ),
            Section(
                id="section2",
                title="Section 2",
                type=SectionType.CONTENT,
                slides=[
                    Slide(
                        id="slide2",
                        title="Content Slide 1",
                        layout_name="Titre et texte",
                        blocks=[
                            SlideBlock(
                                id="block2",
                                content=SlideContent(
                                    content_type=ContentType.TEXT,
                                    text="Text content 1"
                                )
                            )
                        ]
                    ),
                    Slide(
                        id="slide3",
                        title="Content Slide 2",
                        layout_name="Titre et texte",
                        blocks=[
                            SlideBlock(
                                id="block3",
                                content=SlideContent(
                                    content_type=ContentType.TEXT,
                                    text="Text content 2"
                                )
                            )
                        ]
                    )
                ]
            )
        ]
    )
    
    # Create a builder
    builder = PPTBuilder()
    
    # Mock methods to avoid side effects
    with patch.object(builder, '_create_slide', return_value=mock_pptx.slides.add_slide.return_value), \
         patch.object(builder, '_fill_slide'), \
         patch.object(builder, '_clear_template_slides'), \
         patch.object(builder, '_validate_layout_for_content', side_effect=lambda slide: slide.layout_name), \
         patch.object(builder, '_needs_section_header', return_value=False), \
         patch('pathlib.Path.mkdir'):
        
        # Build the presentation
        output_path = builder.build(presentation, "output.pptx")
        
        # Verify slides were created for all slides in all sections
        assert builder._create_slide.call_count == 3
        
        # Verify slides were filled for all slides
        assert builder._fill_slide.call_count == 3


def test_build_with_actual_temp_files():
    """Test building a presentation with actual temporary files."""
    # This test will create actual temporary files, so we need to clean up afterwards
    with tempfile.TemporaryDirectory() as tmpdirname:
        # Create temporary paths
        template_path = Path(tmpdirname) / "template.pptx"
        output_path = Path(tmpdirname) / "output.pptx"
        
        # Mock template creation
        with patch('doc2pptx.ppt.builder_v3.TemplateLoader.load_template') as mock_load_template, \
             patch('doc2pptx.ppt.builder_v3.TemplateLoader.analyze_template') as mock_analyze_template:
            
            # Create a mock PPTX that will be "saved"
            mock_pptx = MagicMock(spec=PptxPresentation)
            mock_load_template.return_value = mock_pptx
            mock_analyze_template.return_value = "template_info"
            
            # Setup slide layouts
            layout1 = MagicMock()
            layout1.name = "Diapositive de titre"
            layout2 = MagicMock()
            layout2.name = "Titre et texte"
            mock_pptx.slide_layouts = [layout1, layout2]
            
            # Create a mock slide
            mock_slide = MagicMock(spec=PptxSlide)
            mock_pptx.slides.add_slide.return_value = mock_slide
            
            # Setup slides collection
            mock_pptx.slides._sldIdLst = []
            
            # Create a simple presentation
            presentation = Presentation(
                id="test_pres",
                title="Test Presentation",
                template_path=template_path,
                sections=[
                    Section(
                        id="section1",
                        title="Test Section",
                        type=SectionType.CONTENT,
                        slides=[
                            Slide(
                                id="slide1",
                                title="Test Slide",
                                layout_name="Titre et texte",
                                blocks=[
                                    SlideBlock(
                                        id="block1",
                                        content=SlideContent(
                                            content_type=ContentType.TEXT,
                                            text="Test content"
                                        )
                                    )
                                ]
                            )
                        ]
                    )
                ]
            )
            
            # Create a builder
            builder = PPTBuilder()
            
            # Mock methods to avoid side effects but allow actual filesystem operations
            with patch.object(builder, '_create_slide', return_value=mock_slide), \
                 patch.object(builder, '_fill_slide'), \
                 patch.object(builder, '_clear_template_slides'), \
                 patch.object(builder, '_validate_layout_for_content', return_value="Titre et texte"):
                
                # Build the presentation
                result_path = builder.build(presentation, output_path)
                
                # Verify template was loaded
                mock_load_template.assert_called_once_with(template_path)
                
                # Verify presentation was saved
                mock_pptx.save.assert_called_once_with(output_path)
                
                # Verify returned path is correct
                assert result_path == output_path


def test_build_with_real_file_operations(mock_template_loader):
    """Test building a presentation with real file operations."""
    MockTemplateLoader, mock_pptx = mock_template_loader
    
    # Create a presentation
    presentation = Presentation(
        id="test_pres",
        title="Test Presentation",
        template_path=Path("test_template.pptx"),
        sections=[
            Section(
                id="section1",
                title="Test Section",
                type=SectionType.CONTENT,
                slides=[
                    Slide(
                        id="slide1",
                        title="Test Slide",
                        layout_name="Titre et texte",
                        blocks=[
                            SlideBlock(
                                id="block1",
                                content=SlideContent(
                                    content_type=ContentType.TEXT,
                                    text="Test content"
                                )
                            )
                        ]
                    )
                ]
            )
        ]
    )
    
    # Create a builder
    builder = PPTBuilder()
    
    # Create a real temporary directory
    with tempfile.TemporaryDirectory() as tmpdirname:
        output_path = Path(tmpdirname) / "output.pptx"
        
        # Mock methods to avoid side effects but keep real file operations
        with patch.object(builder, '_create_slide', return_value=mock_pptx.slides.add_slide.return_value), \
             patch.object(builder, '_fill_slide'), \
             patch.object(builder, '_clear_template_slides'), \
             patch.object(builder, '_validate_layout_for_content', return_value="Titre et texte"):
            
            # Build the presentation
            result_path = builder.build(presentation, output_path)
            
            # Verify presentation was saved
            mock_pptx.save.assert_called_once_with(output_path)
            
            # Verify returned path is correct
            assert result_path == output_path