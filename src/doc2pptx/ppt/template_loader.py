"""
Template loader for doc2pptx.

This module provides functionality to load PowerPoint templates and analyze
their layouts for use in presentation generation.
"""
import os
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Optional, Set, Tuple, Union, Any

from pptx import Presentation as PptxPresentation
from pptx.shapes.autoshape import Shape
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from doc2pptx.llm.optimizer import PresentationOptimizer

logger = logging.getLogger(__name__)

@dataclass
class LayoutInfo:
    """Information about a PowerPoint slide layout."""
    name: str
    idx: int
    placeholder_types: List[int]
    placeholder_names: List[str]
    placeholder_indices: List[int]
    supports_title: bool
    supports_content: bool
    supports_image: bool
    supports_chart: bool
    supports_table: bool
    max_content_blocks: int
    
    # AI-enhanced attributes
    ai_description: str = ""
    best_used_for: List[str] = field(default_factory=list)
    ideal_content_types: List[str] = field(default_factory=list)
    limitations: str = ""
    recommendation_score: int = 5
    
    # Placeholder mapping for easy access
    placeholder_map: Dict[str, int] = field(default_factory=dict)


@dataclass
class TemplateInfo:
    """Information about a PowerPoint template."""
    path: Path
    layouts: List[LayoutInfo]
    layout_map: Dict[str, LayoutInfo]
    title_layouts: List[str]
    content_layouts: List[str]
    image_layouts: List[str]
    chart_layouts: List[str]
    table_layouts: List[str]
    two_content_layouts: List[str]


class TemplateLoader:
    """
    Loads and analyzes PowerPoint templates.
    
    This class handles loading PowerPoint templates and extracting information
    about their available layouts, placeholders, and capabilities.
    """
    
    # Mapping from PPT placeholder types to capabilities
    PLACEHOLDER_TYPE_MAP = {
        PP_PLACEHOLDER.TITLE: "title",
        PP_PLACEHOLDER.BODY: "content",
        PP_PLACEHOLDER.CENTER_TITLE: "title",
        PP_PLACEHOLDER.SUBTITLE: "subtitle",
        PP_PLACEHOLDER.PICTURE: "image",
        PP_PLACEHOLDER.CHART: "chart",
        PP_PLACEHOLDER.TABLE: "table",
        PP_PLACEHOLDER.OBJECT: "object",
        # Add more as needed
    }
    
    # Placeholder types that can contain text content
    TEXT_PLACEHOLDER_TYPES = {
        PP_PLACEHOLDER.TITLE,
        PP_PLACEHOLDER.BODY,
        PP_PLACEHOLDER.CENTER_TITLE,
        PP_PLACEHOLDER.SUBTITLE,
        PP_PLACEHOLDER.OBJECT,
    }
    
    def __init__(self):
        """Initialize a template loader."""
        # Cache of template info by path
        self._template_cache: Dict[str, TemplateInfo] = {}
    
    def load_template(self, template_path: Union[str, Path]) -> PptxPresentation:
        """
        Load a PowerPoint template from a file.
        
        Args:
            template_path: Path to the PowerPoint template file
        
        Returns:
            The loaded PowerPoint presentation
            
        Raises:
            FileNotFoundError: If the template file does not exist
            ValueError: If the template file is invalid
        """
        path = Path(template_path)
        if not path.exists():
            raise FileNotFoundError(f"Template file not found: {path}")
        
        try:
            presentation = PptxPresentation(path)
            return presentation
        except Exception as e:
            raise ValueError(f"Invalid PowerPoint template: {e}")
    
    def analyze_template(self, template_path: Union[str, Path]) -> TemplateInfo:
        """
        Analyze a PowerPoint template and extract layout information.
        
        This method loads the template and extracts detailed information about
        each layout, including placeholder types, indices, and capabilities.
        
        Args:
            template_path: Path to the PowerPoint template file
        
        Returns:
            TemplateInfo containing details about the template's layouts
            
        Raises:
            FileNotFoundError: If the template file does not exist
            ValueError: If the template file is invalid
        """
        path = Path(template_path)
        
        # Check if we've already analyzed this template
        if str(path) in self._template_cache:
            return self._template_cache[str(path)]
        
        # Load the template
        presentation = self.load_template(path)
        
        # Analyze layouts
        layouts = []
        layout_map = {}
        
        for idx, layout in enumerate(presentation.slide_layouts):
            # Collect information about placeholders
            placeholder_types = []
            placeholder_names = []
            placeholder_indices = []
            
            for shape in layout.placeholders:
                placeholder_types.append(shape.placeholder_format.type)
                placeholder_names.append(shape.name)
                placeholder_indices.append(shape.placeholder_format.idx)
            
            # Determine layout capabilities
            supports_title = any(ph_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE] 
                               for ph_type in placeholder_types)
            supports_content = any(ph_type == PP_PLACEHOLDER.BODY for ph_type in placeholder_types)
            supports_image = any(ph_type == PP_PLACEHOLDER.PICTURE for ph_type in placeholder_types)
            supports_chart = any(ph_type == PP_PLACEHOLDER.CHART for ph_type in placeholder_types)
            supports_table = any(ph_type == PP_PLACEHOLDER.TABLE for ph_type in placeholder_types)
            
            # Count how many content blocks this layout can support
            content_placeholders = [ph_type for ph_type in placeholder_types 
                                  if ph_type in self.TEXT_PLACEHOLDER_TYPES]
            # Subtract 1 for the title placeholder if it exists
            max_content_blocks = len(content_placeholders)
            if supports_title:
                max_content_blocks = max(0, max_content_blocks - 1)
            
            # Create layout info
            layout_info = LayoutInfo(
                name=layout.name,
                idx=idx,
                placeholder_types=placeholder_types,
                placeholder_names=placeholder_names,
                placeholder_indices=placeholder_indices,
                supports_title=supports_title,
                supports_content=supports_content,
                supports_image=supports_image,
                supports_chart=supports_chart,
                supports_table=supports_table,
                max_content_blocks=max_content_blocks
            )
            
            layouts.append(layout_info)
            layout_map[layout.name] = layout_info
        
        # Categorize layouts by capability
        title_layouts = [layout.name for layout in layouts if layout.supports_title]
        content_layouts = [layout.name for layout in layouts if layout.supports_content]
        image_layouts = [layout.name for layout in layouts if layout.supports_image]
        chart_layouts = [layout.name for layout in layouts if layout.supports_chart]
        table_layouts = [layout.name for layout in layouts if layout.supports_table]
        two_content_layouts = [layout.name for layout in layouts if layout.max_content_blocks >= 2]
        
        # Create template info
        template_info = TemplateInfo(
            path=path,
            layouts=layouts,
            layout_map=layout_map,
            title_layouts=title_layouts,
            content_layouts=content_layouts,
            image_layouts=image_layouts,
            chart_layouts=chart_layouts,
            table_layouts=table_layouts,
            two_content_layouts=two_content_layouts
        )
        
        # Cache the template info
        self._template_cache[str(path)] = template_info
        
        return template_info
    
    def get_best_layout(self, template_info: TemplateInfo, 
                        needs_title: bool = True,
                        needs_content: bool = False,
                        needs_image: bool = False,
                        needs_chart: bool = False,
                        needs_table: bool = False,
                        num_content_blocks: int = 1) -> str:
        """
        Find the best layout in a template that meets the specified needs.
        
        Args:
            template_info: Template information from analyze_template()
            needs_title: Whether the layout needs to support a title
            needs_content: Whether the layout needs to support body/content
            needs_image: Whether the layout needs to support images
            needs_chart: Whether the layout needs to support charts
            needs_table: Whether the layout needs to support tables
            num_content_blocks: Number of content blocks needed
        
        Returns:
            Name of the best matching layout, or the first layout if none match
        """
        # Filter layouts by requirements
        candidates = []
        
        for layout in template_info.layouts:
            if needs_title and not layout.supports_title:
                continue
            if needs_content and not layout.supports_content:
                continue
            if needs_image and not layout.supports_image:
                continue
            if needs_chart and not layout.supports_chart:
                continue
            if needs_table and not layout.supports_table:
                continue
            if layout.max_content_blocks < num_content_blocks:
                continue
            
            candidates.append(layout)
        
        if not candidates:
            # Return default layout if no matches
            return template_info.layouts[0].name if template_info.layouts else "Titre et texte"
        
        # Sort candidates: first by how close the number of content blocks is to what's needed,
        # then by whether they have exactly the right capabilities (no extras)
        def layout_score(layout: LayoutInfo) -> Tuple[int, int]:
            # Score 1: How close the max_content_blocks is to num_content_blocks (lower is better)
            block_diff = abs(layout.max_content_blocks - num_content_blocks)
            
            # Score 2: How many extra capabilities the layout has (lower is better)
            extra_caps = 0
            if layout.supports_title and not needs_title:
                extra_caps += 1
            if layout.supports_content and not needs_content:
                extra_caps += 1
            if layout.supports_image and not needs_image:
                extra_caps += 1
            if layout.supports_chart and not needs_chart:
                extra_caps += 1
            if layout.supports_table and not needs_table:
                extra_caps += 1
            
            return (block_diff, extra_caps)
        
        # Sort candidates by score (lower is better)
        candidates.sort(key=layout_score)
        
        return candidates[0].name
    
    def get_placeholder_mapping(self, template_info: TemplateInfo, layout_name: str) -> Dict[str, int]:
        """
        Get a mapping from placeholder capabilities to placeholder indices.
        
        Args:
            template_info: Template information from analyze_template()
            layout_name: Name of the layout to get placeholders for
        
        Returns:
            Dictionary mapping capability names to placeholder indices
            
        Raises:
            ValueError: If the layout does not exist in the template
        """
        if layout_name not in template_info.layout_map:
            raise ValueError(f"Layout '{layout_name}' not found in template")
        
        layout_info = template_info.layout_map[layout_name]
        
        # Create mapping from capabilities to placeholder indices
        mapping = {}
        
        for i, ph_type in enumerate(layout_info.placeholder_types):
            capability = self.PLACEHOLDER_TYPE_MAP.get(ph_type, None)
            if capability:
                # Only add the first occurrence of each capability
                if capability not in mapping:
                    mapping[capability] = layout_info.placeholder_indices[i]
        
        return mapping
    

    def analyze_template_with_ai(self, template_path_or_object: Union[str, Path, 'PptxPresentation']) -> TemplateInfo:
        """
        Analyze a PowerPoint template and enhance with AI-generated descriptions.
        
        This method uses AI to analyze layouts and suggest optimal use cases.
        
        Args:
            template_path_or_object: Path to the PowerPoint template file or PptxPresentation object
        
        Returns:
            TemplateInfo enhanced with AI-generated descriptions
        """
        # Handle both path and presentation object
        is_presentation_object = False
        
        # Check if it's a Presentation object by checking for common attributes
        if hasattr(template_path_or_object, 'slide_layouts') and hasattr(template_path_or_object, 'slides'):
            is_presentation_object = True
            presentation = template_path_or_object
            path = getattr(presentation, "_path", Path("unknown"))
            # Create template_info from the presentation
            layouts = []
            layout_map = {}
            
            for idx, layout in enumerate(presentation.slide_layouts):
                # Collect information about placeholders
                placeholder_types = []
                placeholder_names = []
                placeholder_indices = []
                
                for shape in layout.placeholders:
                    placeholder_types.append(shape.placeholder_format.type)
                    placeholder_names.append(shape.name)
                    placeholder_indices.append(shape.placeholder_format.idx)
                
                # Determine layout capabilities
                supports_title = any(ph_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE] 
                                for ph_type in placeholder_types)
                supports_content = any(ph_type == PP_PLACEHOLDER.BODY for ph_type in placeholder_types)
                supports_image = any(ph_type == PP_PLACEHOLDER.PICTURE for ph_type in placeholder_types)
                supports_chart = any(ph_type == PP_PLACEHOLDER.CHART for ph_type in placeholder_types)
                supports_table = any(ph_type == PP_PLACEHOLDER.TABLE for ph_type in placeholder_types)
                
                # Count how many content blocks this layout can support
                content_placeholders = [ph_type for ph_type in placeholder_types 
                                    if ph_type in self.TEXT_PLACEHOLDER_TYPES]
                # Subtract 1 for the title placeholder if it exists
                max_content_blocks = len(content_placeholders)
                if supports_title:
                    max_content_blocks = max(0, max_content_blocks - 1)
                
                # Create layout info
                layout_info = LayoutInfo(
                    name=layout.name,
                    idx=idx,
                    placeholder_types=placeholder_types,
                    placeholder_names=placeholder_names,
                    placeholder_indices=placeholder_indices,
                    supports_title=supports_title,
                    supports_content=supports_content,
                    supports_image=supports_image,
                    supports_chart=supports_chart,
                    supports_table=supports_table,
                    max_content_blocks=max_content_blocks
                )
                
                layouts.append(layout_info)
                layout_map[layout.name] = layout_info
            
            # Categorize layouts by capability
            title_layouts = [layout.name for layout in layouts if layout.supports_title]
            content_layouts = [layout.name for layout in layouts if layout.supports_content]
            image_layouts = [layout.name for layout in layouts if layout.supports_image]
            chart_layouts = [layout.name for layout in layouts if layout.supports_chart]
            table_layouts = [layout.name for layout in layouts if layout.supports_table]
            two_content_layouts = [layout.name for layout in layouts if layout.max_content_blocks >= 2]
            
            # Create template info
            template_info = TemplateInfo(
                path=path,
                layouts=layouts,
                layout_map=layout_map,
                title_layouts=title_layouts,
                content_layouts=content_layouts,
                image_layouts=image_layouts,
                chart_layouts=chart_layouts,
                table_layouts=table_layouts,
                two_content_layouts=two_content_layouts
            )
        else:
            # Get basic template info first using the path
            template_path = Path(template_path_or_object)
            template_info = self.analyze_template(template_path)
        
        try:
            # Initialize the AI optimizer
            optimizer = PresentationOptimizer()
            
            if not optimizer.client:
                logger.warning("AI client not available. Skipping AI analysis of template.")
                return template_info
                
            # Get enhanced information about the layouts
            enhanced_layout_info = self._get_layout_descriptions(optimizer, template_info)
            
            # Update template_info with enhanced information
            self._update_template_info_with_ai_insights(template_info, enhanced_layout_info)
            
            return template_info
            
        except Exception as e:
            logger.warning(f"Error in AI analysis of template: {e}. Using basic analysis.")
            return template_info
            
    def _get_layout_descriptions(self, optimizer: PresentationOptimizer, 
                                template_info: TemplateInfo) -> Dict[str, Dict[str, Any]]:
        """
        Use AI to get descriptions and recommended uses for layouts.
        
        Args:
            optimizer: The PresentationOptimizer instance
            template_info: Basic template info
            
        Returns:
            Dictionary mapping layout names to enhanced information
        """
        # Prepare layout information for AI analysis
        layout_info_dict = {}
        for layout in template_info.layouts:
            layout_info_dict[layout.name] = {
                "supports_title": layout.supports_title,
                "supports_content": layout.supports_content,
                "supports_image": layout.supports_image,
                "supports_chart": layout.supports_chart,
                "supports_table": layout.supports_table,
                "max_content_blocks": layout.max_content_blocks,
                "placeholder_types": [self.PLACEHOLDER_TYPE_MAP.get(ph_type, str(ph_type)) 
                                    for ph_type in layout.placeholder_types]
            }
        
        # Check if optimizer has a client
        if not optimizer.client:
            logger.warning("OpenAI client not available. Skipping layout analysis.")
            return {}
        
        try:
            # Use the optimizer's analyze_template_layouts method
            enhanced_info = optimizer.analyze_template_layouts(layout_info_dict)
            
            # Validate and clean results
            validated_info = {}
            for layout_name, layout_data in enhanced_info.items():
                # Ensure we have all required fields
                validated_layout = {
                    "description": layout_data.get("description", ""),
                    "best_used_for": layout_data.get("best_used_for", []),
                    "ideal_content_types": layout_data.get("ideal_content_types", []),
                    "limitations": layout_data.get("limitations", ""),
                    "recommendation_score": layout_data.get("recommendation_score", 5)
                }
                validated_info[layout_name] = validated_layout
                
            return validated_info
        except Exception as e:
            logger.error(f"Error getting layout descriptions: {e}")
            return {}
    
    def _update_template_info_with_ai_insights(self, template_info: TemplateInfo, 
                                              enhanced_info: Dict[str, Dict[str, Any]]) -> None:
        """
        Update template_info with AI-generated layout insights.
        
        Args:
            template_info: Template info to update
            enhanced_info: AI-generated layout insights
        """
        # Add AI descriptions to each layout in layout_map
        for layout_name, layout_info in enhanced_info.items():
            if layout_name in template_info.layout_map:
                # Create a new attribute for the AI information
                template_info.layout_map[layout_name].ai_description = layout_info.get("description", "")
                template_info.layout_map[layout_name].best_used_for = layout_info.get("best_used_for", [])
                template_info.layout_map[layout_name].ideal_content_types = layout_info.get("ideal_content_types", [])
                template_info.layout_map[layout_name].limitations = layout_info.get("limitations", "")
                template_info.layout_map[layout_name].recommendation_score = layout_info.get("recommendation_score", 5)