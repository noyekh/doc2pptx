"""
Overflow detection and handling for doc2pptx.

This module provides functionality to detect and handle content overflow
in PowerPoint slide placeholders.
"""
import logging
import re
from typing import Dict, List, Optional, Tuple, Union, Any

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.text import CT_TextParagraph
from pptx.shapes.base import BaseShape
from pptx.text.text import TextFrame, _Paragraph
from pptx.util import Pt, Inches, Emu

from doc2pptx.llm.optimizer import PresentationOptimizer
from doc2pptx.core.models import Section, Slide, SlideBlock, SlideContent, ContentType

logger = logging.getLogger(__name__)


class OverflowHandler:
    """
    Detects and handles content overflow in PowerPoint slides.
    
    This class provides methods to check if content will overflow a placeholder,
    and strategies to handle overflow when it occurs.
    """
    HARD_LENGTH_THRESHOLD: int = 2048
    
    def __init__(self, 
                 max_chars_per_line: int = 90,
                 max_lines_per_text_box: int = 15,
                 avg_char_width_pt: float = 7.0,
                 avg_line_height_pt: float = 18.0,
                 use_ai: bool = True):
        """
        Initialize an overflow handler with default settings.
        
        Args:
            max_chars_per_line: Maximum number of characters per line for estimation.
            max_lines_per_text_box: Maximum number of lines per text box for estimation.
            avg_char_width_pt: Average character width in points for estimation.
            avg_line_height_pt: Average line height in points for estimation.
        """
        self.max_chars_per_line = max_chars_per_line
        self.max_lines_per_text_box = max_lines_per_text_box
        self.avg_char_width_pt = avg_char_width_pt
        self.avg_line_height_pt = avg_line_height_pt
        self.use_ai = use_ai
        
        # Initialize AI optimizer if enabled
        if self.use_ai:
            try:
                self.optimizer = PresentationOptimizer()
            except Exception as e:
                logger.warning(f"Could not initialize AI optimizer: {e}. Falling back to rule-based overflow handling.")
                self.use_ai = False
    
    def handle_slide_overflow(self, section: Section, slide: Slide, shape_width: int, shape_height: int) -> List[Slide]:
        """
        Handle content overflow by creating additional slides if necessary.
        """
        if not self.use_ai:
            # If AI is not enabled, just return the original slide
            return [slide]
        
        try:
            # Use AI to analyze overflow
            analysis = self.optimizer.analyze_content_overflow(slide, shape_width, shape_height)
            
            if not analysis["may_overflow"]:
                # No overflow or not enough to warrant a new slide
                return [slide]
            
            # Create additional slides
            result_slides = []
            
            # Modify the original slide with the first part of the content
            original_slide = slide
            original_blocks = slide.blocks
            
            # Process each part of the content from split_recommendation instead of content_division
            split_recommendation = analysis.get("split_recommendation", [])
            
            if not split_recommendation or len(split_recommendation) <= 1:
                # No valid split recommendation
                return [slide]
                
            # Process each division suggested by AI
            for i, content_part in enumerate(split_recommendation):
                if i == 0:
                    # First part goes in the original slide
                    new_blocks = self._create_blocks_from_content(content_part, original_blocks)
                    original_slide.blocks = new_blocks
                    result_slides.append(original_slide)
                else:
                    # Create a new slide for each additional part
                    new_slide_id = f"{slide.id}-part-{i+1}"
                    new_slide_title = f"{slide.title} (continued)"
                    
                    # Create blocks for the new slide
                    new_blocks = self._create_blocks_from_content(content_part, original_blocks)
                    
                    # Create the new slide with the same layout as the original
                    new_slide = Slide(
                        id=new_slide_id,
                        title=new_slide_title,
                        layout_name=slide.layout_name,
                        blocks=new_blocks,
                        notes=f"Continuation of slide {slide.title}"
                    )
                    
                    result_slides.append(new_slide)
            
            return result_slides
                    
        except Exception as e:
            logger.error(f"Error handling slide overflow: {e}")
            # Return the original slide on error
            return [slide]
       
    def _create_blocks_from_content(self, content: str, original_blocks: List[SlideBlock]) -> List[SlideBlock]:
        """
        Create slide blocks from the given content, preserving the structure of the original blocks.
        
        Args:
            content: Content to add to the new blocks.
            original_blocks: Original blocks to use as templates.
            
        Returns:
            List of new slide blocks.
        """
        # If we have multiple blocks, we need to distribute the content intelligently
        if len(original_blocks) > 1:
            # For now, just put all content in the first block
            # In a real implementation, we'd need more sophisticated distribution
            new_blocks = []
            first_block = True
            
            for block in original_blocks:
                new_block = SlideBlock(
                    id=f"{block.id}-new",
                    title=block.title,
                    style=block.style,
                    position=block.position
                )
                
                if first_block and block.content:
                    # Put the content in the first content block
                    content_type = block.content.content_type
                    
                    if content_type == ContentType.TEXT:
                        new_block.content = SlideContent(
                            content_type=ContentType.TEXT,
                            text=content
                        )
                    elif content_type == ContentType.BULLET_POINTS:
                        # Split content into bullet points
                        bullet_points = [line.strip() for line in content.split('\n') if line.strip()]
                        new_block.content = SlideContent(
                            content_type=ContentType.BULLET_POINTS,
                            bullet_points=bullet_points,
                            as_bullets=block.content.as_bullets
                        )
                    else:
                        # For other content types, just copy the original content
                        new_block.content = block.content
                    
                    first_block = False
                else:
                    # Copy original content for other blocks
                    new_block.content = block.content
                
                new_blocks.append(new_block)
            
            return new_blocks
        else:
            # Single block case
            original_block = original_blocks[0]
            content_type = original_block.content.content_type if original_block.content else ContentType.TEXT
            
            new_block = SlideBlock(
                id=f"{original_block.id}-new",
                title=original_block.title,
                style=original_block.style,
                position=original_block.position
            )
            
            if content_type == ContentType.TEXT:
                new_block.content = SlideContent(
                    content_type=ContentType.TEXT,
                    text=content
                )
            elif content_type == ContentType.BULLET_POINTS:
                # Split content into bullet points
                bullet_points = [line.strip() for line in content.split('\n') if line.strip()]
                new_block.content = SlideContent(
                    content_type=ContentType.BULLET_POINTS,
                    bullet_points=bullet_points,
                    as_bullets=True if original_block.content else True
                )
            else:
                # For other content types, just copy the original content
                new_block.content = original_block.content
            
            return [new_block]
            
    def will_text_overflow(
        self,
        shape: BaseShape,
        text: str,
        font_size_pt: float = 12.0,
    ) -> bool:
        """
        Estime si `text` débordera du placeholder `shape`.

        • Conversion correcte des dimensions (EMU ➜ points) ;
        • prise en compte du corps de police ;
        • garde-fou « longueur brute » : au-delà de
          `HARD_LENGTH_THRESHOLD`, l’overflow est déclaré certain.
        """
        # rien à tester ?
        if not text:
            return False

        # garde-fou immédiat
        if len(text) > self.HARD_LENGTH_THRESHOLD:
            logger.warning(
                "Text length (%d) exceeds hard threshold (%d) – overflow assumed.",
                len(text),
                self.HARD_LENGTH_THRESHOLD,
            )
            return True

        # le shape doit contenir un text_frame
        if not hasattr(shape, "text_frame"):
            logger.warning(
                "Shape missing text_frame – falling back to simple estimation."
            )
            return self._estimate_overflow(text)

        # ---------------------------------------------------------------- geometry
        try:
            # `shape.width` / `height` sont en EMU (int) ; on les convertit en pt
            width_pt = Emu(shape.width).pt
            height_pt = Emu(shape.height).pt
        except Exception:
            logger.warning(
                "Could not convert shape dimensions (EMU ➜ pt) – using estimation."
            )
            return self._estimate_overflow(text)

        # ---------------------------------------------------------------- ratios
        char_width = self.avg_char_width_pt * (font_size_pt / 12.0)
        line_height = self.avg_line_height_pt * (font_size_pt / 12.0)

        usable_width = width_pt * 0.95   # marges ≈ 5 %
        usable_height = height_pt * 0.95

        max_chars = max(1, int(usable_width / char_width))
        max_lines = max(1, int(usable_height / line_height))

        # ---------------------------------------------------------------- calc
        lines_needed = self._estimate_lines_needed(text, max_chars)
        is_overflow = lines_needed > max_lines

        if is_overflow:
            logger.warning(
                "Text will likely overflow (%d lines needed, %d fit).",  # noqa: G004
                lines_needed,
                max_lines,
            )

        return is_overflow

    
    def _estimate_lines_needed(self, text: str, chars_per_line: int) -> int:
        """
        Estimate the number of lines needed for the given text.
        
        Args:
            text: Text to estimate lines for.
            chars_per_line: Maximum number of characters per line.
        
        Returns:
            Estimated number of lines needed.
        """
        # Si le texte est vide, aucune ligne n'est nécessaire
        if not text:
            return 0
            
        # Split text into paragraphs
        paragraphs = text.split('\n')
        
        total_lines = 0
        
        for paragraph in paragraphs:
            if not paragraph.strip():
                # Empty paragraph, count as one line
                total_lines += 1
                continue
            
            # Count words
            words = paragraph.split()
            
            if not words:
                # Empty paragraph, count as one line
                total_lines += 1
                continue
            
            # Estimate lines needed for this paragraph
            current_line_length = 0
            paragraph_lines = 1
            
            for word in words:
                # Length of word plus space
                word_length = len(word) + 1
                
                if current_line_length + word_length > chars_per_line:
                    # Word won't fit on current line, start a new line
                    paragraph_lines += 1
                    current_line_length = word_length
                else:
                    # Word fits on current line
                    current_line_length += word_length
            
            total_lines += paragraph_lines
        
        return total_lines
    
    def _estimate_overflow(self, text: str) -> bool:
        """
        Perform a simple estimation of overflow based on text length.
        
        This is used as a fallback when shape dimensions cannot be determined.
        
        Args:
            text: Text to check for overflow.
        
        Returns:
            True if the text will likely overflow, False otherwise.
        """
        # Split text into paragraphs
        paragraphs = text.split('\n')
        
        # Count total characters and estimate lines
        total_chars = sum(len(p) for p in paragraphs)
        estimated_lines = self._estimate_lines_needed(text, self.max_chars_per_line)
        
        # Check if exceeds maximum
        chars_limit = self.max_chars_per_line * self.max_lines_per_text_box
        
        is_overflow = (estimated_lines > self.max_lines_per_text_box) or (total_chars > chars_limit)
        
        if is_overflow:
            logger.warning(
                f"Text may overflow based on simple estimation. "
                f"Text has approximately {estimated_lines} lines, "
                f"but default maximum is {self.max_lines_per_text_box} lines."
            )
        
        return is_overflow
    
    def split_text_for_overflow(self, text: str, max_chars_per_slide: int = 1000) -> List[str]:
        """
        Split text into multiple chunks to handle overflow.
        
        This can be used to create multiple slides for long content.
        
        Args:
            text: Text to split.
            max_chars_per_slide: Maximum characters per slide chunk.
        
        Returns:
            List of text chunks, each suitable for a single slide.
        """
        # Si le texte est vide, retourner une liste avec une seule chaîne vide
        if not text:
            return [""]
            
        # Split text into paragraphs
        paragraphs = text.split('\n')
        
        chunks = []
        current_chunk = ""
        
        for paragraph in paragraphs:
            # Si ce paragraphe est trop long en soi, le diviser en morceaux plus petits
            if len(paragraph) > max_chars_per_slide:
                # Ajouter le chunk courant s'il existe
                if current_chunk:
                    chunks.append(current_chunk)
                    current_chunk = ""
                    
                # Diviser le paragraphe long en morceaux de max_chars_per_slide
                for i in range(0, len(paragraph), max_chars_per_slide):
                    chunk = paragraph[i:i + max_chars_per_slide]
                    chunks.append(chunk)
            # Si l'ajout de ce paragraphe dépasserait la limite, commencer un nouveau chunk
            elif len(current_chunk) + len(paragraph) + 1 > max_chars_per_slide:
                if current_chunk:  # Ne pas ajouter de chunks vides
                    chunks.append(current_chunk)
                current_chunk = paragraph
            else:
                # Ajouter le paragraphe au chunk courant
                if current_chunk:
                    current_chunk += '\n' + paragraph
                else:
                    current_chunk = paragraph
        
        # Ajouter le dernier chunk s'il n'est pas vide
        if current_chunk:
            chunks.append(current_chunk)
        
        return chunks or [""]  # Retourner au moins un chunk, même s'il est vide
    
    def split_bullet_points_for_overflow(self, bullet_points: List[str], max_points_per_slide: int = 10) -> List[List[str]]:
        """
        Split bullet points into multiple chunks to handle overflow.
        
        This can be used to create multiple slides for long bullet lists.
        
        Args:
            bullet_points: List of bullet points to split.
            max_points_per_slide: Maximum number of bullet points per slide.
        
        Returns:
            List of bullet point chunks, each suitable for a single slide.
        """
        chunks = []
        
        for i in range(0, len(bullet_points), max_points_per_slide):
            chunk = bullet_points[i:i + max_points_per_slide]
            chunks.append(chunk)
        
        return chunks or [[]]  # Return at least one chunk, even if empty
    
    def truncate_text_for_placeholder(self, shape: BaseShape, 
                                     text: str,
                                     add_ellipsis: bool = True) -> str:
        """
        Truncate text to fit within a placeholder.
        
        Args:
            shape: PowerPoint shape to fit text into.
            text: Text to truncate.
            add_ellipsis: Whether to add "..." at the end of truncated text.
        
        Returns:
            Truncated text that will fit in the placeholder.
        """
        if not hasattr(shape, "text_frame"):
            logger.warning("Shape does not have a text frame. Cannot truncate text.")
            return text

        try:
            width_obj = shape.width
            height_obj = shape.height
            width = getattr(width_obj, "pt", width_obj)
            height = getattr(height_obj, "pt", height_obj)
            if not isinstance(width, (int, float)) or not isinstance(height, (int, float)):
                raise TypeError
        except (AttributeError, TypeError, ValueError):
            logger.warning("Could not get shape dimensions. Using estimation.")
            return self._truncate_text_by_chars(text, add_ellipsis)

        
        # Calculate max chars per line and max lines
        char_width = self.avg_char_width_pt
        line_height = self.avg_line_height_pt
        
        # Allow for margins (estimate 5% of width and height)
        usable_width = width * 0.95
        usable_height = height * 0.95
        
        max_chars_per_line = int(usable_width / char_width)
        max_lines = int(usable_height / line_height)
        
        # Initialize truncated text
        truncated_text = ""
        lines_used = 0
        
        # Process text line by line
        paragraphs = text.split('\n')
        
        for paragraph in paragraphs:
            if lines_used >= max_lines:
                break
            
            # If paragraph is empty, add a newline and continue
            if not paragraph.strip():
                truncated_text += '\n'
                lines_used += 1
                continue
            
            # Split paragraph into words
            words = paragraph.split()
            
            current_line = ""
            
            for word in words:
                # Try adding the word to the current line
                if len(current_line) + len(word) + 1 <= max_chars_per_line:
                    # Word fits on current line
                    if current_line:
                        current_line += ' ' + word
                    else:
                        current_line = word
                else:
                    # Word doesn't fit, add current line to truncated text and start new line
                    truncated_text += current_line + '\n'
                    lines_used += 1
                    
                    if lines_used >= max_lines:
                        break
                    
                    # Start new line with current word
                    current_line = word
            
            # Add the last line of the paragraph
            if current_line and lines_used < max_lines:
                truncated_text += current_line + '\n'
                lines_used += 1
        
        # Remove trailing newline
        if truncated_text.endswith('\n'):
            truncated_text = truncated_text[:-1]
        
        # Add ellipsis if text was truncated and ellipsis is requested
        if add_ellipsis and truncated_text != text:
            truncated_text = truncated_text + " ..."
        
        return truncated_text
    
    def _truncate_text_by_chars(self, text: str, add_ellipsis: bool = True) -> str:
        """
        Truncate text based on character count.
        
        This is a simpler fallback method when shape dimensions are not available.
        
        Args:
            text: Text to truncate.
            add_ellipsis: Whether to add "..." at the end of truncated text.
        
        Returns:
            Truncated text.
        """
        max_chars = self.max_chars_per_line * self.max_lines_per_text_box
        
        if len(text) < max_chars:
            return text
        
        # Truncate to max chars, but try to break at a word boundary
        # Réserver de l'espace pour l'ellipse si nécessaire
        ellipsis_len = 4 if add_ellipsis else 0
        truncated = text[:max_chars - ellipsis_len]
        
        # Try to find the last space to avoid cutting words
        last_space = truncated.rfind(' ')
        if last_space > max_chars * 0.8:  # Only use the last space if it's reasonably close to the end
            truncated = truncated[:last_space]
        
        # Add ellipsis if requested
        if add_ellipsis:
            truncated += " ..."
        
        return truncated